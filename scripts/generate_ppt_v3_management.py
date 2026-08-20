#!/usr/bin/env python3
"""Generate SEMS_功能介绍.pptx V2026.08 (25 slides) — MANAGEMENT WORKFLOW edition.

Replaces P17-P25 with **management (not operational) workflows**:
  P17 管理流程总览：9大管理闭环 × 角色分工矩阵
  P18 设备全生命周期管理流程（T0-T3资产 + 状态 + 润滑）
  P19 故障响应管理流程（DOWN→工单→SLA→升级→关单考核）
  P20 预防性维护管理流程（PM计划 + 点检巡检 + 漏做追责）
  P21 文控合规管理流程（三级审批+复审+分发+水印+修正）
  P22 电子表单管理流程（模板治理+审核锁+附加修正）
  P23 品质改进管理流程（8D团队+时限+FMEA RPN阈值+整改闭环）
  P24 知识沉淀与复发治理流程（双归档+审核+复发计数+根因根治）
  P25 系统治理与灾备管理流程（账号权限+审计+备份3-2-1+季度演练）
"""

import io
import zipfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Design System ──
DARK_BLUE  = RGBColor(0x0B, 0x3D, 0x91)
TEAL       = RGBColor(0x00, 0x96, 0x88)
LIGHT_GRAY = RGBColor(0xF5, 0xF7, 0xFA)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x2C, 0x3E, 0x50)
GRAY_TEXT  = RGBColor(0x7F, 0x8C, 0x8D)
ACCENT     = RGBColor(0x1A, 0x73, 0xE8)
CARD_BG    = RGBColor(0xEE, 0xF2, 0xF7)
RED        = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE     = RGBColor(0xF3, 0x9C, 0x12)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
PURPLE     = RGBColor(0x8E, 0x44, 0xAD)
FONT       = 'Microsoft YaHei'

SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

# ── Helpers ──
def _set_run(run, text, size, color, bold=False, font_name=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name

def _add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def _add_circle(slide, left, top, size, fill_color, text, text_color=WHITE, text_size=14):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    _set_run(run, text, text_size, text_color, bold=True)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape

def _add_arrow(slide, left, top, width, height, color=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def _add_down_arrow(slide, left, top, width, height, color=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def _add_text(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    _set_run(run, text, size, color, bold)
    return txBox

def _add_bullets(slide, left, top, width, height, items, size=14, color=DARK_TEXT, spacing=6):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        if isinstance(item, tuple):
            prefix, rest = item
            r1 = p.add_run()
            _set_run(r1, prefix, size, TEAL, bold=True)
            r2 = p.add_run()
            _set_run(r2, rest, size, color)
        else:
            r = p.add_run()
            _set_run(r, item, size, color)
    return txBox

def _add_top_bar(slide, title, subtitle=None):
    _add_rect(slide, 0, 0, SLIDE_W, Emu(914400), DARK_BLUE)
    _add_text(slide, Emu(457200), Emu(137160), Emu(9000000), Emu(457200),
              title, 24, WHITE, bold=True)
    if subtitle:
        _add_text(slide, Emu(457200), Emu(548640), Emu(9000000), Emu(274320),
                  subtitle, 13, RGBColor(0xAE, 0xC5, 0xE8))

def _add_card(slide, left, top, width, height, title, items, title_color=DARK_BLUE, item_size=12):
    _add_rect(slide, left, top, width, height, LIGHT_GRAY)
    _add_text(slide, left + Emu(182880), top + Emu(91440),
              width - Emu(365760), Emu(365760),
              title, 14, title_color, bold=True)
    _add_bullets(slide, left + Emu(182880), top + Emu(457200),
                 width - Emu(365760), height - Emu(548640),
                 items, size=item_size, spacing=4)

def _extract_screenshot(orig_path):
    z = zipfile.ZipFile(orig_path)
    for name in z.namelist():
        if name.startswith('ppt/media/') and name.endswith('.png'):
            return z.read(name)
    return None

# ── Swim-lane style management stage builder ──
def _add_mgmt_stages(slide, left, top, width, stages):
    """
    Horizontal management stages with RACI-ish role tags.
    stages = [(stage_name, responsible_role, action_desc, deliverable, color), ...]
    """
    n = len(stages)
    stage_w = (width - Emu(400000) * (n - 1)) // n
    x = left
    for i, (name, role, action, output, color) in enumerate(stages):
        # Stage card
        _add_rect(slide, x, top, stage_w, Emu(2600000), LIGHT_GRAY)
        # Top color bar
        _add_rect(slide, x, top, stage_w, Emu(140000), color)
        # Stage number circle
        _add_circle(slide, x + Emu(137160), top + Emu(200000), Emu(400000), color, str(i+1), WHITE, 13)
        # Stage name
        _add_text(slide, x + Emu(570000), top + Emu(210000), stage_w - Emu(600000), Emu(300000),
                  name, 13, WHITE if False else color, bold=True)
        # Role pill
        _add_rect(slide, x + Emu(137160), top + Emu(700000), Emu(1800000), Emu(360000), color)
        tf = slide.shapes[-1].text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); _set_run(run, f"责任：{role}", 10, WHITE, bold=True)
        # Action
        _add_text(slide, x + Emu(137160), top + Emu(1150000), stage_w - Emu(274320), Emu(700000),
                  f"动作：{action}", 11, DARK_TEXT)
        # Output / deliverable
        _add_text(slide, x + Emu(137160), top + Emu(1900000), stage_w - Emu(274320), Emu(650000),
                  f"输出：{output}", 10, TEAL, bold=True)
        if i < n - 1:
            _add_arrow(slide, x + stage_w, top + Emu(1200000), Emu(300000), Emu(200000), GRAY_TEXT)
        x += stage_w + Emu(400000)

def _add_mgmt_footer(slide, metrics_left, metrics_right, top=Emu(5700000)):
    """Bottom two-column KPIs / exception handling / governance rules."""
    _add_rect(slide, Emu(457200), top, Emu(5600000), Emu(900000), CARD_BG)
    _add_text(slide, Emu(640080), top + Emu(91440), Emu(5400000), Emu(300000),
              "📊 KPI / 考核指标", 12, DARK_BLUE, bold=True)
    _add_bullets(slide, Emu(640080), top + Emu(390000), Emu(5400000), Emu(520000),
                 metrics_left, size=10, spacing=3)
    _add_rect(slide, Emu(6200000), top, Emu(5550000), Emu(900000), LIGHT_GRAY)
    _add_text(slide, Emu(6350000), top + Emu(91440), Emu(5400000), Emu(300000),
              "⚠ 异常处理 / 治理规则", 12, RED, bold=True)
    _add_bullets(slide, Emu(6350000), top + Emu(390000), Emu(5400000), Emu(520000),
                 metrics_right, size=10, spacing=3)

# ════════════════════════════════════════════════════════════════════════════
# Original 16 slides (unchanged structure, same as v2)
# ════════════════════════════════════════════════════════════════════════════

def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_BLUE)
    _add_rect(slide, Emu(4572000), Emu(3200400), Emu(3047695), Emu(38100), TEAL)
    _add_text(slide, Emu(914400), Emu(1828800), Emu(10362895), Emu(914400),
              "SEMS 半导体设备管理系统", 40, WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, Emu(914400), Emu(2833920), Emu(10362895), Emu(548640),
              "Semiconductor Equipment Management System", 20,
              RGBColor(0xAE, 0xC5, 0xE8), align=PP_ALIGN.CENTER)
    _add_text(slide, Emu(914400), Emu(4572000), Emu(10362895), Emu(457200),
              "功能介绍 + 管理流程  |  V2026.08  |  2026-08", 18,
              RGBColor(0xAE, 0xC5, 0xE8), align=PP_ALIGN.CENTER)
    _add_text(slide, Emu(914400), Emu(5486400), Emu(10362895), Emu(365760),
              "FastAPI  +  Vue 3  +  Element Plus  +  SQLite", 14,
              RGBColor(0x8A, 0xAA, 0xC8), align=PP_ALIGN.CENTER)

def build_why_sems(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "为什么需要 SEMS？", "Why SEMS?")
    pain_items = [
        "设备状态靠电话/微信群，台账分散在 Excel/纸档",
        "故障处置经验沉淀差，同一根因反复复发",
        "点检/PM/润滑执行难追踪，容易漏做",
        "备件库存不透明，等要修了才发现缺件",
        "工艺文件版本与审批靠邮件，追溯困难",
        "8D 报告写完就丢，知识库与品管工具完全脱节",
    ]
    _add_bullets(slide, Emu(457200), Emu(1280160), Emu(7000000), Emu(4500000),
                 [("• ", p) for p in pain_items], size=15, spacing=10)
    _add_card(slide, Emu(7700000), Emu(1280160), Emu(4000000), Emu(2500000),
              "一句话价值",
              ["把设备台账、日常运维、工艺文控、品管合规、数据价值",
               "五个核心闭环落在一套系统里，",
               "真正做到数据可查、经验可复用。"], title_color=TEAL)

def build_five_loops(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "五大业务闭环", "Five Business Loops")
    loops = [
        ("① 设备台账 & 状态闭环", "设备台账 / 状态切换 / 设备全生命周期 T0-T3 / 润滑管理", "机台状态一眼看清，DOWN 自动派工单"),
        ("② 日常运维闭环", "工单 / PM 计划 / 点检巡检 / 备件管理 / 工单 SLA 升级", "从建单→维修→验证→关单一气呵成，超时自动升级"),
        ("③ 工艺文控闭环", "工艺文件 / 表单模板 / 电子表单 / 文控审批链 / 分发收回 / 复审告警", "符合体系标准的受控文控，水印+指纹+留痕"),
        ("④ 品管合规闭环", "8D / FMEA / 安全检查 / 环境核查", "质量分析合规可追溯，证书到期提前告警"),
        ("⑤ 数据价值闭环", "故障知识库(工单归档+8D归档) / OEE / 设备成本 LCC", "故障经验不随人流失，LCC 一眼算清单台成本"),
    ]
    y = Emu(1100000)
    for title, modules, value in loops:
        _add_rect(slide, Emu(457200), y, Emu(11277000), Emu(914400), LIGHT_GRAY)
        _add_text(slide, Emu(640080), y + Emu(91440), Emu(3200000), Emu(365760), title, 14, DARK_BLUE, bold=True)
        _add_text(slide, Emu(3900000), y + Emu(91440), Emu(5000000), Emu(731520), modules, 11, DARK_TEXT)
        _add_text(slide, Emu(9000000), y + Emu(91440), Emu(2600000), Emu(731520), value, 11, TEAL)
        y += Emu(1000000)

def build_v2026_updates(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "V2026.08 版本核心更新", "Key Updates in V2026.08")
    cards = [
        ("📚 8D → 知识库 一键归档", ["D0+D2 → 故障现象 / D4 → 根因", "D5 → 处置 / D7 → 预防", "关联 source_d8_report_id 溯源", "同时保留「工单归档」入口"]),
        ("📋 8 大分组分级菜单", ["① 总览 ② 设备管理 ③ 运维工单", "④ 安全与环境 ⑤ 备件与人员", "⑥ 工艺文控 ⑦ 分析改进", "⑧ 系统配置（永远在最下方）"]),
        ("🎯 角色定制看板", ["6 种角色登录→不同看板组合", "QA 优先：文控复审 + 8D 进度", "操作员优先：设备状态 + 今日点检"]),
        ("🧾 审计日志入库 + 🔐 bcrypt 修复", ["11 类敏感操作写入 audit_logs 表", "操作人/目标/IP/UA/时间戳/详情", "SHA-256 预哈希解决 bcrypt 72 字节截断", "旧值透明回退 + 登录自动升级重哈希"]),
    ]
    positions = [
        (Emu(457200), Emu(1100000), Emu(5600000), Emu(2500000)),
        (Emu(6200000), Emu(1100000), Emu(5600000), Emu(2500000)),
        (Emu(457200), Emu(3700000), Emu(5600000), Emu(2500000)),
        (Emu(6200000), Emu(3700000), Emu(5600000), Emu(2500000)),
    ]
    for (title, items), pos in zip(cards, positions):
        _add_card(slide, pos[0], pos[1], pos[2], pos[3], title, items)

def build_dashboard_roles(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "看板总览 & 角色定制", "Dashboard & Role Customization")
    roles = [
        ("管理员", "SLA 达成率 → 文控复审告警 → 8D 进度 → 知识库新增 → 低库存备件 → 安全检查到期 → 系统健康"),
        ("工程师", "我的工单 → 待验证工单 → 到期 PM → 点检 → 8D 进行中 → 相似故障推荐"),
        ("QA", "审核中文档 → 复审到期 → 附加修正审批 → 8D 报告 → 安全检查 → 表单审核清单"),
        ("操作员", "设备状态总览 → 我的工单 → 今日点检 → 润滑到期 → 常用设备快速入口"),
    ]
    y = Emu(1100000)
    for role, cards in roles:
        _add_rect(slide, Emu(457200), y, Emu(11277000), Emu(914400), LIGHT_GRAY)
        _add_text(slide, Emu(640080), y + Emu(137160), Emu(1500000), Emu(640080), role, 14, DARK_BLUE, bold=True)
        _add_text(slide, Emu(2200000), y + Emu(137160), Emu(9400000), Emu(640080), cards, 12, DARK_TEXT)
        y += Emu(1000000)
    _add_text(slide, Emu(457200), Emu(5400000), Emu(11277000), Emu(457200),
              "支持：明色青绿 / 暗色霓虹 / 跟随系统 三套模式一键切换；暗色模式做了亮度压暗专项补丁，夜班大屏久看不累。",
              13, TEAL, bold=True)

def build_device_mgmt(prs, screenshot=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "① 设备管理 & DOWN→工单自动派发", "Equipment Management")
    if screenshot:
        slide.shapes.add_picture(io.BytesIO(screenshot),
            Emu(365760), Emu(1280000), Emu(5600000), Emu(4800000))
        panel_left = Emu(6200000); panel_w = Emu(5600000)
    else:
        panel_left = Emu(457200); panel_w = Emu(11277000)
    _add_card(slide, panel_left, Emu(1280000), panel_w, Emu(1400000),
              "设备台账 & 状态切换",
              ["厂区/区域分类、附件上传（说明书/图纸/SOP）",
               "任意角色切 DOWN → 自动创建 REPAIR 工单 → 自动派工程师",
               "状态：RUN / IDLE / DOWN / PM / ENGINEERING / PROCESS_VALIDATION / OFFLINE"])
    _add_card(slide, panel_left, Emu(2800000), panel_w, Emu(1400000),
              "设备全生命周期 T0-T3",
              ["T0 选型(URS/候选供应商) → T1 采购(PO/金额/交付)",
               "T2 安装调试(FAT/SAT) → T3 量产移交(验收结果)",
               "时间线视图一眼看全流程进度"])
    _add_card(slide, panel_left, Emu(4320000), panel_w, Emu(1400000),
              "润滑管理（五定）",
              ["定点/定人/定时/定质/定量",
               "自动推算下次润滑日，到期提醒高亮"])

def build_daily_ops(prs, screenshot=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "② 日常运维闭环 & 工单 SLA", "Daily Operations & Work Order SLA")
    if screenshot:
        slide.shapes.add_picture(io.BytesIO(screenshot),
            Emu(365760), Emu(1280000), Emu(5600000), Emu(4800000))
        panel_left = Emu(6200000); panel_w = Emu(5600000)
    else:
        panel_left = Emu(457200); panel_w = Emu(11277000)
    _add_card(slide, panel_left, Emu(1280000), panel_w, Emu(1300000),
              "工单管理",
              ["创建/跟踪、紧急度标签、持续时长实时显示、关键词全文检索",
               "详情：5Why 根因分析 + 备件领用 + 状态流转(OPEN→CLOSED)"])
    _add_card(slide, panel_left, Emu(2700000), panel_w, Emu(1300000),
              "SLA & 超期升级",
              ["SLA 目标响应/解决时长（按紧急度预设）",
               "实际时长自动计算 → SLA 达成率统计 → 超期自动升级指派"])
    _add_card(slide, panel_left, Emu(4120000), panel_w, Emu(1600000),
              "PM / 点检 / 备件",
              ["PM：周/双周/月/季 周期，批量生成到期 PM 工单",
               "点检：模板+检查项+历史记录",
               "备件：库存/出入库/低库存告警/设备易损件绑定"])

def build_process_control(prs, screenshot=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "③ 工艺文控闭环", "Process Document Control")
    if screenshot:
        slide.shapes.add_picture(io.BytesIO(screenshot),
            Emu(365760), Emu(1280000), Emu(5600000), Emu(4800000))
        panel_left = Emu(6200000); panel_w = Emu(5600000)
    else:
        panel_left = Emu(457200); panel_w = Emu(11277000)
    _add_card(slide, panel_left, Emu(1280000), panel_w, Emu(1300000),
              "工艺文件 & 结构化电子表单",
              ["指导性文件：版本管理，生效版唯一，同组旧版自动作废",
               "管理员定义模板（9 种字段类型）→ 操作员动态填写 → JSON/CSV 导出"])
    _add_card(slide, panel_left, Emu(2700000), panel_w, Emu(1500000),
              "文控系统（符合体系标准）",
              ["三级电子签名审批链：编制→审核→批准（二次密码校验 + SHA256 指纹）",
               "状态机白名单：草稿↔审核中→生效→作废",
               "修订记录：字段级 before/after 对比"])
    _add_card(slide, panel_left, Emu(4320000), panel_w, Emu(1400000),
              "受控管理 & 复审告警",
              ["分发收回台账、PDF 自动加盖受控章水印",
               "表单审核锁定 + 附加修正流程",
               "复审周期告警：30 天内到期 / 已过期 Badge"])

def build_quality(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "④ 品管合规闭环", "Quality & Compliance")
    _add_card(slide, Emu(457200), Emu(1100000), Emu(5600000), Emu(2400000),
              "8D 报告",
              [("D1-D8 ", "8 步完整建模（团队→问题→遏制→根因→措施→验证→预防→表彰）"),
               ("关联工单 ", "（可选），状态：DRAFT / IN_PROGRESS / CLOSED"),
               ("⭐ 一键归档知识库 ", "：D0+D2→现象 / D4→根因 / D5→处置 / D7→预防")])
    _add_card(slide, Emu(6200000), Emu(1100000), Emu(5600000), Emu(2400000),
              "FMEA & 环境核查",
              [("FMEA ", "：失效模式 / 影响 / S×O×D → RPN，RPN>100 高亮"),
               ("环境核查 ", "：温度/湿度/洁净度/VOC/ESD 参数录入与趋势")])
    _add_card(slide, Emu(457200), Emu(3600000), Emu(11277000), Emu(2400000),
              "安全检查（四类）",
              ["safety_device 安全装置（防护罩/急停/联锁）/ 特种设备 / 环保 / 消防",
               "按频率自动推算下次检查日，30 天内黄底、已过期红字",
               "特种设备证书到期 Badge，检查发现 / 整改措施跟踪"])

def build_knowledge_base(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "⑤ 故障知识库（双路径归档）", "Knowledge Base (Dual-Path Archive)")
    y = Emu(1100000)
    paths = [
        ("工单", "知识库→「从工单归档」", "description→现象 / fault_category→分类 / root_cause→根因 / solution→处置 / prevention→预防"),
        ("8D 报告", "知识库→「从 8D 归档」或 8D 列表→「归档至知识库」", "D0+D2→现象 / D4→根因 / D5→处置 / D7→预防"),
        ("手动", "新建条目", "任意填写"),
    ]
    _add_rect(slide, Emu(457200), y, Emu(11277000), Emu(365760), CARD_BG)
    _add_text(slide, Emu(640080), y + Emu(45720), Emu(1400000), Emu(274320), "来源", 12, DARK_BLUE, bold=True)
    _add_text(slide, Emu(2100000), y + Emu(45720), Emu(4200000), Emu(274320), "入口", 12, DARK_BLUE, bold=True)
    _add_text(slide, Emu(6400000), y + Emu(45720), Emu(5200000), Emu(274320), "自动映射", 12, DARK_BLUE, bold=True)
    y += Emu(365760)
    for source, entry, mapping in paths:
        _add_rect(slide, Emu(457200), y, Emu(11277000), Emu(548640), LIGHT_GRAY)
        _add_text(slide, Emu(640080), y + Emu(91440), Emu(1400000), Emu(365760), source, 12, TEAL, bold=True)
        _add_text(slide, Emu(2100000), y + Emu(91440), Emu(4200000), Emu(365760), entry, 11, DARK_TEXT)
        _add_text(slide, Emu(6400000), y + Emu(91440), Emu(5200000), Emu(365760), mapping, 11, DARK_TEXT)
        y += Emu(548640)
    _add_card(slide, Emu(457200), Emu(3300000), Emu(5600000), Emu(2700000),
              "相似案例推荐",
              [("同设备 + 同故障分类 ", "优先 → 退化匹配"),
               ("按复发次数/浏览量倒序", ""),
               ("新故障一眼看到以前怎么修的", "")])
    _add_card(slide, Emu(6200000), Emu(3300000), Emu(5600000), Emu(2700000),
              "复发追踪",
              [("同一根因复发 → ", "点「标记复发」→ recurrence_count +1"),
               ("看板/列表一眼看出哪些根因反复出现", ""),
               ("需要重点改进的问题一目了然", "")])

def build_lcc_oee(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "⑤ 数据价值（续）：LCC & OEE", "LCC & OEE Analysis")
    _add_card(slide, Emu(457200), Emu(1100000), Emu(5600000), Emu(4800000),
              "设备成本 LCC（全生命周期成本）",
              [("6 种成本类型：", ""),
               ("  采购 / 维护 / 备件 / 能耗 / 折旧 / 报废", ""),
               ("单设备汇总：", "各类成本占比 + 总成本"),
               ("全设备汇总：", "按类型统计（饼图）+ Top10 高成本设备排名"),
               ("单设备年度趋势：", "按年汇总（折线图）")])
    _add_card(slide, Emu(6200000), Emu(1100000), Emu(5600000), Emu(4800000),
              "OEE 分析",
              [("设备综合效率 = ", "可用率 × 性能率 × 良品率"),
               ("按设备 / 时间维度聚合", ""),
               ("按日/周/月维度趋势图展示", ""),
               ("设备对比分析视图", ""),
               ("OEE 组成拆解，瓶颈定位", "")])

def build_security(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "安全 & 审计 & 灾备", "Security, Audit & Backup")
    _add_card(slide, Emu(457200), Emu(1100000), Emu(5600000), Emu(2400000),
              "安全加固（面向局域网部署）",
              [("JWT 双令牌", "：2h access + 7d refresh，5 分钟前自动续期"),
               ("bcrypt", "(SHA-256 预哈希，修复 72 字节截断)"),
               ("密码 3/4 类字符复杂度 / 失败锁定 5 次/15min", ""),
               ("首次登录强制改密 / 用户名枚举防护", ""),
               ("CSP / X-Frame / CORS 白名单", "")])
    _add_card(slide, Emu(6200000), Emu(1100000), Emu(5600000), Emu(2400000),
              "审计日志入库（审核留痕）",
              [("11 类动作全入库：", ""),
               ("  LOGIN_OK / FAIL / LOCKED / LOGOUT", ""),
               ("  PASSWORD_CHANGED / RESET", ""),
               ("  USER_CREATE / UPDATE / DELETE / UNLOCK", ""),
               ("  RESTORE_BACKUP", ""),
               ("每条含：", "action / actor / target / ip / UA / detail / time")])
    _add_card(slide, Emu(457200), Emu(3600000), Emu(11277000), Emu(2400000),
              "灾备（3-2-1 策略）",
              [("应用内定时备份 + SQLite 热快照 + ZIP + AES-256(Fernet/PBKDF2 20万轮) + NAS/SMB/U盘异地副本", ""),
               ("系统级旁路备份脚本（sh/bat，不依赖后端进程也能跑）", ""),
               ("备份后自动烟雾还原测试：解压→打开 sqlite→查表行数→校验 uploads，永远知道备份能不能真还原", "")])

def build_deployment(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "部署方案选型", "Deployment Options")
    y = Emu(1100000)
    headers = [("场景", 2800000), ("推荐部署", 3000000), ("自启", 1600000), ("崩溃重启", 1800000), ("健康检查", 2077000)]
    x = Emu(457200)
    _add_rect(slide, x, y, Emu(11277000), Emu(365760), CARD_BG)
    for h, w in headers:
        _add_text(slide, x + Emu(91440), y + Emu(45720), Emu(w), Emu(274320), h, 11, DARK_BLUE, bold=True)
        x += Emu(w)
    y += Emu(365760)
    rows = [
        ("工厂 Linux 有 root", "systemd（方案 A）", "✅ 开机自启", "✅ on-failure", "✅ cron 2min→重启"),
        ("普通用户 / 受限 Linux", "watchdog_user.sh（方案 B）", "✅ @reboot", "✅ 每分钟 tick", "✅ 同上"),
        ("Windows Server / 工控机", "NSSM + 任务计划（方案 C）", "✅ 延迟自启", "✅ NSSM + SC failure", "✅ 2min 计划任务"),
        ("容器化 / 快速体验", "Docker Compose（方案 D）", "✅ unless-stopped", "✅ compose restart", "✅ healthcheck"),
    ]
    for row in rows:
        x = Emu(457200)
        _add_rect(slide, x, y, Emu(11277000), Emu(548640), LIGHT_GRAY)
        for val, (_, w) in zip(row, headers):
            _add_text(slide, x + Emu(91440), y + Emu(91440), Emu(w), Emu(365760), val, 10, DARK_TEXT)
            x += Emu(w)
        y += Emu(548640)
    _add_text(slide, Emu(457200), y + Emu(91440), Emu(11277000), Emu(731520),
              "入口级加固：端口占用友好提示 / SQLite WAL checkpoint 钩子 / uvicorn 优雅退出 / 日志双写（journald + 按天滚动 14/30 天）",
              12, TEAL, bold=True)

def build_accounts(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "默认账号 & 演示数据", "Default Accounts & Demo Data")
    accounts = [
        ("admin", "管理员", "admin123", "全功能 / 用户与权限 / 备份与恢复"),
        ("engineer1", "张工 工程师", "eng123", "工单处置 / 8D / 知识库录入"),
        ("process1", "周工艺 工艺员", "proc123", "工艺文件 / 电子表单 / 提交文控审核"),
        ("qa1", "陈品管 QA", "qa123", "文控审核批准 / 8D / FMEA / 安全检查"),
        ("operator1", "王操作 操作员", "op123", "切换设备状态 / 点检 / 电子表单填写"),
        ("viewer1", "孙查看 查看者", "view123", "只读浏览（设备/工单/文档/知识）"),
    ]
    y = Emu(1100000)
    _add_rect(slide, Emu(457200), y, Emu(11277000), Emu(365760), CARD_BG)
    cols = [("账号", 2000000), ("角色", 2500000), ("默认密码", 2000000), ("典型工作内容", 4777000)]
    x = Emu(457200)
    for h, w in cols:
        _add_text(slide, x + Emu(91440), y + Emu(45720), Emu(w), Emu(274320), h, 11, DARK_BLUE, bold=True)
        x += Emu(w)
    y += Emu(365760)
    for acc, role, pwd, work in accounts:
        x = Emu(457200)
        _add_rect(slide, x, y, Emu(11277000), Emu(457200), LIGHT_GRAY)
        for val, (_, w) in zip([acc, role, pwd, work], cols):
            _add_text(slide, x + Emu(91440), y + Emu(91440), Emu(w), Emu(274320), val, 10, DARK_TEXT)
            x += Emu(w)
        y += Emu(457200)
    _add_text(slide, Emu(457200), y + Emu(182880), Emu(11277000), Emu(731520),
              "演示数据开箱即有：10 台设备 + 27 条状态日志 + 16 附件 + 14 备件 + 12 工单 + 4 份 8D 报告 + 14 工艺文件 + 12 安全检查 + 9 润滑点 + 9 条故障知识 + 31 成本记录 等 30+ 类记录。",
              12, TEAL, bold=True)

def build_tech_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "技术栈 & 扩展路线", "Tech Stack & Roadmap")
    _add_card(slide, Emu(457200), Emu(1100000), Emu(5600000), Emu(3500000),
              "当前技术栈",
              [("后端：", "Python 3.10+ / FastAPI / SQLAlchemy 2.0 / SQLite"),
               ("前端：", "Vue 3 / Vite / Element Plus / ECharts / Pinia"),
               ("安全：", "JWT / bcrypt(+SHA256 prehash) / CSP / CORS 白名单"),
               ("备份：", "SQLite hot backup / ZIP / Fernet AES-256"),
               ("打包：", "Docker / systemd / NSSM / PyInstaller")])
    _add_card(slide, Emu(6200000), Emu(1100000), Emu(5600000), Emu(1600000),
              "短期扩展路线（Q4 规划）",
              ["① 审计日志前端界面化（查询/导出 CSV）",
               "② 知识条目审核流（草稿→QA 审核→生效）",
               "③ 8D 报告通知（钉钉/飞书 Webhook）",
               "④ 备品低库存自动建采购申请单"])
    _add_card(slide, Emu(6200000), Emu(2900000), Emu(5600000), Emu(1700000),
              "中期扩展",
              ["对接 MES：设备状态 / 工艺参数自动上传",
               "对接 SCADA：关键机台参数趋势、阈值报警",
               "移动端（微信小程序/钉钉微应用）：点检 / 工单"])

def build_contact(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "文档索引 & 更新记录", "Documentation & Changelog")
    _add_card(slide, Emu(457200), Emu(1100000), Emu(5600000), Emu(2500000),
              "文档索引",
              [("README.md", " — 部署 & 架构详情 & 安全说明 & 灾备"),
               ("用户使用教程.md", " — 日常使用操作手册（按角色分步骤）"),
               ("SEMS_功能介绍_最新版.md", " — PPT 文字源（本文件）"),
               ("SEMS_UI设计评审与方案对比.pptx", " — UI 设计方案评审对比")])
    _add_card(slide, Emu(6200000), Emu(1100000), Emu(5600000), Emu(2500000),
              "更新记录",
              [("V2026.08：", "8D→知识库双归档、8 组分级菜单、角色看板、审计日志入库、bcrypt 截断修复"),
               ("V2026.06：", "文控系统(三级签名)、安全检查、润滑管理、工单 SLA、设备生命周期、知识库、LCC"),
               ("V2026.04：", "B+C 双皮肤主题、表单模板结构化电子表单、DOWN→工单自动派发")])
    _add_text(slide, Emu(457200), Emu(4400000), Emu(11277000), Emu(731520),
              "默认账号：admin / admin123（首次登录强制改密）",
              16, DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, Emu(457200), Emu(5200000), Emu(11277000), Emu(731520),
              "感谢观看！", 28, TEAL, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# NEW P17-P25: MANAGEMENT WORKFLOWS (not operational step-by-step)
# ════════════════════════════════════════════════════════════════════════════

def build_mgmt_overview(prs):
    """P17: 管理流程总览 · 8大业务管理闭环 × 角色分工矩阵"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "【管理 0】业务管理流程总览 · 8 大管理闭环 × 角色分工矩阵",
                 "Business Management Workflow Overview · 8 Loops × RACI Matrix")
    # Left: 8 大业务管理闭环列表（账号权限/备份/灾备等系统治理不属于业务流程，不纳入）
    loops = [
        ("①", "设备全生命周期管理", "选型→采购→验收→移交→运行→报废；状态切换/DOWN派工/润滑五定", "管理员/工程师"),
        ("②", "故障响应管理", "DOWN触发→派工→5Why处置→SLA升级→验证→关单→月度考核", "全员/工程师/QA"),
        ("③", "预防性维护管理", "PM计划审批→到期生成工单→执行→验收→漏做追责；点检异常联动", "管理员/工程师/操作"),
        ("④", "文控合规管理", "编制→审核→批准→分发→复审→作废；三级签名+受控水印+指纹", "工艺/QA/管理员"),
        ("⑤", "电子表单治理", "模板立项→字段配置→QA启用审批→填写→审核锁定→附加修正→导出", "管理员/工艺/QA/操作"),
        ("⑥", "品质改进管理", "8D团队+时限→FMEA RPN阈值处置→安全检查+证书→整改闭环", "QA/工程师/管理员"),
        ("⑦", "知识沉淀与复发治理", "双路径归档→QA审核发布→相似推荐→复发计数→根因根治立项", "QA/工程师/管理员"),
        ("⑧", "备件与库存管理", "申购→入库→领用绑定工单→低库存告警→盘点→报废；易损件绑定", "管理员/仓库/工程师"),
    ]
    y = Emu(1050000)
    _add_rect(slide, Emu(457200), y, Emu(7800000), Emu(365760), CARD_BG)
    headers = [("编号", 700000), ("管理闭环", 3300000), ("范围（业务闭环）", 3000000), ("主责角色", 800000)]
    x = Emu(457200)
    for (h, w) in headers:
        _add_text(slide, x + Emu(91440), y + Emu(45720), Emu(w), Emu(274320), h, 11, DARK_BLUE, bold=True)
        x += Emu(w)
    y += Emu(365760)
    colors_loop = [DARK_BLUE, RED, ORANGE, GREEN, ACCENT, PURPLE, TEAL, DARK_BLUE]
    for i, (num, name, scope, role) in enumerate(loops):
        _add_rect(slide, Emu(457200), y, Emu(7800000), Emu(550000), LIGHT_GRAY if i % 2 == 0 else WHITE)
        # 编号色块
        _add_rect(slide, Emu(457200), y, Emu(700000), Emu(550000), colors_loop[i])
        tf = slide.shapes[-1].text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = p.add_run(); _set_run(run, num, 13, WHITE, bold=True)
        _add_text(slide, Emu(1200000), y + Emu(120000), Emu(3300000), Emu(300000), name, 11, colors_loop[i], bold=True)
        _add_text(slide, Emu(4500000), y + Emu(100000), Emu(3000000), Emu(420000), scope, 9, DARK_TEXT)
        _add_text(slide, Emu(7500000), y + Emu(120000), Emu(750000), Emu(300000), role, 10, DARK_BLUE)
        y += Emu(550000)
    # Right: RACI summary card
    _add_card(slide, Emu(8400000), Emu(1050000), Emu(3350000), Emu(4800000),
              "🎯 角色 × 管理项（R=负责/A=审批/C=咨询/I=知会）",
              [("管理员 A：", "⑧全责，①②③④⑤⑥⑦ A/R 审批"),
               ("工程师 R：", "①②③⑦ R（执行主责），⑥ C，⑧ C 领用备件"),
               ("工艺员 R：", "④文件编制 R，⑤模板立项 C"),
               ("QA A/R：", "④文控审核批准 A/R，⑤表单审核 A，⑥品质改进 R，⑦知识审核 A"),
               ("操作员 R：", "②故障上报 R，③PM/点检执行 R，⑤表单填写 R，①状态切换 R"),
               ("查看者 I：", "全部 8 个业务模块只读浏览 I"),
               ("", ""),
               ("业务红线：", "文控批准 / 知识发布 / 备件报废 / 批量关单 — 必须管理员或 QA 双签才能生效")],
              title_color=PURPLE, item_size=10)
    # Bottom
    _add_rect(slide, Emu(457200), Emu(5750000), Emu(11277000), Emu(800000), LIGHT_GRAY)
    _add_text(slide, Emu(640080), Emu(5790000), Emu(11000000), Emu(365760),
              "📌 说明：以下 P18-P25 每页对应一个业务管理闭环，统一结构为「阶段泳道（责任角色/动作/交付物）→ 专项治理矩阵 → KPI考核指标 → 异常处理红线规则」（账号权限/备份/灾备等系统治理工作，不属于本业务管理流程范围）",
              12, DARK_BLUE, bold=True)

def build_mgmt_equipment_lifecycle(prs):
    """P18: 设备全生命周期管理流程"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "【管理 1】设备全生命周期管理流程（资产 + 状态 + 润滑）",
                 "Equipment Lifecycle Management (Asset · Status · Lubrication)")
    # 6 stages
    stages = [
        ("T0 选型立项", "管理员/采购", "URS撰写→候选供应商筛选→技术评审→预算审批", "URS文档 / 选型报告 / 预算签核", PURPLE),
        ("T1 采购签约", "管理员/采购", "PO下单→商务条款→交期约定→金额/成本登记", "采购合同 / PO单号 / 成本台账(采购类)", ACCENT),
        ("T2 安装调试", "工程师/供应商", "到货验收→FAT厂验→SAT现场→Commissioning→问题整改", "FAT/SAT报告 / 调试记录 / 整改清单", ORANGE),
        ("T3 量产移交", "QA/管理员", "最终验收→设备编号/台账入档→文件归档→责任人分配", "验收报告 / 设备档案 / 附件(说明书/图纸)", GREEN),
        ("T3+ 运行治理", "工程师/操作员", "状态7色实时切换→DOWN自动派工→润滑五定执行→点检记录", "状态日志 / 工单关联 / 润滑执行台账", DARK_BLUE),
        ("T4 退役报废", "管理员/财务", "LCC成本汇总→技术鉴定→资产报废→备件回收→状态OFFLINE", "报废审批 / 成本决算 / 残值处理台账", RED),
    ]
    _add_mgmt_stages(slide, Emu(457200), Emu(1050000), Emu(11277000), stages)
    # Lubrication governance inset
    y_mid = Emu(3800000)
    _add_rect(slide, Emu(457200), y_mid, Emu(11277000), Emu(1700000), CARD_BG)
    _add_text(slide, Emu(640080), y_mid + Emu(91440), Emu(11000000), Emu(300000),
              "🛢 润滑子流程（五定治理）：管理员定标准 → 工程师分配责任人 → 操作员按频执行 → 漏做高亮 → 逾期追责",
              12, DARK_BLUE, bold=True)
    lubrication_rows = [
        ("定点", "位置+点位编号→录入系统绑定设备", "工程师确认点位"),
        ("定人", "每点唯一责任人，离岗需交接", "交接清单管理员确认"),
        ("定时", "周期(日/周/月)→系统算next_date", "看板到期Badge黄/红"),
        ("定质", "油品/粘度/等级→写在点位说明", "换油需附MSDS附件"),
        ("定量", "用量(ml/g)→写在点位说明", "异常消耗进LCC维护类"),
    ]
    x_row = Emu(640080)
    col_w = Emu(2200000)
    for (k, v, chk) in lubrication_rows:
        _add_rect(slide, x_row, y_mid + Emu(440000), Emu(2100000), Emu(1150000), LIGHT_GRAY)
        _add_text(slide, x_row + Emu(91440), y_mid + Emu(480000), Emu(1900000), Emu(280000), f"▌{k}", 11, TEAL, bold=True)
        _add_text(slide, x_row + Emu(91440), y_mid + Emu(780000), Emu(1900000), Emu(340000), v, 9, DARK_TEXT)
        _add_text(slide, x_row + Emu(91440), y_mid + Emu(1140000), Emu(1900000), Emu(340000), f"✅ {chk}", 9, GREEN, bold=True)
        x_row += col_w
    # KPI + rules
    _add_mgmt_footer(slide,
        metrics_left=[
            "• DOWN机数/占比（按厂区/区域排名）",
            "• 状态切换合规率（原因必填率 ≥95%）",
            "• 润滑按期执行率 ≥98%（逾期追责到人）",
            "• 设备台账完整率（附件+易损件绑定率）",
        ],
        metrics_right=[
            "• 连续7日不更新状态 → 管理员看板红标",
            "• DOWN不填原因 → 审计日志标疑似漏填",
            "• 润滑漏做超3次/季度 → 责任人绩效扣项",
            "• 报废未走T4 → 资产账与系统账差异稽核",
        ],
        top=Emu(5750000))

def build_mgmt_fault_response(prs):
    """P19: 故障响应管理流程（SLA治理）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "【管理 2】故障响应管理流程（DOWN → SLA → 升级 → 考核）",
                 "Fault Response Management (DOWN → SLA → Escalation → KPI)")
    stages = [
        ("发现与上报", "操作员/全员", "发现异常→切DOWN→填现象→系统自动建REPAIR工单", "DOWN状态日志 / 自动工单 / 通知工程师", RED),
        ("派工与响应", "工程师/管理员", "接单(或自动派)→actual_start时间戳→现场确认", "响应时长 = actual_start - created_at", ORANGE),
        ("处置与5Why", "工程师", "备件领用→5Why根因→处置→预防→填解决方案", "5Why记录 / 备件出库 / 故障分类+根因", DARK_BLUE),
        ("SLA超时升级", "管理员/主管", "响应或解决超目标→一键升级→改派高级工程师", "升级记录(谁/何时/目标/原因) / 超期工单统计", PURPLE),
        ("验证与关闭", "工艺/操作员/QA", "转PENDING_VERIFY→现场验证→CLOSED或退回IN_PROGRESS", "验证人/时间 / 实际解决时长actual_end", ACCENT),
        ("复盘与考核", "管理员/QA", "SLA达成率统计→超时追责→复发次数→月度复盘", "月度SLA报告 / 工程师排名 / 改善项", TEAL),
    ]
    _add_mgmt_stages(slide, Emu(457200), Emu(1050000), Emu(11277000), stages)
    # SLA Matrix
    y_mid = Emu(3800000)
    _add_rect(slide, Emu(457200), y_mid, Emu(11277000), Emu(1750000), CARD_BG)
    _add_text(slide, Emu(640080), y_mid + Emu(91440), Emu(11000000), Emu(300000),
              "⏱ SLA 目标矩阵（按紧急度）· 超任一目标 = 超期工单（SLA违约）",
              12, DARK_BLUE, bold=True)
    matrix = [
        ("紧急度", "响应时长", "解决时长", "升级到"),
        ("URGENT 紧急", "≤ 15 分钟", "≤ 2 小时", "设备主管"),
        ("HIGH 高", "≤ 30 分钟", "≤ 4 小时", "主管工程师"),
        ("NORMAL 中", "≤ 1 小时", "≤ 8 小时", "工程师组长"),
        ("LOW 低", "≤ 2 小时", "≤ 24 小时", "不升级，记录原因"),
    ]
    col_w = [Emu(2500000), Emu(2500000), Emu(2500000), Emu(3500000)]
    x_h = Emu(640080)
    for (j, h) in enumerate(matrix[0]):
        _add_text(slide, x_h, y_mid + Emu(440000), col_w[j], Emu(280000), h, 11, TEAL, bold=True)
        x_h += col_w[j]
    colors_row = [RED, ORANGE, ACCENT, GRAY_TEXT]
    for i in range(1, 5):
        x_r = Emu(640080)
        y_r = y_mid + Emu(720000) + Emu(250000) * (i - 1)
        for j in range(4):
            w = col_w[j]
            _add_rect(slide, x_r, y_r, w - Emu(50000), Emu(230000), LIGHT_GRAY if i % 2 == 0 else WHITE)
            color = colors_row[i-1] if j == 0 else DARK_TEXT
            _add_text(slide, x_r + Emu(91440), y_r + Emu(45000), w - Emu(180000), Emu(200000), matrix[i][j], 10, color, bold=(j==0))
            x_r += w
    # Footer
    _add_mgmt_footer(slide,
        metrics_left=[
            "• 工单SLA达成率（按紧急度/按工程师/按月）",
            "• 平均响应时长(MTTR) / 平均解决时长(MTBF分母)",
            "• 超期工单数量 / 占比（月度趋势下降）",
            "• 5Why填写率 / 根因空值率（目标 0%）",
        ],
        metrics_right=[
            "• 30分钟无人接单 → 管理员看板通知 + 邮件",
            "• 超SLA未升级 → 管理员操作合规扣分",
            "• 根因空关单 → 退回补填，QA复核",
            "• 同一设备同根因复发≥3 → 强制开8D",
        ],
        top=Emu(5700000))

def build_mgmt_pm_inspection(prs):
    """P20: 预防性维护管理流程"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "【管理 3】预防性维护管理流程（PM计划 + 点检巡检 + 漏做追责）",
                 "Preventive Maintenance Management (PM Plan · Inspection · Trace)")
    stages = [
        ("PM计划编制", "工程师/管理员", "按设备重要度分级→定周期(周/双周/月/季/半年/年)→列维护内容清单→审批", "PM计划清单 / 周期标准 / 备件预需求", DARK_BLUE),
        ("点检模板治理", "QA/工程师", "定义检查项(9种类型)+标准值+上下限→启用审批→变更留痕", "点检模板库 / 标准值变更记录", ACCENT),
        ("到期生成工单", "管理员/PM责任人", "系统检测next_due→批量生成PM类型工单→自动派给计划责任人", "PM工单 / 到期Badge看板", ORANGE),
        ("执行与记录", "工程师/操作员", "按PM清单/点检模板逐项执行→备件消耗→异常项登记→可直接触发出工单", "执行记录 / 异常报告 / 关联REPAIR工单", PURPLE),
        ("验收与关闭", "QA/管理员", "抽检比例≥10%→照片核对→签字→关单；漏填/造假→退回重填", "验收记录 / 抽检结果台账", GREEN),
        ("分析与改进", "管理员", "PM完成率/点检异常率→高异常设备→升级PM频度或进8D", "月度PM分析 / 频度调整记录", TEAL),
    ]
    _add_mgmt_stages(slide, Emu(457200), Emu(1050000), Emu(11277000), stages)
    # 漏做追责矩阵
    y_mid = Emu(3800000)
    _add_rect(slide, Emu(457200), y_mid, Emu(5600000), Emu(1750000), CARD_BG)
    _add_text(slide, Emu(640080), y_mid + Emu(91440), Emu(5400000), Emu(300000),
              "🔴 漏做 / 逾期 分级处置",
              12, RED, bold=True)
    penalty = [
        ("逾期 < 1天", "责任人口头提醒，当日完成"),
        ("逾期 1-3 天", "工程师主管书面警告 + 补做"),
        ("逾期 > 3 天", "管理员介入，绩效扣分 + 设备风险评估"),
        ("漏做 ≥ 3 次/季度", "责任人年度绩效评级下调 1 档"),
        ("造假记录", "一经查实 → 直接重大违规处理"),
    ]
    _add_bullets(slide, Emu(640080), y_mid + Emu(420000), Emu(5400000), Emu(1300000),
                 [("★ " + t + ": ", d) for (t, d) in penalty], size=10, spacing=4)
    # 异常联动
    _add_rect(slide, Emu(6200000), y_mid, Emu(5550000), Emu(1750000), LIGHT_GRAY)
    _add_text(slide, Emu(6350000), y_mid + Emu(91440), Emu(5400000), Emu(300000),
              "🔗 点检异常 → 联动处置链",
              12, DARK_BLUE, bold=True)
    links = [
        ("① 发现异常项", "操作员点检时勾选异常 + 备注描述"),
        ("② 自动建议派工", "系统检测异常→提示工程师「是否创建REPAIR工单」"),
        ("③ 严重异常 → 设备切 DOWN", "关键项异常(安全类)→强制先切DOWN→再派工"),
        ("④ 追溯漏检", "设备DOWN后→自动查该设备近30天点检执行率"),
        ("⑤ 改进模板", "异常高发项→QA更新模板(增项/加严上下限)"),
    ]
    _add_bullets(slide, Emu(6350000), y_mid + Emu(420000), Emu(5400000), Emu(1300000),
                 [("▸ " + t + ": ", d) for (t, d) in links], size=10, spacing=4)
    _add_mgmt_footer(slide,
        metrics_left=[
            "• PM按期完成率 ≥95%（季度滚动）",
            "• 点检执行率（按人/按设备）≥98%",
            "• 点检异常发现率 / 异常→工单转化率",
            "• 抽检合格率（QA/管理员抽检）≥90%",
        ],
        metrics_right=[
            "• 计划到期未生成 → 管理员系统健康告警",
            "• 异常项未跟进处置 → QA月度品管会议议题",
            "• 模板标准值≥半年未更新 → 复审提醒",
            "• DOWN机追溯发现漏检 → 追加责任人处分",
        ],
        top=Emu(5700000))

def build_mgmt_doc_control(prs):
    """P21: 文控合规管理流程"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "【管理 4】文控合规管理流程（三级审批 + 复审 + 分发 + 水印 + 修正）",
                 "Document Compliance Management (3-Level Approval · Review · Distribute)")
    stages = [
        ("编制与提交", "工艺员/工程师", "编号规则预登记→文件上传→填分类/复审周期→二次密码提交审核→SHA256指纹", "草稿→审核中 / 编制签名 / 指纹记录", ACCENT),
        ("审核（合规性）", "QA审核员", "合规性/完整性/格式→通过或驳回(必填原因)→二次密码→指纹", "审核签名 / 驳回原因台账 / 同文档不可重审", ORANGE),
        ("批准（生效）", "QA主管/管理员", "二次密码批准→状态生效→生效日期+next_review_date→同组旧版自动作废", "批准签名 / 生效日期 / 作废旧版台账", GREEN),
        ("受控分发收回", "文控/管理员", "按USER/DEPT分发→份数登记→下载PDF自动盖受控章(编号/状态/下载人/日期)→收回", "分发台账 / 收回台账 / 水印PDF下载日志", PURPLE),
        ("复审与修订", "QA/工艺员", "30天内到期黄Badge→已过期红Badge→复审通过或作废→修订字段before/after", "复审记录 / 修订记录对比表 / 过期Badge清单", DARK_BLUE),
        ("作废与归档", "管理员/QA", "批准作废→状态作废(终态不可逆)→关联电子表单同步标记→PDF归档保留", "作废审批 / 终态归档 / 关联表单锁定", RED),
    ]
    _add_mgmt_stages(slide, Emu(457200), Emu(1050000), Emu(11277000), stages)
    # 三级签名规则
    y_mid = Emu(3800000)
    _add_rect(slide, Emu(457200), y_mid, Emu(5600000), Emu(1750000), CARD_BG)
    _add_text(slide, Emu(640080), y_mid + Emu(91440), Emu(5400000), Emu(300000),
              "🔏 三级电子签名 — 合规留痕要点",
              12, DARK_BLUE, bold=True)
    sig_rules = [
        ("每级必须二次密码校验", "防代签；密码校验失败不可签署"),
        ("签署意见可选+留痕", "系统记录空意见或具体意见"),
        ("SHA256指纹自动计算", "对(文档ID+版本+签署人+时间+密码哈希)算指纹"),
        ("状态机白名单单向流转", "非白名单跳转直接被后端拒绝"),
        ("同文档不可重复审核/批准", "已签角色不能再签，防绕过"),
        ("驳回必须填原因", "原因写进修订记录并通知编制人"),
    ]
    _add_bullets(slide, Emu(640080), y_mid + Emu(420000), Emu(5400000), Emu(1300000),
                 [("▸ " + t + ": ", d) for (t, d) in sig_rules], size=10, spacing=4)
    # 受控水印
    _add_rect(slide, Emu(6200000), y_mid, Emu(5550000), Emu(1750000), LIGHT_GRAY)
    _add_text(slide, Emu(6350000), y_mid + Emu(91440), Emu(5400000), Emu(300000),
              "🖨 受控章水印 + 复审告警矩阵",
              12, TEAL, bold=True)
    watermark = [
        ("PDF 必盖 4 要素", "文档编号 / 状态 / 下载人 / 下载日期时间"),
        ("非PDF不加水印但留痕", "下载日志入库：谁/何时/哪个文件/版本"),
        ("复审黄标", "next_review_date 在 30 天内 → 黄 Badge + 看板"),
        ("复审红标", "next_review_date 已过期 → 红 Badge + 看板置顶"),
        ("作废文件", "PDF下载额外加盖「作废」红章；操作员列表默认隐藏"),
        ("分发收回对账", "月度文控会议：分发数 vs 收回数 vs 在岗人数核对"),
    ]
    _add_bullets(slide, Emu(6350000), y_mid + Emu(420000), Emu(5400000), Emu(1300000),
                 [("▸ " + t + ": ", d) for (t, d) in watermark], size=10, spacing=4)
    _add_mgmt_footer(slide,
        metrics_left=[
            "• 文控在审/已生效/已作废文档数量",
            "• 复审到期率（当月到期 / 已过期数）",
            "• 平均审批时长（编制→批准总天数）",
            "• 驳回率（按文档类型 / 按编制人）",
        ],
        metrics_right=[
            "• 红标过期文档 ≥ 3天未处理 → QA月度议题",
            "• 无指纹签署 → 审计日志异常 → 安全追溯",
            "• 审批跳过QA直接批准 → 管理员权限核查",
            "• 作废文件恢复（非法操作）→ 审计标红",
        ],
        top=Emu(5700000))

def build_mgmt_forms(prs):
    """P22: 电子表单治理流程"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "【管理 5】电子表单治理流程（模板定义 → 审核锁 → 附加修正 → 导出）",
                 "E-Form Governance (Template · Lock · Amendment · Export)")
    stages = [
        ("模板立项", "工艺/QA/管理员", "业务部门提交需求→分析适用机台+字段→归类(作业记录类/通用类)", "模板需求单 / 字段清单", DARK_BLUE),
        ("模板配置", "管理员", "9种字段类型(Key+显示名+必填/选项/单位/上下限/排序)→上传参考空白模板", "模板配置 / 参考文件 / Key命名规范", ACCENT),
        ("启用审批", "QA", "QA审核模板合规→字段标准一致→通过启用；不合规→退回修改", "启用签名 / 变更记录", ORANGE),
        ("填写提交", "操作员", "选模板+机台→批次/班次/生产日期→逐字段→保存草稿或提交", "电子表单记录 / 结构化JSON", PURPLE),
        ("审核锁定", "QA/管理员", "文控审核通过→状态已审核→禁止原地修改；不通过→退回修改", "审核签名 / 锁定状态 / 审核清单", GREEN),
        ("附加修正 / 导出", "QA/管理员/全员", "已审核→修正审批(PENDING→QA审批APPROVED)；按模板导出JSON/CSV做分析", "修正台账 / JSON·CSV导出日志", TEAL),
    ]
    _add_mgmt_stages(slide, Emu(457200), Emu(1050000), Emu(11277000), stages)
    # 附加修正流程
    y_mid = Emu(3800000)
    _add_rect(slide, Emu(457200), y_mid, Emu(5600000), Emu(1750000), CARD_BG)
    _add_text(slide, Emu(640080), y_mid + Emu(91440), Emu(5400000), Emu(300000),
              "🧩 附加修正审批链（已审核表单不可原地修改）",
              12, PURPLE, bold=True)
    amend = [
        ("① 提出修正", "原填写人/工艺员 → 选字段+填原值+修正值+修正原因"),
        ("② 二次密码校验", "提出人必须再次输入登录密码，防越权代提"),
        ("③ 状态 PENDING", "修正项暂不生效；表单显示「修正待审批」Badge"),
        ("④ QA 审批", "QA/管理员 → 同意(APPROVED)或拒绝(REJECTED)→必须填意见"),
        ("⑤ 生效或退回", "APPROVED→修正值写入原值字段旁+保留原值；REJECTED→不生效"),
        ("⑥ 追溯", "修正前后值+原因+提出人+审批人+时间全留痕"),
    ]
    _add_bullets(slide, Emu(640080), y_mid + Emu(420000), Emu(5400000), Emu(1300000),
                 [("▸ " + t + ": ", d) for (t, d) in amend], size=10, spacing=4)
    # 模板治理
    _add_rect(slide, Emu(6200000), y_mid, Emu(5550000), Emu(1750000), LIGHT_GRAY)
    _add_text(slide, Emu(6350000), y_mid + Emu(91440), Emu(5400000), Emu(300000),
              "📐 模板治理规范",
              12, DARK_BLUE, bold=True)
    template_rules = [
        ("Key命名规范", "snake_case小写英文，同一模板内唯一；含中文→后端拒绝"),
        ("字段必填率", "关键参数(温度/压力/批次)必须必填；避免全可选"),
        ("上下限设置", "数值类必须配合理上下限，超出时提交警告并记异常"),
        ("参考模板对照", "作业记录类必须上传空白PDF/Excel参考件"),
        ("版本兼容性", "模板启用后修改字段→旧表单仍以旧版本结构导出"),
        ("停用与归档", "模板停用→不可新建；历史表单保留可读可导出"),
    ]
    _add_bullets(slide, Emu(6350000), y_mid + Emu(420000), Emu(5400000), Emu(1300000),
                 [("▸ " + t + ": ", d) for (t, d) in template_rules], size=10, spacing=4)
    _add_mgmt_footer(slide,
        metrics_left=[
            "• 按模板统计：填写量 / 异常率 / 超期未提交数",
            "• 已审核表单占比（目标 ≥85%）",
            "• 附加修正通过率 / 平均修正字段数",
            "• 结构化导出频次（工艺分析/质量分析）",
        ],
        metrics_right=[
            "• 关键数值超上下限 → 看板异常高亮",
            "• 模板字段Key非法 → 管理员启用前必须整改",
            "• 无审批直接改已审核表单 → 审计异常标红",
            "• 连续3个月未使用模板 → QA复核是否停用",
        ],
        top=Emu(5700000))

def build_mgmt_quality(prs):
    """P23: 品质改进管理流程"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "【管理 6】品质改进管理流程（8D + FMEA RPN + 安全检查 + 整改闭环）",
                 "Quality Improvement Management (8D · FMEA RPN · Safety · Closure)")
    stages = [
        ("触发与立项", "QA/管理员", "重大停机/同根因复发≥3/客户投诉→启动8D→D1组建团队(D1)+指定团长", "8D立项 / 团队名单 / 关联工单/设备", RED),
        ("问题与遏制", "8D团队", "D0问题描述+D2精确定义(5W2H)→D3临时遏制措施(短期围堵)→遏制责任人+期限", "D0-D3记录 / 遏制追踪表", ORANGE),
        ("根因与措施", "8D团队", "D4根因分析(5Why/鱼骨图/DOE)→D5永久纠正措施→负责人+计划完成日", "D4-D5记录 / 根因分类统计", ACCENT),
        ("验证与预防", "8D团队/QA", "D6实施+效果验证(数据对比)→D7预防再发(标准化/更新SOP/FMEA)→横向展开", "D6-D7记录 / 验证数据 / 预防展开清单", PURPLE),
        ("FMEA治理", "QA/工程师", "S/O/D打分→RPN=S×O×D→RPN>100高亮→必须改进行动→季度重评", "FMEA台账 / RPN分布 / 改进计划", DARK_BLUE),
        ("安全/环境/证书", "QA/安全员", "4类安全检查+环境参数+证书到期→整改→复查→关闭；过期红Badge→看板", "检查记录 / 整改台账 / 证书到期清单", GREEN),
    ]
    _add_mgmt_stages(slide, Emu(457200), Emu(1050000), Emu(11277000), stages)
    # 8D 时限标准
    y_mid = Emu(3800000)
    _add_rect(slide, Emu(457200), y_mid, Emu(5600000), Emu(1750000), CARD_BG)
    _add_text(slide, Emu(640080), y_mid + Emu(91440), Emu(5400000), Emu(300000),
              "⏳ 8D 各阶段时限标准（目标值·可调）",
              12, DARK_BLUE, bold=True)
    ddl = [
        ("D0-D3 (遏制)", "立项后 24 小时内完成围堵措施"),
        ("D4 (根因)", "立项后 3 工作日内输出根因初稿"),
        ("D5 (永久措施)", "D4确认后 5 工作日内批准措施方案"),
        ("D6 (验证)", "措施实施后 10 工作日内给出验证数据"),
        ("D7 (预防)", "验证通过后 5 工作日内完成SOP/FMEA更新"),
        ("D8 (总结关闭)", "QA复核通过后 2 工作日内表彰+关闭"),
    ]
    _add_bullets(slide, Emu(640080), y_mid + Emu(420000), Emu(5400000), Emu(1300000),
                 [("▸ " + t + ": ", d) for (t, d) in ddl], size=10, spacing=4)
    # RPN阈值
    _add_rect(slide, Emu(6200000), y_mid, Emu(5550000), Emu(1750000), LIGHT_GRAY)
    _add_text(slide, Emu(6350000), y_mid + Emu(91440), Emu(5400000), Emu(300000),
              "🎯 FMEA RPN 分级处置 + 安全 4 类检查矩阵",
              12, RED, bold=True)
    rpn_safety = [
        ("RPN < 50", "低风险 → 常规记录，年度复核"),
        ("50 ≤ RPN < 100", "中风险 → 半年内改进计划 + 责任人"),
        ("RPN ≥ 100", "高风险 → 30天内改进行动 + QA跟踪"),
        ("4类检查(装置/特种设备/环保/消防)", "按频率(next_date)自动提醒；黄/红Badge"),
        ("特种设备证书到期", "30天黄标；过期红标+看板置顶+整改单"),
        ("环境超标(VOC/温湿度/ESD)", "超阈值→自动记异常+生成整改追踪"),
    ]
    _add_bullets(slide, Emu(6350000), y_mid + Emu(420000), Emu(5400000), Emu(1300000),
                 [("▸ " + t + ": ", d) for (t, d) in rpn_safety], size=10, spacing=4)
    _add_mgmt_footer(slide,
        metrics_left=[
            "• 8D按期关闭率（按上述时限）≥85%",
            "• RPN≥100项目数 / 平均RPN（季度下降）",
            "• 安全检查按期率 + 证书过期数（目标 0）",
            "• 整改完成率 / 平均整改天数（趋势）",
        ],
        metrics_right=[
            "• 8D超期未推进 → 团长升级主管+看板红标",
            "• RPN≥100无改进 → QA月度品质会议必提",
            "• 证书过期仍在使用 → 立即暂停+审计追责",
            "• 复发≥3未开8D → 管理员强制发起+团长问责",
        ],
        top=Emu(5700000))

def build_mgmt_knowledge(prs):
    """P24: 知识沉淀与复发治理流程"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "【管理 7】知识沉淀与复发治理流程（双归档 + 审核 + 复发 + 根因根治）",
                 "Knowledge Management (Archive · Audit · Recurrence · Root Cause Fix)")
    stages = [
        ("沉淀触发", "QA/工程师/管理员", "工单CLOSED(根因完整) / 8D CLOSED / 手动创建 → 触发归档动作", "触发来源记录 / 关联工单/8D ID", ORANGE),
        ("双路径归档", "QA/工程师", "工单→知识(描述→现象/根因→根因)；8D→知识(D0+D2→现象/D4→根因/D5→处置/D7→预防)", "草稿知识条目 / 自动映射字段", ACCENT),
        ("知识审核发布", "QA", "草稿→QA审核(分类/标签是否准确/根因是否完整)→发布或驳回；默认打标签8D/工单编号", "已发布知识 / 审核签名 / 驳回原因", GREEN),
        ("检索与推荐", "全员", "关键词全文(标题/现象/根因/处置)；优先同设备+同分类→退化；按浏览/复发排序", "浏览记录 / views统计 / 推荐命中", TEAL),
        ("复发标记计数", "工程师/QA", "确认同根因再次出现→二次确认→recurrence_count+1→自动关联本次工单", "复发台账 / 复发次数排名", RED),
        ("根因根治项目", "管理员/QA", "复发≥3的条目→强制立项专项改善→更新预防措施→更新知识条目", "专项立项 / 预防措施修订 / 知识更新记录", PURPLE),
    ]
    _add_mgmt_stages(slide, Emu(457200), Emu(1050000), Emu(11277000), stages)
    # 双路径归档字段映射
    y_mid = Emu(3800000)
    _add_rect(slide, Emu(457200), y_mid, Emu(5600000), Emu(1750000), CARD_BG)
    _add_text(slide, Emu(640080), y_mid + Emu(91440), Emu(5400000), Emu(300000),
              "📥 双路径自动字段映射表（不可篡改留痕）",
              12, DARK_BLUE, bold=True)
    map_rows = [
        ("知识字段", "← 工单归档 (source_work_order_id)", "← 8D归档 (source_d8_report_id)"),
        ("symptom 现象", "description + title", "D0(问题) + D2(问题定义)"),
        ("root_cause 根因", "root_cause 字段", "D4(根本原因)"),
        ("solution 处置", "solution / 处置措施", "D5(永久纠正措施)"),
        ("prevention 预防", "prevention / 预防措施", "D7(预防再发)"),
        ("fault_category", "fault_category 字段", "按设备/关键词自动推荐，QA可改"),
    ]
    for i, row in enumerate(map_rows):
        y_row = y_mid + Emu(420000) + Emu(220000) * i
        bg = CARD_BG if i == 0 else (LIGHT_GRAY if i % 2 == 0 else WHITE)
        if i == 0:
            color = TEAL; bold = True
        else:
            color = DARK_TEXT; bold = False
        ws = [Emu(1600000), Emu(1950000), Emu(1950000)]
        x_r = Emu(640080)
        for j, val in enumerate(row):
            w = ws[j]
            _add_rect(slide, x_r, y_row, w, Emu(200000), bg)
            _add_text(slide, x_r + Emu(45720), y_row + Emu(30000), w - Emu(91440), Emu(180000), val, 9, color, bold=bold)
            x_r += w
    # 复发治理
    _add_rect(slide, Emu(6200000), y_mid, Emu(5550000), Emu(1750000), LIGHT_GRAY)
    _add_text(slide, Emu(6350000), y_mid + Emu(91440), Emu(5400000), Emu(300000),
              "🔁 复发治理分级处置（recurrence_count 驱动）",
              12, RED, bold=True)
    recur = [
        ("count = 0 (首次)", "常规发布，工程师维修时可见相似案例"),
        ("count = 1 (再发1次)", "QA复核根因是否完整 / 预防是否到位"),
        ("count = 2 (再发2次)", "通知设备主管 → 预防措施升级审核"),
        ("count ≥ 3 (反复发)", "强制启动8D或专项改善 → 管理员月度议题"),
        ("复发TOP榜", "看板按复发次数倒序 → 一眼看清需重点根治的根因"),
        ("防误标机制", "标记复发必须二次确认 + 填写关联工单ID佐证"),
    ]
    _add_bullets(slide, Emu(6350000), y_mid + Emu(420000), Emu(5400000), Emu(1300000),
                 [("▸ " + t + ": ", d) for (t, d) in recur], size=10, spacing=4)
    _add_mgmt_footer(slide,
        metrics_left=[
            "• 知识库新增条目数（月度增长）",
            "• 工单/8D → 知识 转化率（目标 ≥70%）",
            "• 复发TOP10根因清单 / 复发总次数",
            "• 相似案例命中率（维修前浏览率）",
        ],
        metrics_right=[
            "• 有根因却未归档 → QA月度品质会议追责",
            "• 知识发布后30天无人浏览 → 复审是否冗余",
            "• 复发≥3未立项 → 管理员直接发起专项",
            "• 无关联工单/8D ID的复发标记 → 驳回",
        ],
        top=Emu(5700000))

def build_mgmt_spare_and_system(prs):
    """P25: 备件与库存管理流程（纯业务流程，不含账号权限/备份等系统治理内容）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "【管理 8】备件与库存管理流程（申购 · 入库 · 领用 · 低库存告警 · 盘点 · 报废）",
                 "Spare Parts Inventory Management (Procure · Stock · Issue · Alert · Audit · Scrap)")
    stages = [
        ("申购与审批", "工程师/管理员", "维修/PM预需求→填写申购单(品号/数量/预算/交期)→管理员或主管审批", "申购单 / 审批记录 / 采购预通知", DARK_BLUE),
        ("到货入库", "仓库/管理员", "到货核对→三单匹配(申购单/送货单/发票)→批次/有效期/供应商录入→入库上架", "入库台账 / 批次标签 / 库存(数量+位置)更新", ACCENT),
        ("领用与出库", "工程师/操作员", "工单详情→领用备件(选品号/数量)→自动扣减库存；无工单领用→管理员双签批准", "出库记录 / 关联工单ID / LCC备件类成本自动累加", ORANGE),
        ("低库存与补货", "管理员/仓库", "低于安全库存→看板标红Badge→建议补货量=历史月均×采购周期×安全系数→触发申购", "低库存清单 / 补货建议 / 采购周期看板", RED),
        ("易损件绑定", "工程师/管理员", "每台设备建立专属易损件清单→关联品号/平均寿命→DOWN时自动弹窗提示库存", "设备易损件绑定表 / DOWN机库存提示", PURPLE),
        ("盘点与报废", "管理员/财务", "季度实物盘点→差异审批→报废(过期/损坏/淘汰)→调账→报废值进LCC报废类", "盘点差异台账 / 报废审批 / LCC成本更新", GREEN),
    ]
    _add_mgmt_stages(slide, Emu(457200), Emu(1050000), Emu(11277000), stages)
    # 中部卡片：左-无工单出库治理；右-易损件与维修联动
    y_mid = Emu(3800000)
    _add_rect(slide, Emu(457200), y_mid, Emu(5600000), Emu(1750000), CARD_BG)
    _add_text(slide, Emu(640080), y_mid + Emu(91440), Emu(5400000), Emu(300000),
              "🚫 出库合规红线（防止备件流失）",
              12, RED, bold=True)
    issue_rules = [
        ("工单绑定（主渠道）", "95%以上出库必须从工单详情发起，自动记录关联工单ID"),
        ("无工单领用（例外）", "必须填写用途说明+管理员二次密码双签+异常台账标记"),
        ("数量异常", "单次领用超该品号月均2倍 → 自动弹窗要求主管审批"),
        ("追溯", "反向追溯：备件批次号→入库→出库→哪台设备哪张工单→哪个责任人"),
        ("防呆", "系统禁止零库存出库；禁止同一张工单领用同一品号3次以上"),
        ("月结", "月度出库汇总→按设备/按工程师排名→异常消耗Top10分析"),
    ]
    _add_bullets(slide, Emu(640080), y_mid + Emu(420000), Emu(5400000), Emu(1300000),
                 [("▸ " + t + ": ", d) for (t, d) in issue_rules], size=10, spacing=4)
    # 右：补货安全库存模型 + 易损件寿命联动
    _add_rect(slide, Emu(6200000), y_mid, Emu(5550000), Emu(1750000), LIGHT_GRAY)
    _add_text(slide, Emu(6350000), y_mid + Emu(91440), Emu(5400000), Emu(300000),
              "📐 安全库存模型 + 易损件寿命联动",
              12, DARK_BLUE, bold=True)
    stock_model = [
        ("安全库存公式", "SS = 月均消耗量 × (采购周期月 + 安全月)，默认安全系数 1.3"),
        ("采购周期分级", "本地供应商7天 / 外地30天 / 进口90天，周期越久SS越高"),
        ("ABC分类管理", "A类(价值高)严格按SS；C类(螺丝/标签)可宽松，避免缺件误修"),
        ("寿命预警1：PM", "易损件绑定设备后，累计PM次数≥寿命→自动建议下次更换"),
        ("寿命预警2：运行时", "备件已用次数/寿命 > 80% → 看板黄标；>100% → 红标"),
        ("DOWN机智能提示", "设备切DOWN→系统自动弹出：此设备易损件X当前库存=Y 是否领用？"),
    ]
    _add_bullets(slide, Emu(6350000), y_mid + Emu(420000), Emu(5400000), Emu(1300000),
                 [("▸ " + t + ": ", d) for (t, d) in stock_model], size=10, spacing=4)
    # KPI / Governance
    _add_mgmt_footer(slide,
        metrics_left=[
            "• 备件库存准确率（盘点差异率，目标 ≤2%）",
            "• 低库存告警响应率 / 平均补货周期（天）",
            "• 工单绑定出库率（目标 ≥95%）",
            "• 易损件DOWN机命中率（目标 ≥80%：DOWN时弹对的件）",
        ],
        metrics_right=[
            "• 无工单无双签出库 → 审计异常 + 责任人追责",
            "• 低库存超14天未补货 → 管理员月度设备会议必提",
            "• 盘点差异率≥5% → 复盘流程，管理员书面解释",
            "• 备件私自借出不入系统 → 一经查实按重大违规处理",
        ],
        top=Emu(5700000))


# ════════════════════════════════════════════════════════════════════════════
# Main
# ════════════════════════════════════════════════════════════════════════════
def main():
    orig_path = '/workspace/SEMS_功能介绍_old_backup.pptx'
    out_path = '/workspace/SEMS_功能介绍.pptx'
    screenshot = None
    import os
    for p in ['/workspace/SEMS_功能介绍.pptx', orig_path]:
        if os.path.exists(p):
            try:
                screenshot = _extract_screenshot(p)
                if screenshot: break
            except: pass

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Original 16 slides (feature intro)
    build_cover(prs)                          # P1
    build_why_sems(prs)                       # P2
    build_five_loops(prs)                     # P3
    build_v2026_updates(prs)                  # P4
    build_dashboard_roles(prs)                # P5
    build_device_mgmt(prs, screenshot)        # P6
    build_daily_ops(prs, screenshot)          # P7
    build_process_control(prs, screenshot)    # P8
    build_quality(prs)                        # P9
    build_knowledge_base(prs)                 # P10
    build_lcc_oee(prs)                        # P11
    build_security(prs)                       # P12
    build_deployment(prs)                     # P13
    build_accounts(prs)                       # P14
    build_tech_stack(prs)                     # P15
    build_contact(prs)                        # P16

    # NEW P17-P25: MANAGEMENT WORKFLOWS
    build_mgmt_overview(prs)                  # P17 总览+RACI矩阵
    build_mgmt_equipment_lifecycle(prs)       # P18 设备T0-T4+润滑
    build_mgmt_fault_response(prs)            # P19 DOWN→SLA→升级→考核
    build_mgmt_pm_inspection(prs)             # P20 PM+点检+漏做追责
    build_mgmt_doc_control(prs)               # P21 文控三级+复审+水印
    build_mgmt_forms(prs)                     # P22 电子表单+附加修正
    build_mgmt_quality(prs)                   # P23 8D时限+FMEA RPN+安全
    build_mgmt_knowledge(prs)                 # P24 双路径归档+复发治理
    build_mgmt_spare_and_system(prs)          # P25 备件+权限+审计+灾备

    prs.save(out_path)
    print(f"Generated {out_path} with {len(prs.slides)} slides "
          f"(V2026.08 feature intro P1-P16 + MANAGEMENT workflows P17-P25)")

if __name__ == '__main__':
    main()
