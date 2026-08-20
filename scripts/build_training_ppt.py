"""生成《生产系统从定义到使用》培训 PPT。

以半导体光刻工艺段为示范，覆盖：初始化 → 工艺设计 → 计划 → 执行 → 沉淀 五阶段。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- 主题色 ----------
COLOR_PRIMARY = RGBColor(0x1F, 0x4E, 0x79)      # 深蓝
COLOR_ACCENT = RGBColor(0x00, 0x78, 0xD4)        # 亮蓝
COLOR_DARK = RGBColor(0x26, 0x26, 0x26)          # 深灰文字
COLOR_GREY = RGBColor(0x59, 0x59, 0x59)          # 中灰
COLOR_LIGHT_BG = RGBColor(0xF2, 0xF6, 0xFB)      # 浅蓝背景
COLOR_CODE_BG = RGBColor(0xF5, 0xF5, 0xF5)        # 代码块灰底
COLOR_GREEN = RGBColor(0x2E, 0x7D, 0x32)         # 绿
COLOR_ORANGE = RGBColor(0xE6, 0x7E, 0x22)       # 橙
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ---------- 工具函数 ----------
def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # 移到最底层
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(slide, text, left, top, width, height, font_size=18, bold=False, color=COLOR_DARK, align=PP_ALIGN.LEFT, font_name="微软雅黑"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tb


def add_title_bar(slide, title_text, subtitle=None):
    """页眉：左上角彩色色块 + 标题文字"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()
    bar.shadow.inherit = False
    add_text(slide, title_text, Inches(0.4), Inches(0.18), Inches(11), Inches(0.6),
             font_size=26, bold=True, color=COLOR_WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.65), Inches(11), Inches(0.35),
                 font_size=12, color=RGBColor(0xCC, 0xDD, 0xEE))


def add_footer(slide, page_no, total):
    add_text(slide, f"SEMS 生产管理 · 培训材料   |   {page_no} / {total}",
             Inches(0.3), Inches(7.0), Inches(9.4), Inches(0.3),
             font_size=9, color=COLOR_GREY, align=PP_ALIGN.CENTER)


def add_code_block(slide, code_text, left, top, width, height, font_size=11):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_CODE_BG
    box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    box.line.width = Pt(0.5)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)
    # 多行
    lines = code_text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = "Consolas"
        p.font.color.rgb = COLOR_DARK
        p.line_spacing = 1.15
    return box


def add_bullets(slide, items, left, top, width, height, font_size=16, color=COLOR_DARK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "•  " + item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "微软雅黑"
        p.space_after = Pt(6)
        p.line_spacing = 1.25
    return tb


def add_stage_chip(slide, label, color, left, top, width=Inches(2.0), height=Inches(0.5)):
    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    chip.fill.solid()
    chip.fill.fore_color.rgb = color
    chip.line.fill.background()
    chip.shadow.inherit = False
    tf = chip.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = "微软雅黑"
    p.alignment = PP_ALIGN.CENTER
    return chip


def add_arrow(slide, left, top, width=Inches(0.4), height=Inches(0.5)):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLOR_GREY
    arrow.line.fill.background()
    arrow.shadow.inherit = False
    return arrow


# ---------- 构造 PPT ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TOTAL_PAGES = 14

# ============ Slide 1: 封面 ============
s = add_blank_slide(prs)
set_bg(s, COLOR_PRIMARY)
# 大色块
deco = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.5), prs.slide_width, Inches(2.0))
deco.fill.solid()
deco.fill.fore_color.rgb = COLOR_ACCENT
deco.line.fill.background()
deco.shadow.inherit = False

add_text(s, "SEMS 生产管理系统", Inches(0.8), Inches(1.8), Inches(12), Inches(0.8),
         font_size=20, color=RGBColor(0xCC, 0xDD, 0xEE))
add_text(s, "从定义到使用", Inches(0.8), Inches(2.6), Inches(12), Inches(1.6),
         font_size=54, bold=True, color=COLOR_WHITE)
add_text(s, "—— 以半导体光刻工艺段为示范 ——", Inches(0.8), Inches(4.2), Inches(12), Inches(0.6),
         font_size=22, color=RGBColor(0xFF, 0xD7, 0x00))
add_text(s, "工段库定义  /  工艺路由设计  /  生产订单计划  /  派工执行  /  数据沉淀",
         Inches(0.8), Inches(5.8), Inches(12), Inches(0.6),
         font_size=16, color=COLOR_WHITE)
add_text(s, "培训对象：工艺工程师 / 生产管理员 / 班组长 / 操作员",
         Inches(0.8), Inches(6.4), Inches(12), Inches(0.5),
         font_size=14, color=COLOR_WHITE)
add_footer(s, 1, TOTAL_PAGES)

# ============ Slide 2: 目录 ============
s = add_blank_slide(prs)
add_title_bar(s, "目录", "Agenda")
add_text(s, "本次培训将沿光刻工艺段从 0 到 1 走完全流程", Inches(0.4), Inches(1.2),
         Inches(12.5), Inches(0.5), font_size=14, color=COLOR_GREY)

agenda = [
    ("01", "总体流程概览", "5 个阶段、4 类角色、1 条数据主线", COLOR_PRIMARY),
    ("02", "阶段一：初始化基础数据", "设备台账 / 工艺数据采集表单模板", COLOR_ACCENT),
    ("03", "阶段二：工艺设计（核心）", "定义光刻工段 + 编排产品工艺路线", COLOR_GREEN),
    ("04", "阶段三：计划——生产订单", "MO 创建与下发", COLOR_ORANGE),
    ("05", "阶段四：执行——派工 + 自动初始化", "本系统最具差异化的环节", RGBColor(0xC2, 0x18, 0x5B)),
    ("06", "阶段五：填写工艺数据 + 报工", "操作员录入 + 状态联动", COLOR_PRIMARY),
    ("07", "全链路追溯能力 + 设计要点", "数据主线 + 关键设计决策", COLOR_ACCENT),
]
top = Inches(1.85)
for num, title, sub, color in agenda:
    chip = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), top, Inches(0.55), Inches(0.55))
    chip.fill.solid()
    chip.fill.fore_color.rgb = color
    chip.line.fill.background()
    chip.shadow.inherit = False
    tf = chip.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    add_text(s, title, Inches(1.35), top + Emu(20000), Inches(6.5), Inches(0.4),
             font_size=18, bold=True, color=COLOR_DARK)
    add_text(s, sub, Inches(1.35), top + Emu(380000), Inches(8), Inches(0.3),
             font_size=12, color=COLOR_GREY)
    top = top + Inches(0.74)

add_footer(s, 2, TOTAL_PAGES)

# ============ Slide 3: 总体流程图 ============
s = add_blank_slide(prs)
add_title_bar(s, "01  总体流程概览", "光刻工艺段的全生命周期")
add_text(s, "5 个阶段串成数据主线，各阶段由不同角色负责，确保工艺定义与生产执行解耦",
         Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.4), font_size=13, color=COLOR_GREY)

# 五个阶段色块
stages = [
    ("初始化", "管理员", "设备/产品/表单模板", COLOR_PRIMARY),
    ("工艺设计", "工艺工程师", "工段库 + 工序路由", COLOR_ACCENT),
    ("计划", "计划员", "生产订单 MO", COLOR_GREEN),
    ("执行", "班组长/操作员", "派工 + 填写数据", COLOR_ORANGE),
    ("沉淀", "QA/工艺", "审核 + 追溯分析", RGBColor(0xC2, 0x18, 0x5B)),
]
chip_w = Inches(2.2)
gap = Inches(0.18)
start_left = Inches(0.4)
top = Inches(1.85)
for i, (name, role, content, color) in enumerate(stages):
    left = start_left + i * (chip_w + gap)
    # 阶段块
    block = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, chip_w, Inches(1.6))
    block.fill.solid()
    block.fill.fore_color.rgb = color
    block.line.fill.background()
    block.shadow.inherit = False
    tf = block.text_frame
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "——"
    p2.font.size = Pt(10)
    p2.font.color.rgb = COLOR_WHITE
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = role
    p3.font.size = Pt(12)
    p3.font.color.rgb = COLOR_WHITE
    p3.alignment = PP_ALIGN.CENTER
    # 箭头
    if i < len(stages) - 1:
        add_arrow(s, left + chip_w, top + Inches(0.55), gap, Inches(0.5))

# 底部数据主线
add_text(s, "▼ 数据主线（一条不可断的链）",
         Inches(0.4), Inches(3.85), Inches(12.5), Inches(0.4),
         font_size=14, bold=True, color=COLOR_PRIMARY)

main_line = (
    "Equipment  →  FormTemplate  →  ProcessSection  →  Routing / RoutingStep  →  "
    "ProductionOrder  →  Dispatch  →  FormRecord  →  LaborReport"
)
add_code_block(s, main_line, Inches(0.4), Inches(4.3), Inches(12.5), Inches(0.7), font_size=12)

# 关键能力一览
add_text(s, "关键能力", Inches(0.4), Inches(5.3), Inches(12.5), Inches(0.4),
         font_size=14, bold=True, color=COLOR_PRIMARY)
caps = [
    "工段库独立可复用，一处定义全局生效",
    "派工创建时自动按工段模板初始化空白工艺数据表单",
    "派工状态与生产订单状态双向联动",
    "设备 DOWN 自动暂停派工 + 自动创建维修工单",
    "工艺记录审核后不可原地改，修正走留痕 Amendment",
]
add_bullets(s, caps, Inches(0.4), Inches(5.7), Inches(12.5), Inches(1.6), font_size=13)
add_footer(s, 3, TOTAL_PAGES)

# ============ Slide 4: 阶段一 - 设备台账 ============
s = add_blank_slide(prs)
add_title_bar(s, "02  阶段一：初始化基础数据（1/2）", "设备台账")
add_text(s, "光刻工段对应一组光刻机 + 配套显影机，归入同一设备组字符串",
         Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.4), font_size=13, color=COLOR_GREY)

# 设备组示例
add_text(s, "设备组 litho_group", Inches(0.4), Inches(1.7), Inches(5), Inches(0.4),
         font_size=15, bold=True, color=COLOR_ACCENT)
eq_code = """设备组：litho_group
  ├─ LITHO-001  ASML PAS5500  状态：RUNNING
  ├─ LITHO-002  ASML PAS5500  状态：RUNNING
  └─ DEV-001    东京精密显影机  状态：RUNNING"""
add_code_block(s, eq_code, Inches(0.4), Inches(2.15), Inches(6.3), Inches(2.0), font_size=12)

# 模板层 vs 执行层
add_text(s, "关键设计：设备组绑定在模板层，具体设备在执行层选",
         Inches(0.4), Inches(4.5), Inches(12.5), Inches(0.4),
         font_size=14, bold=True, color=COLOR_PRIMARY)

# 模板层/执行层对比卡
def add_layer_card(slide, title, content, color, left):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(5.0), Inches(6.0), Inches(1.8))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_LIGHT_BG
    card.line.color.rgb = color
    card.line.width = Pt(1.5)
    card.shadow.inherit = False
    add_text(slide, title, left + Inches(0.2), Inches(5.1), Inches(5.6), Inches(0.4),
             font_size=14, bold=True, color=color)
    add_text(slide, content, left + Inches(0.2), Inches(5.55), Inches(5.6), Inches(1.2),
             font_size=12, color=COLOR_DARK)

add_layer_card(s, "模板层（工段定义时）",
                "ProcessSection.equipment_group = \"litho_group\"\n同组设备共用同一套工艺参数采集模板，\n管理员不必为每台机器重复定义工段",
                COLOR_ACCENT, Inches(0.4))
add_layer_card(s, "执行层（派工时选）",
                "Dispatch.equipment_id = 25  (LITHO-001)\n生产人员根据当时哪台空闲、哪台资质可用，\n在组内选一台具体设备执行",
                COLOR_ORANGE, Inches(6.9))
add_footer(s, 4, TOTAL_PAGES)

# ============ Slide 5: 阶段一 - 工艺数据采集表单模板 ============
s = add_blank_slide(prs)
add_title_bar(s, "02  阶段一：初始化基础数据（2/2）", "工艺数据采集表单模板 FormTemplate")
add_text(s, "光刻工艺必须采集的关键参数（按 SPC 规范），通过 FormTemplate 定义字段结构",
         Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.4), font_size=13, color=COLOR_GREY)

form_json = """{
  "name": "光刻工艺参数记录表",
  "code": "FT-LITHO-01",
  "category": "record",
  "field_schema": [
    {"key":"exposure_energy","type":"number","label":"曝光能量(mJ/cm²)","required":true,"unit":"mJ/cm²","min":30,"max":80,"seq":1},
    {"key":"focus_offset",   "type":"number","label":"焦距偏移(μm)",    "required":true,"unit":"μm","min":-2,"max":2,"seq":2},
    {"key":"alignment_error","type":"number","label":"套刻误差(nm)",     "required":true,"unit":"nm","min":0,"max":50,"seq":3},
    {"key":"develop_time",  "type":"number","label":"显影时间(秒)",      "required":true,"unit":"s","default_value":60,"seq":4},
    {"key":"reticle_id",     "type":"text",  "label":"掩膜版编号",        "required":true,"seq":5},
    {"key":"abnormal",       "type":"radio", "label":"是否异常",
       "options":[{"label":"正常","value":"NORMAL"},{"label":"异常","value":"ABNORMAL"}],"seq":6}
  ]
}"""
add_code_block(s, form_json, Inches(0.4), Inches(1.6), Inches(8.5), Inches(5.3), font_size=10)

# 右侧说明
add_text(s, "字段类型支持", Inches(9.1), Inches(1.6), Inches(4), Inches(0.4),
         font_size=14, bold=True, color=COLOR_PRIMARY)
add_bullets(s, [
    "text 文本",
    "textarea 长文本",
    "number 数值（带 min/max 范围）",
    "select 下拉单选",
    "radio 单选按钮",
    "date / datetime / time",
    "boolean 布尔",
], Inches(9.1), Inches(2.05), Inches(4), Inches(2.5), font_size=12)

add_text(s, "校验能力", Inches(9.1), Inches(4.5), Inches(4), Inches(0.4),
         font_size=14, bold=True, color=COLOR_PRIMARY)
add_bullets(s, [
    "required 必填校验",
    "min/max 数值范围校验",
    "default_value 预填默认值",
    "options 选项约束",
    "提交时一次性校验",
], Inches(9.1), Inches(4.95), Inches(4), Inches(2.0), font_size=12)
add_footer(s, 5, TOTAL_PAGES)

# ============ Slide 6: 阶段二 - 定义光刻工段 ============
s = add_blank_slide(prs)
add_title_bar(s, "03  阶段二：工艺设计（1/2）", "定义光刻工段 ProcessSection")
add_text(s, "工段是管理员维护的可复用工艺单元，绑定设备组 + 关联工艺数据采集模板",
         Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.4), font_size=13, color=COLOR_GREY)

section_code = """POST /api/v1/process-sections
{
  "name": "光刻工段",
  "code": "SEC-LITHO",
  "equipment_group": "litho_group",        ← 模板层：允许这一组设备
  "form_template_id": 5,                    ← 关联"光刻工艺参数记录表"
  "standard_cycle_min": 45.0,              ← 标准工时 45 分钟/片
  "theoretical_uph": 60,                    ← 理论每小时 60 片
  "required_skill_level": "L3",             ← 需 L3 光刻操作资质
  "acceptance_criteria": "CD 偏差±2nm,套刻误差<35nm",
  "description": "光刻曝光+显影主工序"
}"""
add_code_block(s, section_code, Inches(0.4), Inches(1.65), Inches(12.5), Inches(3.2), font_size=11)

# 工段的复用价值
add_text(s, "工段库的核心价值", Inches(0.4), Inches(5.0), Inches(12.5), Inches(0.4),
         font_size=14, bold=True, color=COLOR_PRIMARY)
add_bullets(s, [
    "独立工段库：ProcessSection 主表，被多个产品的多道工序引用，一处改全局生效",
    "引用阻断硬删除：被 RoutingStep/Dispatch 引用时禁止 delete，只能 is_active=False 停用",
    "模板停用温和降级：关联的 FormTemplate 停用后，新派工不再自动初始化表单，但不阻断派工本身",
], Inches(0.4), Inches(5.4), Inches(12.5), Inches(1.7), font_size=12)
add_footer(s, 6, TOTAL_PAGES)

# ============ Slide 7: 阶段二 - 编排工艺路线 ============
s = add_blank_slide(prs)
add_title_bar(s, "03  阶段二：工艺设计（2/2）", "编排产品工艺路线 Routing")
add_text(s, "某 28nm 工艺产品定义 5 道工序，光刻是第 30 道，引用上面定义的工段",
         Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.4), font_size=13, color=COLOR_GREY)

routing_code = """POST /api/v1/routings
{
  "product_id": 12,                        ← 28nm 工艺产品
  "version": "v1.0",
  "steps": [
    {"seq": 10, "step_name": "氧化",     "process_section_id": 5},
    {"seq": 20, "step_name": "CVD沉积",  "process_section_id": 7},
    {"seq": 30, "step_name": "光刻",     "process_section_id": 9},   ← 引用光刻工段
    {"seq": 40, "step_name": "刻蚀",     "process_section_id": 11},
    {"seq": 50, "step_name": "去胶",     "process_section_id": 13}
  ]
}

POST /api/v1/routings/{id}/release        ← 路由必须生效才可用于生产订单
                                           （同产品其他生效版本自动作废）"""
add_code_block(s, routing_code, Inches(0.4), Inches(1.65), Inches(12.5), Inches(4.4), font_size=11)

# 路由状态机
add_text(s, "路由状态机：DRAFT → EFFECTIVE → OBSOLETE",
         Inches(0.4), Inches(6.15), Inches(12.5), Inches(0.4),
         font_size=13, bold=True, color=COLOR_PRIMARY)
add_text(s, "EFFECTIVE 状态的路由不可删除（须先作废）；同产品同时只允许一个 EFFECTIVE 版本",
         Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.4),
         font_size=12, color=COLOR_GREY)
add_footer(s, 7, TOTAL_PAGES)

# ============ Slide 8: 阶段三 - 生产订单 ============
s = add_blank_slide(prs)
add_title_bar(s, "04  阶段三：计划——生产订单 MO", "ProductionOrder")
add_text(s, "生产计划员按客户 PO 开生产订单，初始 DRAFT 状态，下发后才能派工",
         Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.4), font_size=13, color=COLOR_GREY)

mo_code = """POST /api/v1/production-orders
{
  "product_id": 12,
  "routing_id": 4,                          ← 采用生效的 v1.0 路由
  "batch_no": "W20260818-01",
  "plan_qty": 1000,                         ← 1000 片晶圆
  "planned_start": "2026-08-18T08:00:00",
  "planned_end":   "2026-08-18T20:00:00",
  "customer_po": "CUST-PO-2026-001"
}
→ 返回 mo_no: MO-20260818-0001, status: DRAFT

PUT /api/v1/production-orders/6  { "status": "RELEASED" }     ← 下发"""
add_code_block(s, mo_code, Inches(0.4), Inches(1.65), Inches(7.8), Inches(4.2), font_size=11)

# 右侧：MO 状态机
add_text(s, "MO 状态机", Inches(8.4), Inches(1.65), Inches(4.5), Inches(0.4),
         font_size=14, bold=True, color=COLOR_PRIMARY)

mo_states = [
    ("DRAFT", "草稿", COLOR_GREY),
    ("RELEASED", "已下发", COLOR_ACCENT),
    ("IN_PROGRESS", "开工", COLOR_ORANGE),
    ("COMPLETED", "完工待结", COLOR_GREEN),
    ("CLOSED", "结案", COLOR_PRIMARY),
]
y = Inches(2.15)
for i, (code, label, color) in enumerate(mo_states):
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.4), y, Inches(4.5), Inches(0.45))
    chip.fill.solid()
    chip.fill.fore_color.rgb = color
    chip.line.fill.background()
    chip.shadow.inherit = False
    tf = chip.text_frame
    tf.margin_top = Pt(2)
    p = tf.paragraphs[0]
    p.text = f"{code}  ·  {label}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    if i < len(mo_states) - 1:
        ar = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(10.4), y + Inches(0.45), Inches(0.5), Inches(0.2))
        ar.fill.solid()
        ar.fill.fore_color.rgb = COLOR_GREY
        ar.line.fill.background()
        ar.shadow.inherit = False
    y = y + Inches(0.7)

add_text(s, "状态联动：派工 RUNNING 时\nMO 自动从 RELEASED → IN_PROGRESS",
         Inches(8.4), Inches(5.85), Inches(4.5), Inches(0.8),
         font_size=11, color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)
add_footer(s, 8, TOTAL_PAGES)

# ============ Slide 9: 阶段四 - 派工 ============
s = add_blank_slide(prs)
add_title_bar(s, "05  阶段四：执行——派工 Dispatch（1/2）", "本系统最具差异化的环节")
add_text(s, "派工创建时，系统自动按工段挂接的模板初始化一份空白工艺数据表单",
         Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.4), font_size=14, bold=True, color=COLOR_PRIMARY)

dispatch_code = """POST /api/v1/dispatches
{
  "mo_id": 6,                              ← MO-20260818-0001
  "step_seq": 30,                          ← 第 30 工序：光刻
  "step_name": "光刻",
  "process_section_id": 9,                 ← 光刻工段（可省略，会自动反查）
  "equipment_id": 25,                     ← 生产人员选具体设备：LITHO-001
  "assigned_operator_id": 8,             ← 张三（L3 资质）
  "assigned_team": "A班",
  "dispatch_qty": 1000
}"""
add_code_block(s, dispatch_code, Inches(0.4), Inches(1.7), Inches(12.5), Inches(2.7), font_size=11)

# 返回结果
add_text(s, "服务端自动处理：", Inches(0.4), Inches(4.55), Inches(12.5), Inches(0.4),
         font_size=14, bold=True, color=COLOR_PRIMARY)
add_bullets(s, [
    "校验 MO 非 DRAFT 状态；校验工段 id 存在",
    "解析工艺数据字段模板（按优先级反查 form_template_id）",
    "创建草稿 FormRecord，按 default_value 预填字段",
    "回填 form_template_id / form_record_id 到 Dispatch",
], Inches(0.4), Inches(4.95), Inches(12.5), Inches(1.2), font_size=13)
add_footer(s, 9, TOTAL_PAGES)

# ============ Slide 10: 模板反查优先级 + 联动 ============
s = add_blank_slide(prs)
add_title_bar(s, "05  阶段四：执行——派工 Dispatch（2/2）", "模板反查优先级 + 设备 DOWN 联动")

add_text(s, "工艺数据字段模板反查优先级（_resolve_form_template_id）",
         Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.4),
         font_size=14, bold=True, color=COLOR_PRIMARY)

prio = [
    ("优先级 1", "入参 process_section_id 指定工段的 form_template_id", COLOR_GREEN),
    ("优先级 2", "按 mo.routing_id + step_seq 找 RoutingStep，从其 process_section_id 反查", COLOR_ACCENT),
    ("优先级 3", "RoutingStep 直接挂接的 param_form_template_id", COLOR_ORANGE),
    ("兜底", "找不到则不强制工艺数据采集，派工照常创建", COLOR_GREY),
]
y = Inches(1.7)
for label, content, color in prio:
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), y, Inches(1.3), Inches(0.45))
    chip.fill.solid()
    chip.fill.fore_color.rgb = color
    chip.line.fill.background()
    chip.shadow.inherit = False
    tf = chip.text_frame
    tf.margin_top = Pt(2)
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    add_text(s, content, Inches(1.85), y + Emu(60000), Inches(11), Inches(0.45),
             font_size=12, color=COLOR_DARK)
    y = y + Inches(0.6)

# 设备 DOWN 联动
add_text(s, "设备 DOWN 时的自动联动", Inches(0.4), Inches(4.2), Inches(12.5), Inches(0.4),
         font_size=14, bold=True, color=RGBColor(0xC2, 0x18, 0x5B))

down_flow = """光刻机 LITHO-001 突发故障 → 设备状态切 DOWN
   ├─ 自动创建一条 REPAIR 维修工单
   ├─ 调用 hold_dispatch_by_equipment_down
   │    把该设备上所有 RUNNING 派工切到 HELD
   │    回填 held_reason = "设备DOWN机-关联工单#WO..."
   │    回填 held_work_order_id
   └─ 维修完成、设备 UP 后，派工可从 HELD 退回 RUNNING 继续"""
add_code_block(s, down_flow, Inches(0.4), Inches(4.7), Inches(12.5), Inches(2.4), font_size=11)
add_footer(s, 10, TOTAL_PAGES)

# ============ Slide 11: 派工状态机 + 状态联动 ============
s = add_blank_slide(prs)
add_title_bar(s, "05  阶段四：派工状态机", "Dispatch 状态流转与双向联动")

add_text(s, "派工状态机（VALID_DISPATCH_TRANSITIONS）",
         Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.4),
         font_size=14, bold=True, color=COLOR_PRIMARY)

# 派工状态流程图
states = [
    ("QUEUED", "排队", COLOR_GREY),
    ("ASSIGNED", "已派", COLOR_ACCENT),
    ("RUNNING", "开工", COLOR_ORANGE),
    ("COMPLETED", "完工", COLOR_GREEN),
]
y = Inches(1.8)
x = Inches(0.5)
for i, (code, label, color) in enumerate(states):
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.4), Inches(0.7))
    chip.fill.solid()
    chip.fill.fore_color.rgb = color
    chip.line.fill.background()
    chip.shadow.inherit = False
    tf = chip.text_frame
    tf.margin_top = Pt(4)
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(10)
    p2.font.color.rgb = COLOR_WHITE
    p2.alignment = PP_ALIGN.CENTER
    if i < len(states) - 1:
        add_arrow(s, x + Inches(2.4), y + Inches(0.15), Inches(0.4), Inches(0.4))
    x = x + Inches(2.8)

# 旁支：HELD
held_chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.1), Inches(3.0), Inches(2.4), Inches(0.7))
held_chip.fill.solid()
held_chip.fill.fore_color.rgb = RGBColor(0xC2, 0x18, 0x5B)
held_chip.line.fill.background()
held_chip.shadow.inherit = False
tf = held_chip.text_frame
tf.margin_top = Pt(4)
p = tf.paragraphs[0]
p.text = "HELD"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLOR_WHITE
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "暂停（DOWN/缺料/工艺异常）"
p2.font.size = Pt(10)
p2.font.color.rgb = COLOR_WHITE
p2.alignment = PP_ALIGN.CENTER

add_text(s, "RUNNING ⇄ HELD 双向，HELD 可来自 RUNNING 或退回 RUNNING",
         Inches(0.4), Inches(3.85), Inches(12.5), Inches(0.4),
         font_size=11, color=COLOR_GREY)

# 联动说明
add_text(s, "状态联动（update_dispatch 中自动处理）",
         Inches(0.4), Inches(4.5), Inches(12.5), Inches(0.4),
         font_size=14, bold=True, color=COLOR_PRIMARY)
add_bullets(s, [
    "QUEUED → ASSIGNED：班组长指派设备和作业员",
    "ASSIGNED → RUNNING：自动回填 actual_start；并把 MO 从 RELEASED 自动升到 IN_PROGRESS",
    "RUNNING → COMPLETED：自动回填 actual_end",
    "RUNNING → HELD：必须填 held_reason（DOWN 机/缺料/工艺异常）",
    "COMPLETED → RUNNING：允许退回重做",
], Inches(0.4), Inches(4.9), Inches(12.5), Inches(2.2), font_size=13)
add_footer(s, 11, TOTAL_PAGES)

# ============ Slide 12: 阶段五 - 填表单 + 报工 ============
s = add_blank_slide(prs)
add_title_bar(s, "06  阶段五：填写工艺数据 + 报工", "操作员录入 + 状态联动")

add_text(s, "张三在 HMI 上完成光刻后，按表单填写实际值：",
         Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.4), font_size=13, color=COLOR_GREY)

fill_code = """PUT /api/v1/form-records/42
{
  "values": [
    {"field_key":"exposure_energy","field_value":55.2},
    {"field_key":"focus_offset",   "field_value":-0.3},
    {"field_key":"alignment_error","field_value":28},              ← < 35nm 阈值，合格
    {"field_key":"develop_time",   "field_value":62},              ← 覆盖默认 60
    {"field_key":"reticle_id",     "field_value":"RET-2026-A07"},
    {"field_key":"abnormal",       "field_value":"NORMAL"}
  ]
}

PATCH /api/v1/form-records/42/submit   →  status: 已提交, submitted_at 自动回填
PATCH /api/v1/form-records/42/audit    →  status: 已审核（含二次密码+电子签名）"""
add_code_block(s, fill_code, Inches(0.4), Inches(1.6), Inches(8.5), Inches(5.2), font_size=10)

# 右侧：审核后的留痕修正
add_text(s, "审核后的记录管理", Inches(9.1), Inches(1.6), Inches(4), Inches(0.4),
         font_size=14, bold=True, color=COLOR_PRIMARY)
add_bullets(s, [
    "审核通过后记录锁定，不可原地改",
    "修正必须走 FormRecordAmendment",
    "需填写：原值、新值、原因",
    "修正人二次密码校验+电子签名",
    "可由审核人批准修正",
    "满足 ISO/车规体系要求",
], Inches(9.1), Inches(2.05), Inches(4), Inches(2.5), font_size=12)

add_text(s, "报工 LaborReport", Inches(9.1), Inches(4.6), Inches(4), Inches(0.4),
         font_size=14, bold=True, color=COLOR_PRIMARY)
add_bullets(s, [
    "关联 dispatch_id",
    "投入/合格/不良数量",
    "session_start/end 时段",
    "man_hours 人×小时",
    "form_record_id 关联工艺数据表",
], Inches(9.1), Inches(5.05), Inches(4), Inches(2.0), font_size=12)
add_footer(s, 12, TOTAL_PAGES)

# ============ Slide 13: 全链路追溯 ============
s = add_blank_slide(prs)
add_title_bar(s, "07  全链路追溯能力", "数据主线一条不可断")

add_text(s, "任何追溯查询都能拿到完整数据链：",
         Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.4), font_size=13, color=COLOR_GREY)

trace_code = """ProductionOrder (MO-20260818-001)
  └─ Dispatch (派工#15 光刻, 设备=LITHO-001, 操作员=张三)
       ├─ LaborReport (报工: 投入1000/合格995/不良5)
       └─ FormRecord (工艺参数: exposure_energy=55.2, alignment_error=28nm, ...)
              ├─ 引用 FormTemplate (光刻工艺参数记录表 v1)
              └─ FormRecordAmendment (如有修正, 留痕)"""
add_code_block(s, trace_code, Inches(0.4), Inches(1.6), Inches(12.5), Inches(3.0), font_size=12)

add_text(s, "典型追溯场景", Inches(0.4), Inches(4.85), Inches(12.5), Inches(0.4),
         font_size=14, bold=True, color=COLOR_PRIMARY)
add_bullets(s, [
    "QA/工艺工程师通过 GET /api/v1/form-records/{id}/export/json 或 /export/csv 导出",
    "SPC 分析：曝光能量是否偏离控制线",
    "不良追溯：某片晶圆 CD 超差，反查当时 LITHO-001 的焦距偏移",
    "归档备审：满足 ISO 9001 / IATF 16949 / 车规 AEC-Q100 体系追溯要求",
], Inches(0.4), Inches(5.25), Inches(12.5), Inches(1.8), font_size=13)
add_footer(s, 13, TOTAL_PAGES)

# ============ Slide 14: 设计要点回顾 ============
s = add_blank_slide(prs)
add_title_bar(s, "07  设计要点回顾", "关键设计决策一览")
add_text(s, "工艺定义与生产执行彻底解耦：工艺工程师聚焦\"采什么数据\"，生产人员派工时聚焦\"在哪台机做\"",
         Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.4), font_size=13, color=COLOR_GREY)

decisions = [
    ("工段绑定设备组（模板层），不绑具体设备", "ProcessSection.equipment_group", COLOR_ACCENT),
    ("派工时由生产人员选具体设备（执行层）", "Dispatch.equipment_id", COLOR_ORANGE),
    ("工段作为独立可复用工段库", "ProcessSection 主表，被多个 RoutingStep 引用", COLOR_PRIMARY),
    ("派工自动初始化空白工艺表单", "_init_form_record_for_dispatch，按 default_value 预填", COLOR_GREEN),
    ("工段被引用时禁止硬删除", "delete_process_section 校验引用", RGBColor(0xC2, 0x18, 0x5B)),
    ("模板停用温和降级", "is_active=False 时仅跳过初始化，不阻断派工", COLOR_GREY),
    ("工艺数据审核后不可原地改", "FormRecordAmendment 留痕修正", COLOR_PRIMARY),
    ("设备 DOWN 自动暂停派工", "hold_dispatch_by_equipment_down 联动钩子", RGBColor(0xC2, 0x18, 0x5B)),
    ("MO 与派工状态联动", "派工 RUNNING 自动把 MO 从 RELEASED 升 IN_PROGRESS", COLOR_ACCENT),
]

# 双列卡片
col_w = Inches(6.1)
row_h = Inches(0.95)
left_positions = [Inches(0.4), Inches(6.8)]
for i, (decision, impl, color) in enumerate(decisions):
    col = i % 2
    row = i // 2
    left = left_positions[col]
    top = Inches(1.75) + row * row_h
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, col_w, Inches(0.85))
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_LIGHT_BG
    card.line.color.rgb = color
    card.line.width = Pt(1.2)
    card.shadow.inherit = False
    # 序号色块
    seq = s.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.1), top + Inches(0.15), Inches(0.5), Inches(0.5))
    seq.fill.solid()
    seq.fill.fore_color.rgb = color
    seq.line.fill.background()
    seq.shadow.inherit = False
    tf = seq.text_frame
    tf.margin_top = Pt(0)
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    # 文字
    add_text(s, decision, left + Inches(0.7), top + Emu(50000), col_w - Inches(0.8), Inches(0.45),
             font_size=12, bold=True, color=COLOR_DARK)
    add_text(s, "→ " + impl, left + Inches(0.7), top + Emu(380000), col_w - Inches(0.8), Inches(0.4),
             font_size=10, color=color)

# 底部结语
add_text(s, "Q&A  /  Thank you",
         Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.5),
         font_size=20, bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)
add_footer(s, 14, TOTAL_PAGES)

# ---------- 保存 ----------
output_path = "/workspace/SEMS生产系统培训_光刻工艺段示范.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
