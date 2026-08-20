#!/usr/bin/env python3
"""Generate SEMS_功能介绍.pptx V2026.08 (25 slides) with detailed business workflows.

Based on generate_ppt.py 16 slides, plus 9 new detailed workflow slides (P17-P25):
  P17 登录 & 主题切换 & 首次改密
  P18 设备 DOWN → 工单自动触发流程
  P19 工单全生命周期（创建→派工→维修→验证→关闭 + SLA 超期升级）
  P20 PM 计划 / 点检巡检流程
  P21 文控审批链（三级电子签名状态机）
  P22 表单模板 + 结构化电子表单（管理员→操作员）
  P23 8D 报告 + 一键归档知识库双路径
  P24 故障知识库 + 相似案例 + 复发追踪
  P25 系统配置：备份/恢复/灾备 3-2-1 流程
"""

import io
import zipfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Design System ──
DARK_BLUE  = RGBColor(0x0B, 0x3D, 0x91)
TEAL       = RGBColor(0x00, 0x96, 0x88)
LIGHT_GRAY = RGBColor(0xF5, 0xF7, 0xFA)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x2C, 0x3E, 0x50)
GRAY_TEXT  = RGBColor(0x7F, 0x8C, 0x8D)
ACCENT     = RGBColor(0x1A, 0x73, 0xE8)
CARD_BG    = RGBColor(0xEE, 0xF2, 0xF7)
RED        = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE     = RGBColor(0xF3, 0x9C, 0x12)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
FONT       = 'Microsoft YaHei'

SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

# ── Helpers ──
def _set_run(run, text, size, color, bold=False, font_name=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name

def _add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def _add_circle(slide, left, top, size, fill_color, text, text_color=WHITE, text_size=14):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    _set_run(run, text, text_size, text_color, bold=True)
    # Vertical center
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape

def _add_arrow(slide, left, top, width, height, color=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def _add_text(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    _set_run(run, text, size, color, bold)
    return txBox

def _add_bullets(slide, left, top, width, height, items, size=14, color=DARK_TEXT, spacing=6):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        if isinstance(item, tuple):
            prefix, rest = item
            r1 = p.add_run()
            _set_run(r1, prefix, size, TEAL, bold=True)
            r2 = p.add_run()
            _set_run(r2, rest, size, color)
        else:
            r = p.add_run()
            _set_run(r, item, size, color)
    return txBox

def _add_top_bar(slide, title, subtitle=None):
    _add_rect(slide, 0, 0, SLIDE_W, Emu(914400), DARK_BLUE)
    _add_text(slide, Emu(457200), Emu(137160), Emu(9000000), Emu(457200),
              title, 24, WHITE, bold=True)
    if subtitle:
        _add_text(slide, Emu(457200), Emu(548640), Emu(9000000), Emu(274320),
                  subtitle, 13, RGBColor(0xAE, 0xC5, 0xE8))

def _add_card(slide, left, top, width, height, title, items, title_color=DARK_BLUE, item_size=12):
    _add_rect(slide, left, top, width, height, LIGHT_GRAY)
    _add_text(slide, left + Emu(182880), top + Emu(91440),
              width - Emu(365760), Emu(365760),
              title, 14, title_color, bold=True)
    _add_bullets(slide, left + Emu(182880), top + Emu(457200),
                 width - Emu(365760), height - Emu(548640),
                 items, size=item_size, spacing=4)

def _extract_screenshot(orig_path):
    z = zipfile.ZipFile(orig_path)
    for name in z.namelist():
        if name.startswith('ppt/media/') and name.endswith('.png'):
            return z.read(name)
    return None

# ── Vertical numbered-step flow builder ──
def _add_step_flow(slide, left, top, width, steps, step_colors=None):
    """
    Render a vertical numbered step flow.
    steps = [(label, description, note?), ...]
    """
    if step_colors is None:
        step_colors = [TEAL]*len(steps)
    step_h = Emu(580000)
    y = top
    for i, (label, desc) in enumerate(steps):
        color = step_colors[i] if i < len(step_colors) else TEAL
        # Circle number
        _add_circle(slide, left, y + Emu(150000), Emu(457200), color, str(i+1), WHITE, 16)
        # Label + description card
        card_left = left + Emu(600000)
        _add_rect(slide, card_left, y, width - Emu(600000), step_h - Emu(30000), LIGHT_GRAY)
        _add_text(slide, card_left + Emu(137160), y + Emu(91440),
                  width - Emu(800000), Emu(274320), label, 13, color, bold=True)
        _add_text(slide, card_left + Emu(137160), y + Emu(365760),
                  width - Emu(800000), Emu(274320), desc, 11, DARK_TEXT)
        # Arrow connector (except last)
        if i < len(steps) - 1:
            _add_rect(slide, left + Emu(210000), y + step_h - Emu(30000),
                      Emu(38100), Emu(90000), color)
        y += step_h

# ════════════════════════════════════════════════════════════════════════════
# Original 16 slides (unchanged structure)
# ════════════════════════════════════════════════════════════════════════════

def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_BLUE)
    _add_rect(slide, Emu(4572000), Emu(3200400), Emu(3047695), Emu(38100), TEAL)
    _add_text(slide, Emu(914400), Emu(1828800), Emu(10362895), Emu(914400),
              "SEMS 半导体设备管理系统", 40, WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, Emu(914400), Emu(2833920), Emu(10362895), Emu(548640),
              "Semiconductor Equipment Management System", 20,
              RGBColor(0xAE, 0xC5, 0xE8), align=PP_ALIGN.CENTER)
    _add_text(slide, Emu(914400), Emu(4572000), Emu(10362895), Emu(457200),
              "功能介绍  |  V2026.08  |  2026-08", 18,
              RGBColor(0xAE, 0xC5, 0xE8), align=PP_ALIGN.CENTER)
    _add_text(slide, Emu(914400), Emu(5486400), Emu(10362895), Emu(365760),
              "FastAPI  +  Vue 3  +  Element Plus  +  SQLite", 14,
              RGBColor(0x8A, 0xAA, 0xC8), align=PP_ALIGN.CENTER)

def build_why_sems(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "为什么需要 SEMS？", "Why SEMS?")
    pain_items = [
        "设备状态靠电话/微信群，台账分散在 Excel/纸档",
        "故障处置经验沉淀差，同一根因反复复发",
        "点检/PM/润滑执行难追踪，容易漏做",
        "备件库存不透明，等要修了才发现缺件",
        "工艺文件版本与审批靠邮件，追溯困难",
        "8D 报告写完就丢，知识库与品管工具完全脱节",
    ]
    _add_bullets(slide, Emu(457200), Emu(1280160), Emu(7000000), Emu(4500000),
                 [("• ", p) for p in pain_items], size=15, spacing=10)
    _add_card(slide, Emu(7700000), Emu(1280160), Emu(4000000), Emu(2500000),
              "一句话价值",
              ["把设备台账、日常运维、工艺文控、品管合规、数据价值",
               "五个核心闭环落在一套系统里，",
               "真正做到数据可查、经验可复用。"], title_color=TEAL)

def build_five_loops(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "五大业务闭环", "Five Business Loops")
    loops = [
        ("① 设备台账 & 状态闭环", "设备台账 / 状态切换 / 设备全生命周期 T0-T3 / 润滑管理", "机台状态一眼看清，DOWN 自动派工单"),
        ("② 日常运维闭环", "工单 / PM 计划 / 点检巡检 / 备件管理 / 工单 SLA 升级", "从建单→维修→验证→关单一气呵成，超时自动升级"),
        ("③ 工艺文控闭环", "工艺文件 / 表单模板 / 电子表单 / 文控审批链 / 分发收回 / 复审告警", "符合体系标准的受控文控，水印+指纹+留痕"),
        ("④ 品管合规闭环", "8D / FMEA / 安全检查 / 环境核查", "质量分析合规可追溯，证书到期提前告警"),
        ("⑤ 数据价值闭环", "故障知识库(工单归档+8D归档) / OEE / 设备成本 LCC", "故障经验不随人流失，LCC 一眼算清单台成本"),
    ]
    y = Emu(1100000)
    for title, modules, value in loops:
        _add_rect(slide, Emu(457200), y, Emu(11277000), Emu(914400), LIGHT_GRAY)
        _add_text(slide, Emu(640080), y + Emu(91440), Emu(3200000), Emu(365760), title, 14, DARK_BLUE, bold=True)
        _add_text(slide, Emu(3900000), y + Emu(91440), Emu(5000000), Emu(731520), modules, 11, DARK_TEXT)
        _add_text(slide, Emu(9000000), y + Emu(91440), Emu(2600000), Emu(731520), value, 11, TEAL)
        y += Emu(1000000)

def build_v2026_updates(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "V2026.08 版本核心更新", "Key Updates in V2026.08")
    cards = [
        ("📚 8D → 知识库 一键归档", ["D0+D2 → 故障现象 / D4 → 根因", "D5 → 处置 / D7 → 预防", "关联 source_d8_report_id 溯源", "同时保留「工单归档」入口"]),
        ("📋 8 大分组分级菜单", ["① 总览 ② 设备管理 ③ 运维工单", "④ 安全与环境 ⑤ 备件与人员", "⑥ 工艺文控 ⑦ 分析改进", "⑧ 系统配置（永远在最下方）"]),
        ("🎯 角色定制看板", ["6 种角色登录→不同看板组合", "QA 优先：文控复审 + 8D 进度", "操作员优先：设备状态 + 今日点检"]),
        ("🧾 审计日志入库 + 🔐 bcrypt 修复", ["11 类敏感操作写入 audit_logs 表", "操作人/目标/IP/UA/时间戳/详情", "SHA-256 预哈希解决 bcrypt 72 字节截断", "旧值透明回退 + 登录自动升级重哈希"]),
    ]
    positions = [
        (Emu(457200), Emu(1100000), Emu(5600000), Emu(2500000)),
        (Emu(6200000), Emu(1100000), Emu(5600000), Emu(2500000)),
        (Emu(457200), Emu(3700000), Emu(5600000), Emu(2500000)),
        (Emu(6200000), Emu(3700000), Emu(5600000), Emu(2500000)),
    ]
    for (title, items), pos in zip(cards, positions):
        _add_card(slide, pos[0], pos[1], pos[2], pos[3], title, items)

def build_dashboard_roles(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "看板总览 & 角色定制", "Dashboard & Role Customization")
    roles = [
        ("管理员", "SLA 达成率 → 文控复审告警 → 8D 进度 → 知识库新增 → 低库存备件 → 安全检查到期 → 系统健康"),
        ("工程师", "我的工单 → 待验证工单 → 到期 PM → 点检 → 8D 进行中 → 相似故障推荐"),
        ("QA", "审核中文档 → 复审到期 → 附加修正审批 → 8D 报告 → 安全检查 → 表单审核清单"),
        ("操作员", "设备状态总览 → 我的工单 → 今日点检 → 润滑到期 → 常用设备快速入口"),
    ]
    y = Emu(1100000)
    for role, cards in roles:
        _add_rect(slide, Emu(457200), y, Emu(11277000), Emu(914400), LIGHT_GRAY)
        _add_text(slide, Emu(640080), y + Emu(137160), Emu(1500000), Emu(640080), role, 14, DARK_BLUE, bold=True)
        _add_text(slide, Emu(2200000), y + Emu(137160), Emu(9400000), Emu(640080), cards, 12, DARK_TEXT)
        y += Emu(1000000)
    _add_text(slide, Emu(457200), Emu(5400000), Emu(11277000), Emu(457200),
              "支持：明色青绿 / 暗色霓虹 / 跟随系统 三套模式一键切换；暗色模式做了亮度压暗专项补丁，夜班大屏久看不累。",
              13, TEAL, bold=True)

def build_device_mgmt(prs, screenshot=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "① 设备管理 & DOWN→工单自动派发", "Equipment Management")
    if screenshot:
        slide.shapes.add_picture(io.BytesIO(screenshot),
            Emu(365760), Emu(1280000), Emu(5600000), Emu(4800000))
        panel_left = Emu(6200000); panel_w = Emu(5600000)
    else:
        panel_left = Emu(457200); panel_w = Emu(11277000)
    _add_card(slide, panel_left, Emu(1280000), panel_w, Emu(1400000),
              "设备台账 & 状态切换",
              ["厂区/区域分类、附件上传（说明书/图纸/SOP）",
               "任意角色切 DOWN → 自动创建 REPAIR 工单 → 自动派工程师",
               "状态：RUN / IDLE / DOWN / PM / ENGINEERING / PROCESS_VALIDATION / OFFLINE"])
    _add_card(slide, panel_left, Emu(2800000), panel_w, Emu(1400000),
              "设备全生命周期 T0-T3",
              ["T0 选型(URS/候选供应商) → T1 采购(PO/金额/交付)",
               "T2 安装调试(FAT/SAT) → T3 量产移交(验收结果)",
               "时间线视图一眼看全流程进度"])
    _add_card(slide, panel_left, Emu(4320000), panel_w, Emu(1400000),
              "润滑管理（五定）",
              ["定点/定人/定时/定质/定量",
               "自动推算下次润滑日，到期提醒高亮"])

def build_daily_ops(prs, screenshot=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "② 日常运维闭环 & 工单 SLA", "Daily Operations & Work Order SLA")
    if screenshot:
        slide.shapes.add_picture(io.BytesIO(screenshot),
            Emu(365760), Emu(1280000), Emu(5600000), Emu(4800000))
        panel_left = Emu(6200000); panel_w = Emu(5600000)
    else:
        panel_left = Emu(457200); panel_w = Emu(11277000)
    _add_card(slide, panel_left, Emu(1280000), panel_w, Emu(1300000),
              "工单管理",
              ["创建/跟踪、紧急度标签、持续时长实时显示、关键词全文检索",
               "详情：5Why 根因分析 + 备件领用 + 状态流转(OPEN→CLOSED)"])
    _add_card(slide, panel_left, Emu(2700000), panel_w, Emu(1300000),
              "SLA & 超期升级",
              ["SLA 目标响应/解决时长（按紧急度预设）",
               "实际时长自动计算 → SLA 达成率统计 → 超期自动升级指派"])
    _add_card(slide, panel_left, Emu(4120000), panel_w, Emu(1600000),
              "PM / 点检 / 备件",
              ["PM：周/双周/月/季 周期，批量生成到期 PM 工单",
               "点检：模板+检查项+历史记录",
               "备件：库存/出入库/低库存告警/设备易损件绑定"])

def build_process_control(prs, screenshot=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "③ 工艺文控闭环", "Process Document Control")
    if screenshot:
        slide.shapes.add_picture(io.BytesIO(screenshot),
            Emu(365760), Emu(1280000), Emu(5600000), Emu(4800000))
        panel_left = Emu(6200000); panel_w = Emu(5600000)
    else:
        panel_left = Emu(457200); panel_w = Emu(11277000)
    _add_card(slide, panel_left, Emu(1280000), panel_w, Emu(1300000),
              "工艺文件 & 结构化电子表单",
              ["指导性文件：版本管理，生效版唯一，同组旧版自动作废",
               "管理员定义模板（9 种字段类型）→ 操作员动态填写 → JSON/CSV 导出"])
    _add_card(slide, panel_left, Emu(2700000), panel_w, Emu(1500000),
              "文控系统（符合体系标准）",
              ["三级电子签名审批链：编制→审核→批准（二次密码校验 + SHA256 指纹）",
               "状态机白名单：草稿↔审核中→生效→作废",
               "修订记录：字段级 before/after 对比"])
    _add_card(slide, panel_left, Emu(4320000), panel_w, Emu(1400000),
              "受控管理 & 复审告警",
              ["分发收回台账、PDF 自动加盖受控章水印",
               "表单审核锁定 + 附加修正流程",
               "复审周期告警：30 天内到期 / 已过期 Badge"])

def build_quality(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "④ 品管合规闭环", "Quality & Compliance")
    _add_card(slide, Emu(457200), Emu(1100000), Emu(5600000), Emu(2400000),
              "8D 报告",
              [("D1-D8 ", "8 步完整建模（团队→问题→遏制→根因→措施→验证→预防→表彰）"),
               ("关联工单 ", "（可选），状态：DRAFT / IN_PROGRESS / CLOSED"),
               ("⭐ 一键归档知识库 ", "：D0+D2→现象 / D4→根因 / D5→处置 / D7→预防")])
    _add_card(slide, Emu(6200000), Emu(1100000), Emu(5600000), Emu(2400000),
              "FMEA & 环境核查",
              [("FMEA ", "：失效模式 / 影响 / S×O×D → RPN，RPN>100 高亮"),
               ("环境核查 ", "：温度/湿度/洁净度/VOC/ESD 参数录入与趋势")])
    _add_card(slide, Emu(457200), Emu(3600000), Emu(11277000), Emu(2400000),
              "安全检查（四类）",
              ["safety_device 安全装置（防护罩/急停/联锁）/ 特种设备 / 环保 / 消防",
               "按频率自动推算下次检查日，30 天内黄底、已过期红字",
               "特种设备证书到期 Badge，检查发现 / 整改措施跟踪"])

def build_knowledge_base(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "⑤ 故障知识库（双路径归档）", "Knowledge Base (Dual-Path Archive)")
    y = Emu(1100000)
    paths = [
        ("工单", "知识库→「从工单归档」", "description→现象 / fault_category→分类 / root_cause→根因 / solution→处置 / prevention→预防"),
        ("8D 报告", "知识库→「从 8D 归档」或 8D 列表→「归档至知识库」", "D0+D2→现象 / D4→根因 / D5→处置 / D7→预防"),
        ("手动", "新建条目", "任意填写"),
    ]
    _add_rect(slide, Emu(457200), y, Emu(11277000), Emu(365760), CARD_BG)
    _add_text(slide, Emu(640080), y + Emu(45720), Emu(1400000), Emu(274320), "来源", 12, DARK_BLUE, bold=True)
    _add_text(slide, Emu(2100000), y + Emu(45720), Emu(4200000), Emu(274320), "入口", 12, DARK_BLUE, bold=True)
    _add_text(slide, Emu(6400000), y + Emu(45720), Emu(5200000), Emu(274320), "自动映射", 12, DARK_BLUE, bold=True)
    y += Emu(365760)
    for source, entry, mapping in paths:
        _add_rect(slide, Emu(457200), y, Emu(11277000), Emu(548640), LIGHT_GRAY)
        _add_text(slide, Emu(640080), y + Emu(91440), Emu(1400000), Emu(365760), source, 12, TEAL, bold=True)
        _add_text(slide, Emu(2100000), y + Emu(91440), Emu(4200000), Emu(365760), entry, 11, DARK_TEXT)
        _add_text(slide, Emu(6400000), y + Emu(91440), Emu(5200000), Emu(365760), mapping, 11, DARK_TEXT)
        y += Emu(548640)
    _add_card(slide, Emu(457200), Emu(3300000), Emu(5600000), Emu(2700000),
              "相似案例推荐",
              [("同设备 + 同故障分类 ", "优先 → 退化匹配"),
               ("按复发次数/浏览量倒序", ""),
               ("新故障一眼看到以前怎么修的", "")])
    _add_card(slide, Emu(6200000), Emu(3300000), Emu(5600000), Emu(2700000),
              "复发追踪",
              [("同一根因复发 → ", "点「标记复发」→ recurrence_count +1"),
               ("看板/列表一眼看出哪些根因反复出现", ""),
               ("需要重点改进的问题一目了然", "")])

def build_lcc_oee(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "⑤ 数据价值（续）：LCC & OEE", "LCC & OEE Analysis")
    _add_card(slide, Emu(457200), Emu(1100000), Emu(5600000), Emu(4800000),
              "设备成本 LCC（全生命周期成本）",
              [("6 种成本类型：", ""),
               ("  采购 / 维护 / 备件 / 能耗 / 折旧 / 报废", ""),
               ("单设备汇总：", "各类成本占比 + 总成本"),
               ("全设备汇总：", "按类型统计（饼图）+ Top10 高成本设备排名"),
               ("单设备年度趋势：", "按年汇总（折线图）")])
    _add_card(slide, Emu(6200000), Emu(1100000), Emu(5600000), Emu(4800000),
              "OEE 分析",
              [("设备综合效率 = ", "可用率 × 性能率 × 良品率"),
               ("按设备 / 时间维度聚合", ""),
               ("按日/周/月维度趋势图展示", ""),
               ("设备对比分析视图", ""),
               ("OEE 组成拆解，瓶颈定位", "")])

def build_security(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "安全 & 审计 & 灾备", "Security, Audit & Backup")
    _add_card(slide, Emu(457200), Emu(1100000), Emu(5600000), Emu(2400000),
              "安全加固（面向局域网部署）",
              [("JWT 双令牌", "：2h access + 7d refresh，5 分钟前自动续期"),
               ("bcrypt", "(SHA-256 预哈希，修复 72 字节截断)"),
               ("密码 3/4 类字符复杂度 / 失败锁定 5 次/15min", ""),
               ("首次登录强制改密 / 用户名枚举防护", ""),
               ("CSP / X-Frame / CORS 白名单", "")])
    _add_card(slide, Emu(6200000), Emu(1100000), Emu(5600000), Emu(2400000),
              "审计日志入库（审核留痕）",
              [("11 类动作全入库：", ""),
               ("  LOGIN_OK / FAIL / LOCKED / LOGOUT", ""),
               ("  PASSWORD_CHANGED / RESET", ""),
               ("  USER_CREATE / UPDATE / DELETE / UNLOCK", ""),
               ("  RESTORE_BACKUP", ""),
               ("每条含：", "action / actor / target / ip / UA / detail / time")])
    _add_card(slide, Emu(457200), Emu(3600000), Emu(11277000), Emu(2400000),
              "灾备（3-2-1 策略）",
              [("应用内定时备份 + SQLite 热快照 + ZIP + AES-256(Fernet/PBKDF2 20万轮) + NAS/SMB/U盘异地副本", ""),
               ("系统级旁路备份脚本（sh/bat，不依赖后端进程也能跑）", ""),
               ("备份后自动烟雾还原测试：解压→打开 sqlite→查表行数→校验 uploads，永远知道备份能不能真还原", "")])

def build_deployment(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "部署方案选型", "Deployment Options")
    y = Emu(1100000)
    headers = [("场景", 2800000), ("推荐部署", 3000000), ("自启", 1600000), ("崩溃重启", 1800000), ("健康检查", 2077000)]
    x = Emu(457200)
    _add_rect(slide, x, y, Emu(11277000), Emu(365760), CARD_BG)
    for h, w in headers:
        _add_text(slide, x + Emu(91440), y + Emu(45720), Emu(w), Emu(274320), h, 11, DARK_BLUE, bold=True)
        x += Emu(w)
    y += Emu(365760)
    rows = [
        ("工厂 Linux 有 root", "systemd（方案 A）", "✅ 开机自启", "✅ on-failure", "✅ cron 2min→重启"),
        ("普通用户 / 受限 Linux", "watchdog_user.sh（方案 B）", "✅ @reboot", "✅ 每分钟 tick", "✅ 同上"),
        ("Windows Server / 工控机", "NSSM + 任务计划（方案 C）", "✅ 延迟自启", "✅ NSSM + SC failure", "✅ 2min 计划任务"),
        ("容器化 / 快速体验", "Docker Compose（方案 D）", "✅ unless-stopped", "✅ compose restart", "✅ healthcheck"),
    ]
    for row in rows:
        x = Emu(457200)
        _add_rect(slide, x, y, Emu(11277000), Emu(548640), LIGHT_GRAY)
        for val, (_, w) in zip(row, headers):
            _add_text(slide, x + Emu(91440), y + Emu(91440), Emu(w), Emu(365760), val, 10, DARK_TEXT)
            x += Emu(w)
        y += Emu(548640)
    _add_text(slide, Emu(457200), y + Emu(91440), Emu(11277000), Emu(731520),
              "入口级加固：端口占用友好提示 / SQLite WAL checkpoint 钩子 / uvicorn 优雅退出 / 日志双写（journald + 按天滚动 14/30 天）",
              12, TEAL, bold=True)

def build_accounts(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "默认账号 & 演示数据", "Default Accounts & Demo Data")
    accounts = [
        ("admin", "管理员", "admin123", "全功能 / 用户与权限 / 备份与恢复"),
        ("engineer1", "张工 工程师", "eng123", "工单处置 / 8D / 知识库录入"),
        ("process1", "周工艺 工艺员", "proc123", "工艺文件 / 电子表单 / 提交文控审核"),
        ("qa1", "陈品管 QA", "qa123", "文控审核批准 / 8D / FMEA / 安全检查"),
        ("operator1", "王操作 操作员", "op123", "切换设备状态 / 点检 / 电子表单填写"),
        ("viewer1", "孙查看 查看者", "view123", "只读浏览（设备/工单/文档/知识）"),
    ]
    y = Emu(1100000)
    _add_rect(slide, Emu(457200), y, Emu(11277000), Emu(365760), CARD_BG)
    cols = [("账号", 2000000), ("角色", 2500000), ("默认密码", 2000000), ("典型工作内容", 4777000)]
    x = Emu(457200)
    for h, w in cols:
        _add_text(slide, x + Emu(91440), y + Emu(45720), Emu(w), Emu(274320), h, 11, DARK_BLUE, bold=True)
        x += Emu(w)
    y += Emu(365760)
    for acc, role, pwd, work in accounts:
        x = Emu(457200)
        _add_rect(slide, x, y, Emu(11277000), Emu(457200), LIGHT_GRAY)
        for val, (_, w) in zip([acc, role, pwd, work], cols):
            _add_text(slide, x + Emu(91440), y + Emu(91440), Emu(w), Emu(274320), val, 10, DARK_TEXT)
            x += Emu(w)
        y += Emu(457200)
    _add_text(slide, Emu(457200), y + Emu(182880), Emu(11277000), Emu(731520),
              "演示数据开箱即有：10 台设备 + 27 条状态日志 + 16 附件 + 14 备件 + 12 工单 + 4 份 8D 报告 + 14 工艺文件 + 12 安全检查 + 9 润滑点 + 9 条故障知识 + 31 成本记录 等 30+ 类记录。",
              12, TEAL, bold=True)

def build_tech_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "技术栈 & 扩展路线", "Tech Stack & Roadmap")
    _add_card(slide, Emu(457200), Emu(1100000), Emu(5600000), Emu(3500000),
              "当前技术栈",
              [("后端：", "Python 3.10+ / FastAPI / SQLAlchemy 2.0 / SQLite"),
               ("前端：", "Vue 3 / Vite / Element Plus / ECharts / Pinia"),
               ("安全：", "JWT / bcrypt(+SHA256 prehash) / CSP / CORS 白名单"),
               ("备份：", "SQLite hot backup / ZIP / Fernet AES-256"),
               ("打包：", "Docker / systemd / NSSM / PyInstaller")])
    _add_card(slide, Emu(6200000), Emu(1100000), Emu(5600000), Emu(1600000),
              "短期扩展路线（Q4 规划）",
              ["① 审计日志前端界面化（查询/导出 CSV）",
               "② 知识条目审核流（草稿→QA 审核→生效）",
               "③ 8D 报告通知（钉钉/飞书 Webhook）",
               "④ 备品低库存自动建采购申请单"])
    _add_card(slide, Emu(6200000), Emu(2900000), Emu(5600000), Emu(1700000),
              "中期扩展",
              ["对接 MES：设备状态 / 工艺参数自动上传",
               "对接 SCADA：关键机台参数趋势、阈值报警",
               "移动端（微信小程序/钉钉微应用）：点检 / 工单"])

def build_contact(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "文档索引 & 更新记录", "Documentation & Changelog")
    _add_card(slide, Emu(457200), Emu(1100000), Emu(5600000), Emu(2500000),
              "文档索引",
              [("README.md", " — 部署 & 架构详情 & 安全说明 & 灾备"),
               ("用户使用教程.md", " — 日常使用操作手册（按角色分步骤）"),
               ("SEMS_功能介绍_最新版.md", " — PPT 文字源（本文件）"),
               ("SEMS_UI设计评审与方案对比.pptx", " — UI 设计方案评审对比")])
    _add_card(slide, Emu(6200000), Emu(1100000), Emu(5600000), Emu(2500000),
              "更新记录",
              [("V2026.08：", "8D→知识库双归档、8 组分级菜单、角色看板、审计日志入库、bcrypt 截断修复"),
               ("V2026.06：", "文控系统(三级签名)、安全检查、润滑管理、工单 SLA、设备生命周期、知识库、LCC"),
               ("V2026.04：", "B+C 双皮肤主题、表单模板结构化电子表单、DOWN→工单自动派发")])
    _add_text(slide, Emu(457200), Emu(4400000), Emu(11277000), Emu(731520),
              "默认账号：admin / admin123（首次登录强制改密）",
              16, DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, Emu(457200), Emu(5200000), Emu(11277000), Emu(731520),
              "感谢观看！", 28, TEAL, bold=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# NEW: 9 detailed workflow slides (P17-P25)
# ════════════════════════════════════════════════════════════════════════════

def build_workflow_login(prs):
    """P17: 登录 & 首次改密 & 主题切换"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "【流程 1】登录 · 首次改密 · 主题切换", "Login & First-PW-Change & Theme Switch")
    # Left column: login flow
    _add_card(slide, Emu(457200), Emu(1050000), Emu(5600000), Emu(1100000),
              "🔐 登录 + 首次强制改密",
              [], item_size=0)
    steps = [
        ("打开系统登录页", "浏览器访问 http://<服务器IP>:8080（开发环境 localhost:5173）"),
        ("输入账号密码登录", "默认 admin/admin123；其他账号见 P14 账号表"),
        ("首次登录 → 跳转改密页", "要求 ≥8 位 + 3/4 类字符（大写/小写/数字/符号）"),
        ("改密成功，进入系统", "后续 JWT 2h 自动续期，长期不操作需重新登录"),
    ]
    _add_step_flow(slide, Emu(550000), Emu(1500000), Emu(5400000), steps,
                   step_colors=[DARK_BLUE, TEAL, ORANGE, GREEN])
    # Right column: theme switch
    _add_card(slide, Emu(6200000), Emu(1050000), Emu(5600000), Emu(1100000),
              "🎨 三模式主题一键切换",
              [], item_size=0)
    steps2 = [
        ("点击右上角主题按钮", "图标显示当前模式：🌞 明色 / 🌜 暗色 / 🖥 跟随系统"),
        ("下拉选择目标模式", "明色青绿(明亮办公室) | 暗色霓虹(夜班大屏) | 跟随系统(OS 浅深联动)"),
        ("立即生效 + 本地保存", "写入 localStorage.sems_theme_mode，不同电脑独立生效"),
    ]
    _add_step_flow(slide, Emu(6300000), Emu(1500000), Emu(5400000), steps2,
                   step_colors=[TEAL, ORANGE, GREEN])
    # Tip at bottom
    _add_rect(slide, Emu(457200), Emu(5750000), Emu(11277000), Emu(800000), LIGHT_GRAY)
    _add_text(slide, Emu(640080), Emu(5790000), Emu(11000000), Emu(365760),
              "💡 提示：若反复提示登录过期 → 清除浏览器 localStorage 中 sems_ 前缀项后重新登录",
              13, RED, bold=True)

def build_workflow_down_wo(prs):
    """P18: 设备 DOWN → 工单自动触发"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "【流程 2】设备 DOWN → 故障工单自动触发", "Equipment DOWN → Auto Work Order Trigger")
    _add_card(slide, Emu(457200), Emu(1050000), Emu(11277000), Emu(2000000),
              "📌 前置条件",
              ["操作角色：admin、engineer、operator 均可切换设备状态",
               "入口：设备台账列表行 / 设备档案页 → 右上角「状态切换」按钮"], title_color=ORANGE)
    steps = [
        ("操作员发现机台故障", "机台异常停机、报警或其他问题，需要工程师介入"),
        ("切设备状态为 DOWN", "打开状态切换对话框 → 选 DOWN → 填写变更原因（选填，建议详细记录现象）"),
        ("系统自动创建 REPAIR 工单", "自动写入：设备ID、标题（设备名 + 状态 DOWN）、现象=变更原因、紧急度=HIGH、assignee=派给相关工程师"),
        ("工程师看板立即看到工单", "管理员/工程师看板：DOWN 机数 ↑ + 我的工单列表 + 超时预警计时开始"),
    ]
    _add_step_flow(slide, Emu(457200), Emu(3150000), Emu(11277000), steps,
                   step_colors=[ORANGE, RED, TEAL, GREEN])
    # 7 statuses
    statuses = [
        ("RUN运行", GREEN), ("IDLE待机", RGBColor(0x95,0xA5,0xA6)),
        ("DOWN故障", RED), ("PM维护", ORANGE),
        ("ENG调试", ACCENT), ("PV工艺验证", RGBColor(0x8E,0x44,0xAD)),
        ("OFFLINE离线", GRAY_TEXT),
    ]
    y = Emu(5750000)
    _add_rect(slide, Emu(457200), y, Emu(11277000), Emu(800000), LIGHT_GRAY)
    _add_text(slide, Emu(640080), y + Emu(91440), Emu(3000000), Emu(274320), "其他可选状态：", 12, DARK_BLUE, bold=True)
    x = Emu(3700000)
    for name, color in statuses:
        _add_rect(slide, x, y + Emu(91440), Emu(1400000), Emu(450000), color)
        tf = slide.shapes[-1].text_frame
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); _set_run(run, name, 10, WHITE, bold=True)
        x += Emu(1500000)

def build_workflow_wo_lifecycle(prs):
    """P19: 工单全生命周期 + SLA"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "【流程 3】工单全生命周期 + SLA 超期升级", "Work Order Full Lifecycle & SLA Escalation")
    # Horizontal status flow
    y_top = Emu(1100000)
    status_flow = [
        ("OPEN\n待处理", ACCENT), ("IN_PROGRESS\n处理中", ORANGE),
        ("PENDING_VERIFY\n待验证", DARK_BLUE), ("CLOSED\n已关闭", GREEN),
    ]
    x = Emu(457200)
    for i, (name, color) in enumerate(status_flow):
        # Box
        _add_rect(slide, x, y_top, Emu(2400000), Emu(914400), color)
        tf = slide.shapes[-1].text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for line in name.split("\n"):
            p = tf.paragraphs[0] if not tf.paragraphs or (tf.paragraphs[0].text == "" and len(tf.paragraphs)==1) else tf.add_paragraph()
            if tf.paragraphs[0].text == "" and p == tf.paragraphs[0]:
                pass
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); _set_run(run, line, 13 if line.isupper() or len(name)<15 else 11, WHITE, bold=True)
        if i < len(status_flow) - 1:
            _add_arrow(slide, x + Emu(2400000), y_top + Emu(350000), Emu(500000), Emu(200000), GRAY_TEXT)
        x += Emu(2900000)
    # SLA calculation note
    _add_rect(slide, Emu(457200), Emu(2200000), Emu(11277000), Emu(500000), LIGHT_GRAY)
    _add_text(slide, Emu(640080), Emu(2230000), Emu(11000000), Emu(365760),
              "⏱ SLA 自动重算机制：派工(actual_start)算响应时长 → 关闭(completed_at/actual_end)算解决时长 → 任一超目标即标记 sla_breach=True（超期工单）",
              12, RED, bold=True)
    # Left: detailed steps
    steps_wo = [
        ("创建工单 / 自动派工", "人工创建：选设备+填紧急度；或 DOWN 状态自动创建 → assignee=工程师"),
        ("工程师接单 → IN_PROGRESS", "派工 → 实际响应时长=actual_start - created_at；SLA 目标按紧急度预设"),
        ("维修：5Why分析 + 备件领用", "详情页填：故障现象 → 根因(5Why) → 处置 → 预防 → 领用备件扣减库存"),
        ("转 PENDING_VERIFY 待验证", "维修完成，请求操作员/工艺验证修复效果"),
        ("CLOSED / CANCELLED", "验证通过→CLOSED；否则→IN_PROGRESS继续修；误报/无法修→CANCELLED"),
    ]
    _add_step_flow(slide, Emu(457200), Emu(2850000), Emu(11277000), steps_wo,
                   step_colors=[ACCENT, ORANGE, DARK_BLUE, ACCENT, GREEN])
    # Bottom SLA escalation
    _add_rect(slide, Emu(457200), Emu(5900000), Emu(11277000), Emu(700000), CARD_BG)
    _add_text(slide, Emu(640080), Emu(5930000), Emu(11000000), Emu(365760),
              "⚠ 超时升级：任一 SLA 超期 → 一键「升级指派」更高级别工程师（选升级目标用户+备注+可选改派），自动记录升级人/时间",
              13, ORANGE, bold=True)

def build_workflow_pm_inspection(prs):
    """P20: PM 计划 & 点检巡检"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "【流程 4】PM 维护计划 & 点检巡检", "PM Maintenance Plan & Inspection")
    # Left: PM
    _add_card(slide, Emu(457200), Emu(1050000), Emu(5600000), Emu(5500000),
              "🔄 PM 维护计划（周期性保养）", [], item_size=0)
    steps_pm = [
        ("管理员创建 PM 计划", "选设备 → 周期(周/双周/月/季度/半年/年) → 维护内容清单 → 开始日期"),
        ("系统自动检测到期项", "看板/列表显示「即将到期/已过期」Badge，PM 工程师看板优先推送"),
        ("一键生成到期 PM 工单", "有 pm_plan.generate_due 权限角色 → 批量创建 PM 类型工单"),
        ("执行 → 关单", "派工 → 按维护清单执行 → 记录备件消耗 → 关单"),
    ]
    _add_step_flow(slide, Emu(550000), Emu(1500000), Emu(5400000), steps_pm,
                   step_colors=[DARK_BLUE, ACCENT, TEAL, GREEN])
    # Right: 点检
    _add_card(slide, Emu(6200000), Emu(1050000), Emu(5600000), Emu(5500000),
              "📋 点检巡检（模板 + 记录）", [], item_size=0)
    steps_insp = [
        ("管理员定义点检模板", "检查项名称 / 类型(text/number/选项) / 标准值 / 上下限 → 启用"),
        ("操作员提交点检记录", "点检巡检页 → 选设备 + 选模板 → 逐项填写检查结果 → 提交"),
        ("异常自动联动", "异常项可直接触发 REPAIR 工单或提醒工程师关注"),
        ("历史追溯", "设备 DOWN 后可追溯该设备点检执行记录，判断是否漏检/异常未报告"),
    ]
    _add_step_flow(slide, Emu(6300000), Emu(1500000), Emu(5400000), steps_insp,
                   step_colors=[ORANGE, TEAL, RED, DARK_BLUE])

def build_workflow_doc_control(prs):
    """P21: 文控审批链（三级电子签名 + 状态机）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "【流程 5】文控三级审批链（编制→审核→批准）", "Document 3-Level E-Signature Approval")
    # Status machine diagram at top
    y = Emu(1080000)
    _add_rect(slide, Emu(457200), y, Emu(11277000), Emu(2100000), LIGHT_GRAY)
    _add_text(slide, Emu(640080), y + Emu(91440), Emu(11000000), Emu(274320), "状态机（白名单单向流转）：", 12, DARK_BLUE, bold=True)
    states = [
        ("草稿", GRAY_TEXT), ("审核中", ORANGE), ("生效", GREEN), ("作废", RED)
    ]
    x = Emu(700000)
    for i, (name, color) in enumerate(states):
        _add_circle(slide, x, y + Emu(600000), Emu(1000000), color, name, WHITE, 13)
        if i == 0 or i == 1:  # arrows: 草稿→审核中 审核中→生效
            _add_arrow(slide, x + Emu(1000000), y + Emu(980000), Emu(1600000), Emu(250000), TEAL)
        if i == 1:  # also rejection arrows
            pass
        x += Emu(2600000)
    _add_text(slide, Emu(640080), y + Emu(1700000), Emu(5000000), Emu(274320),
              "✖ 驳回路径：审核中 → 草稿、生效 → 草稿（必须填驳回原因）", 11, RED)
    _add_text(slide, Emu(6400000), y + Emu(1700000), Emu(5200000), Emu(274320),
              "✖ 作废：生效 → 作废（单向终态，不可逆）", 11, RED)

    # 3-level step detail
    y += Emu(2300000)
    _add_card(slide, Emu(457200), y, Emu(3700000), Emu(3300000),
              "① 编制提交（工艺员）",
              ["上传文件 → 填编号/文控分类/复审周期",
               "「提交审核」→ 二次输入密码 + 签署意见",
               "SHA256 签名指纹留痕",
               "状态 → 审核中"], title_color=GRAY_TEXT)
    _add_card(slide, Emu(4250000), y, Emu(3700000), Emu(3300000),
              "② 审核（QA）",
              ["审核内容：合规性 / 完整性 / 格式",
               "「审核」→ 二次密码 + 意见",
               "通过→可批准；驳回→退回草稿",
               "同文档不可重复审核"], title_color=ORANGE)
    _add_card(slide, Emu(8050000), y, Emu(3657000), Emu(3300000),
              "③ 批准生效（QA/管理员）",
              ["「批准」→ 二次密码 + 意见",
               "状态 → 生效 + 写入生效日期",
               "按复审周期计算 next_review_date",
               "同组旧生效版 → 自动作废"], title_color=GREEN)

    # Bottom extras
    y += Emu(3400000)
    extras = [
        ("🔁 修订记录", "字段级 before/after 对比 + 变更原因"),
        ("📤 分发收回", "USER/DEPT 批量分发 → 收回台账，PDF 下载自动盖受控章水印(编号/状态/下载人/日期)"),
        ("🔒 表单审核", "已审核表单锁死，修改必须走「附加修正」审批流程(修正字段+原值+修正值+原因+二次密码)"),
        ("🔔 复审告警", "next_review_date 30 天内 → 黄色 Badge；已过期 → 红色 Badge"),
    ]
    _add_rect(slide, Emu(457200), y, Emu(11277000), Emu(700000), CARD_BG)
    x_e = Emu(640080)
    for (title, desc) in extras:
        _add_text(slide, x_e, y + Emu(91440), Emu(2700000), Emu(274320), title, 11, TEAL, bold=True)
        _add_text(slide, x_e, y + Emu(365760), Emu(2700000), Emu(274320), desc, 10, DARK_TEXT)
        x_e += Emu(2850000)

def build_workflow_forms(prs):
    """P22: 表单模板 + 结构化电子表单"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "【流程 6】表单模板定义 → 电子表单填写 → 导出", "Form Template Definition → E-Form Fill → Export")
    # Left: Admin flow
    _add_card(slide, Emu(457200), Emu(1050000), Emu(5600000), Emu(5500000),
              "🛠 管理员：创建表单模板", [], item_size=0)
    steps_admin = [
        ("新建模板（基本信息）", "菜单：表单模板管理 → 新建 → 填名称/编码/分类(作业记录类/通用类)/适用机台/说明 → 保存"),
        ("定义字段（9 种类型）", "文本/多行文本/数字/下拉/单选/日期/日期时间/时间/是/否；每项配：Key(英文)/显示名/必填/选项列表/单位/上下限/排序"),
        ("上传参考模板文件", "上传 PDF/Excel/图片空白模板 → 操作员填写时可「下载参考对照」"),
        ("启用模板", "状态切为启用 → 操作员可选此模板生成电子表单"),
    ]
    _add_step_flow(slide, Emu(550000), Emu(1500000), Emu(5400000), steps_admin,
                   step_colors=[ORANGE, ACCENT, TEAL, GREEN])
    # Right: Operator flow
    _add_card(slide, Emu(6200000), Emu(1050000), Emu(5600000), Emu(5500000),
              "👷 操作员：填写结构化电子表单", [], item_size=0)
    steps_op = [
        ("新建电子表单", "工艺文件 → 作业记录文件 Tab → 点蓝色「新建电子表单」按钮"),
        ("选择模板", "选择已启用的作业记录模板（按名称/适用机台筛选）"),
        ("填写元数据 + 字段", "关联机台(必填) + 批次号+班次+生产日期 → 逐项填字段(*必填) → 对照参考文件"),
        ("保存 / 提交 → 导出", "保存草稿(可继续改) / 保存并提交 → 结构化存储；支持导出 JSON / CSV"),
    ]
    _add_step_flow(slide, Emu(6300000), Emu(1500000), Emu(5400000), steps_op,
                   step_colors=[DARK_BLUE, ORANGE, ACCENT, GREEN])
    # State transition
    y = Emu(5750000)
    _add_rect(slide, Emu(457200), y, Emu(11277000), Emu(800000), LIGHT_GRAY)
    _add_text(slide, Emu(640080), y + Emu(91440), Emu(11000000), Emu(365760),
              "表单状态：草稿 → 已提交 → 已审核(锁定) → 附加修正(PENDING→APPROVED) | 作废", 12, DARK_BLUE, bold=True)

def build_workflow_8d(prs):
    """P23: 8D 报告 + 双路径归档知识库"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "【流程 7】8D 报告 → 故障知识库一键归档（双路径）", "8D Report → Knowledge Base Archive (Dual Paths)")
    # Left: 8D Report creation flow
    _add_card(slide, Emu(457200), Emu(1050000), Emu(5600000), Emu(2700000),
              "📝 创建 8D 报告（8 步完整填写）", [], item_size=0)
    steps_8d = [
        ("新建 8D（D0 背景）", "选关联设备(必填)+关联工单(可选) → 填报告编号/标题/问题描述(D0)"),
        ("依次填 D1~D8", "D1团队→D2问题定义→D3临时遏制→D4根因→D5永久措施→D6验证→D7预防→D8表彰"),
        ("关闭 8D (CLOSED)", "状态切为 CLOSED → 即可一键归档知识库"),
    ]
    _add_step_flow(slide, Emu(550000), Emu(1500000), Emu(5400000), steps_8d,
                   step_colors=[DARK_BLUE, ORANGE, GREEN])
    # Right: Two archive methods
    _add_card(slide, Emu(6200000), Emu(1050000), Emu(5600000), Emu(2700000),
              "📚 自动字段映射", [], item_size=0)
    mappings = [
        ("D0(问题) + D2(问题定义)  →  symptom 故障现象", TEAL),
        ("D4(根因)  →  root_cause 根因", DARK_BLUE),
        ("D5(永久措施)  →  solution 处置措施", ACCENT),
        ("D7(预防)  →  prevention 预防措施", ORANGE),
        ("记录 source_d8_report_id + source_work_order_id（溯源）", GREEN),
    ]
    _add_bullets(slide, Emu(6400000), Emu(1500000), Emu(5400000), Emu(2200000),
                 [("★ ", t) for t, _ in mappings], size=12, spacing=7)

    # Bottom dual paths
    y = Emu(3900000)
    _add_card(slide, Emu(457200), y, Emu(5600000), Emu(2500000),
              "方式 A：从 8D 列表直接归档（推荐）",
              ["8D 报告列表 → 已关闭报告 → 「归档至知识库」按钮",
               "→ 确认/调整标题 / 故障分类 / 标签(默认 8D,{编号})",
               "→ 点「归档」→ 自动写入知识库表",
               "✅ 不会修改/删除原始 8D 报告"], title_color=TEAL)
    _add_card(slide, Emu(6200000), y, Emu(5600000), Emu(2500000),
              "方式 B：从故障知识库页面发起",
              ["故障知识库 → 「从 8D 报告归档」按钮",
               "→ 下拉选 8D 报告(按编号/标题筛选)",
               "→ 填分类/标签 → 保存",
               "✅ 同 8D 可归档多条目(多角度提炼)"], title_color=ACCENT)
    # + 工单归档 as third path
    y += Emu(2600000)
    _add_rect(slide, Emu(457200), y, Emu(11277000), Emu(700000), CARD_BG)
    _add_text(slide, Emu(640080), y + Emu(91440), Emu(11000000), Emu(274320),
              "➕ 补充：工单归档路径", 12, DARK_BLUE, bold=True)
    _add_text(slide, Emu(640080), y + Emu(365760), Emu(11000000), Emu(274320),
              "故障知识库 → 「从工单归档」 → 输入工单 ID → 自动提取 description→现象 / fault_category→分类 / root_cause→根因 / solution→处置 / prevention→预防",
              11, DARK_TEXT)

def build_workflow_knowledge(prs):
    """P24: 知识库检索 + 相似案例 + 复发追踪"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "【流程 8】知识库检索 · 相似案例 · 复发追踪", "Knowledge Base: Search · Similar Cases · Recurrence Track")
    # 1. Search flow
    steps_search = [
        ("关键词全文检索", "搜索框输入如「漏真空」「刻蚀」→ 模糊匹配 title/现象/根因/处置四字段"),
        ("多维度筛选 + 统计行", "按故障分类/关联设备过滤；顶部统计：总条目数/本月新增/复发合计"),
        ("查看详情 + 相似案例", "点「查看详情」→ 完整信息展示 + views 自动 +1；底部相似案例：优先同设备同分类→退化同设备或同分类"),
    ]
    _add_card(slide, Emu(457200), Emu(1050000), Emu(11277000), Emu(3400000),
              "🔎 检索与相似案例推荐", [], item_size=0)
    _add_step_flow(slide, Emu(550000), Emu(1500000), Emu(11000000), steps_search,
                   step_colors=[DARK_BLUE, TEAL, ACCENT])

    # 2. Recurrence tracking
    y = Emu(4550000)
    _add_card(slide, Emu(457200), y, Emu(11277000), Emu(1900000),
              "🔄 标记复发（追踪同一根因反复出现的故障）", [], item_size=0)
    steps_recur = [
        ("再次发现同根因故障", "在知识列表中定位到对应历史条目（相同根因、相同现象）"),
        ("点「标记复发」 → 二次确认", "二次确认弹窗：确认确实是同一根因再次出现"),
        ("recurrence_count +1 → 看板统计", "列表一眼看到高复发条目 → 推进 PMP/备件升级/工艺参数调优（根因根治）"),
    ]
    _add_step_flow(slide, Emu(550000), y + Emu(400000), Emu(11000000), steps_recur,
                   step_colors=[RED, ORANGE, GREEN])

def build_workflow_backup(prs):
    """P25: 系统配置：备份+加密+异地副本+烟雾还原"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "【流程 9】系统配置：备份 · 加密 · 异地副本 · 烟雾还原", "System Config: Backup · Encrypt · Offsite · Smoke Restore")
    # Backup flow
    _add_card(slide, Emu(457200), Emu(1050000), Emu(5600000), Emu(3400000),
              "📦 备份 & 加密 & 异地副本（3-2-1 策略）", [], item_size=0)
    steps_bu = [
        ("一键备份（系统配置→备份）", "打包 SQLite DB + uploads 附件 → ZIP"),
        ("AES-256 加密（可选）", "Fernet(PBKDF2 20万轮) / 旁路也可用 openssl 离线加密"),
        ("写入异地副本（第二份）", "配置第二备份目录 → 自动复制到 NAS/SMB/U盘 冷备"),
        ("定时计划（自动）", "cron 表达式配置，按日/按周定时自动备份"),
    ]
    _add_step_flow(slide, Emu(550000), Emu(1500000), Emu(5400000), steps_bu,
                   step_colors=[DARK_BLUE, ACCENT, TEAL, GREEN])

    # Smoke restore flow
    _add_card(slide, Emu(6200000), Emu(1050000), Emu(5600000), Emu(3400000),
              "🧪 烟雾还原测试（确保备份能真的还原）", [], item_size=0)
    steps_smoke = [
        ("备份完成后自动执行（后台）", "每次备份完成后，系统自动解压备份包验证"),
        ("Step 1：打开 SQLite 文件", "尝试连接 sqlite → 主表行数 > 0 → 判有效"),
        ("Step 2：校验 uploads 目录", "随机抽样附件文件可正常读取"),
        ("全部通过 → 备份健康", "日志写入 sems.log；任何失败 → 邮件/告警通知管理员"),
    ]
    _add_step_flow(slide, Emu(6300000), Emu(1500000), Emu(5400000), steps_smoke,
                   step_colors=[ORANGE, ACCENT, TEAL, GREEN])

    # Manual restore flow
    y = Emu(4550000)
    _add_card(slide, Emu(457200), y, Emu(11277000), Emu(1900000),
              "🔧 恢复流程（季度演练 SOP 建议执行一次）", [], item_size=0)
    steps_res = [
        ("Step 1：部署一套空 SEMS 环境", "同网段新机器，空 DB 也可以，验证恢复能力"),
        ("Step 2：如加密 → 先用离线解密 .enc → out.zip", "openssl -pbkdf2 -iter 200000 命令或 decrypt 脚本解出 ZIP"),
        ("Step 3：系统配置 → 恢复备份 → 上传 zip", "上传备份文件 → 系统自动还原 DB + uploads 目录"),
        ("Step 4：校验", "设备条数 / 工单数 / 用户列表与生产一致 → 抽查 2-3 份附件能打开"),
    ]
    _add_step_flow(slide, Emu(550000), y + Emu(400000), Emu(11000000), steps_res,
                   step_colors=[GRAY_TEXT, ACCENT, DARK_BLUE, GREEN])


# ════════════════════════════════════════════════════════════════════════════
# Main
# ════════════════════════════════════════════════════════════════════════════
def main():
    orig_path = '/workspace/SEMS_功能介绍_old_backup.pptx'
    out_path = '/workspace/SEMS_功能介绍.pptx'
    # Try to reuse screenshot; fallback to original SEMS .pptx if backup doesn't exist
    screenshot = None
    import os
    for p in ['/workspace/SEMS_功能介绍.pptx', orig_path]:
        if os.path.exists(p):
            try:
                screenshot = _extract_screenshot(p)
                if screenshot: break
            except: pass

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Original 16 slides
    build_cover(prs)                          # P1
    build_why_sems(prs)                       # P2
    build_five_loops(prs)                     # P3
    build_v2026_updates(prs)                  # P4
    build_dashboard_roles(prs)                # P5
    build_device_mgmt(prs, screenshot)        # P6
    build_daily_ops(prs, screenshot)          # P7
    build_process_control(prs, screenshot)    # P8
    build_quality(prs)                        # P9
    build_knowledge_base(prs)                 # P10
    build_lcc_oee(prs)                        # P11
    build_security(prs)                       # P12
    build_deployment(prs)                     # P13
    build_accounts(prs)                       # P14
    build_tech_stack(prs)                     # P15
    build_contact(prs)                        # P16

    # NEW: 9 detailed workflow slides
    build_workflow_login(prs)                 # P17
    build_workflow_down_wo(prs)               # P18
    build_workflow_wo_lifecycle(prs)          # P19
    build_workflow_pm_inspection(prs)         # P20
    build_workflow_doc_control(prs)           # P21
    build_workflow_forms(prs)                 # P22
    build_workflow_8d(prs)                    # P23
    build_workflow_knowledge(prs)             # P24
    build_workflow_backup(prs)                # P25

    prs.save(out_path)
    print(f"Generated {out_path} with {len(prs.slides)} slides (V2026.08 + 9 detailed workflows)")


if __name__ == '__main__':
    main()
