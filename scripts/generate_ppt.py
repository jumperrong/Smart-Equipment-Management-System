#!/usr/bin/env python3
"""Generate SEMS_功能介绍.pptx V2026.08 from SEMS_功能介绍_最新版.md outline."""

import io
import zipfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Design System ──
DARK_BLUE  = RGBColor(0x0B, 0x3D, 0x91)
TEAL       = RGBColor(0x00, 0x96, 0x88)
LIGHT_GRAY = RGBColor(0xF5, 0xF7, 0xFA)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x2C, 0x3E, 0x50)
GRAY_TEXT  = RGBColor(0x7F, 0x8C, 0x8D)
ACCENT     = RGBColor(0x1A, 0x73, 0xE8)
CARD_BG    = RGBColor(0xEE, 0xF2, 0xF7)
FONT       = 'Microsoft YaHei'

# Slide dimensions (16:9 widescreen)
SLIDE_W = Emu(12192000)  # 13.33 in
SLIDE_H = Emu(6858000)   # 7.5 in

# ── Helpers ──

def _set_run(run, text, size, color, bold=False, font_name=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name

def _add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def _add_text(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    _set_run(run, text, size, color, bold)
    return txBox

def _add_bullets(slide, left, top, width, height, items, size=14, color=DARK_TEXT, spacing=6):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        if isinstance(item, tuple):
            prefix, rest = item
            r1 = p.add_run()
            _set_run(r1, prefix, size, TEAL, bold=True)
            r2 = p.add_run()
            _set_run(r2, rest, size, color)
        else:
            r = p.add_run()
            _set_run(r, item, size, color)
    return txBox

def _add_top_bar(slide, title, subtitle=None):
    """Content slide header: blue top bar + title."""
    _add_rect(slide, 0, 0, SLIDE_W, Emu(914400), DARK_BLUE)
    _add_text(slide, Emu(457200), Emu(137160), Emu(9000000), Emu(457200),
              title, 24, WHITE, bold=True)
    if subtitle:
        _add_text(slide, Emu(457200), Emu(548640), Emu(9000000), Emu(274320),
                  subtitle, 13, RGBColor(0xAE, 0xC5, 0xE8))

def _add_card(slide, left, top, width, height, title, items, title_color=DARK_BLUE):
    """A light-gray card with a title and bullet items."""
    _add_rect(slide, left, top, width, height, LIGHT_GRAY)
    _add_text(slide, left + Emu(182880), top + Emu(91440),
              width - Emu(365760), Emu(365760),
              title, 14, title_color, bold=True)
    _add_bullets(slide, left + Emu(182880), top + Emu(457200),
                 width - Emu(365760), height - Emu(548640),
                 items, size=12, spacing=4)

def _extract_screenshot(orig_path):
    """Extract the screenshot image blob from the original PPT."""
    z = zipfile.ZipFile(orig_path)
    for name in z.namelist():
        if name.startswith('ppt/media/') and name.endswith('.png'):
            return z.read(name)
    return None

# ── Slide builders ──

def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_BLUE)
    _add_rect(slide, Emu(4572000), Emu(3200400), Emu(3047695), Emu(38100), TEAL)
    _add_text(slide, Emu(914400), Emu(1828800), Emu(10362895), Emu(914400),
              "SEMS 半导体设备管理系统", 40, WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, Emu(914400), Emu(2833920), Emu(10362895), Emu(548640),
              "Semiconductor Equipment Management System", 20,
              RGBColor(0xAE, 0xC5, 0xE8), align=PP_ALIGN.CENTER)
    _add_text(slide, Emu(914400), Emu(4572000), Emu(10362895), Emu(457200),
              "功能介绍  |  V2026.08  |  2026-08", 18,
              RGBColor(0xAE, 0xC5, 0xE8), align=PP_ALIGN.CENTER)
    _add_text(slide, Emu(914400), Emu(5486400), Emu(10362895), Emu(365760),
              "FastAPI  +  Vue 3  +  Element Plus  +  SQLite", 14,
              RGBColor(0x8A, 0xAA, 0xC8), align=PP_ALIGN.CENTER)


def build_why_sems(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "为什么需要 SEMS？", "Why SEMS?")
    pain_items = [
        "设备状态靠电话/微信群，台账分散在 Excel/纸档",
        "故障处置经验沉淀差，同一根因反复复发",
        "点检/PM/润滑执行难追踪，容易漏做",
        "备件库存不透明，等要修了才发现缺件",
        "工艺文件版本与审批靠邮件，追溯困难",
        "8D 报告写完就丢，知识库与品管工具完全脱节",
    ]
    _add_bullets(slide, Emu(457200), Emu(1280160), Emu(7000000), Emu(4500000),
                 [("• ", p) for p in pain_items], size=15, spacing=10)
    _add_card(slide, Emu(7700000), Emu(1280160), Emu(4000000), Emu(2500000),
              "一句话价值",
              ["把设备台账、日常运维、工艺文控、品管合规、数据价值",
               "五个核心闭环落在一套系统里，",
               "真正做到数据可查、经验可复用。"],
              title_color=TEAL)


def build_five_loops(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "五大业务闭环", "Five Business Loops")
    loops = [
        ("① 设备台账 & 状态闭环", "设备台账 / 状态切换 / 设备全生命周期 T0-T3 / 润滑管理", "机台状态一眼看清，DOWN 自动派工单"),
        ("② 日常运维闭环", "工单 / PM 计划 / 点检巡检 / 备件管理 / 工单 SLA 升级", "从建单→维修→验证→关单一气呵成，超时自动升级"),
        ("③ 工艺文控闭环", "工艺文件 / 表单模板 / 电子表单 / 文控审批链 / 分发收回 / 复审告警", "符合体系标准的受控文控，水印+指纹+留痕"),
        ("④ 品管合规闭环", "8D / FMEA / 安全检查 / 环境核查", "质量分析合规可追溯，证书到期提前告警"),
        ("⑤ 数据价值闭环", "故障知识库(工单归档+8D归档) / OEE / 设备成本 LCC", "故障经验不随人流失，LCC 一眼算清单台成本"),
    ]
    y = Emu(1100000)
    for title, modules, value in loops:
        _add_rect(slide, Emu(457200), y, Emu(11277000), Emu(914400), LIGHT_GRAY)
        _add_text(slide, Emu(640080), y + Emu(91440), Emu(3200000), Emu(365760),
                  title, 14, DARK_BLUE, bold=True)
        _add_text(slide, Emu(3900000), y + Emu(91440), Emu(5000000), Emu(731520),
                  modules, 11, DARK_TEXT)
        _add_text(slide, Emu(9000000), y + Emu(91440), Emu(2600000), Emu(731520),
                  value, 11, TEAL)
        y += Emu(1000000)


def build_v2026_updates(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "V2026.08 版本核心更新", "Key Updates in V2026.08")
    cards = [
        ("📚 8D → 知识库 一键归档", [
            "D0+D2 → 故障现象 / D4 → 根因",
            "D5 → 处置 / D7 → 预防",
            "关联 source_d8_report_id 溯源",
            "同时保留「工单归档」入口",
        ]),
        ("📋 8 大分组分级菜单", [
            "① 总览 ② 设备管理 ③ 运维工单",
            "④ 安全与环境 ⑤ 备件与人员",
            "⑥ 工艺文控 ⑦ 分析改进",
            "⑧ 系统配置（永远在最下方）",
        ]),
        ("🎯 角色定制看板", [
            "6 种角色登录→不同看板组合",
            "QA 优先：文控复审 + 8D 进度",
            "操作员优先：设备状态 + 今日点检",
        ]),
        ("🧾 审计日志入库 + 🔐 bcrypt 修复", [
            "11 类敏感操作写入 audit_logs 表",
            "操作人/目标/IP/UA/时间戳/详情",
            "SHA-256 预哈希解决 bcrypt 72 字节截断",
            "旧值透明回退 + 登录自动升级重哈希",
        ]),
    ]
    positions = [
        (Emu(457200), Emu(1100000), Emu(5600000), Emu(2500000)),
        (Emu(6200000), Emu(1100000), Emu(5600000), Emu(2500000)),
        (Emu(457200), Emu(3700000), Emu(5600000), Emu(2500000)),
        (Emu(6200000), Emu(3700000), Emu(5600000), Emu(2500000)),
    ]
    for (title, items), pos in zip(cards, positions):
        _add_card(slide, pos[0], pos[1], pos[2], pos[3], title, items)


def build_dashboard_roles(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "看板总览 & 角色定制", "Dashboard & Role Customization")
    roles = [
        ("管理员", "SLA 达成率 → 文控复审告警 → 8D 进度 → 知识库新增 → 低库存备件 → 安全检查到期 → 系统健康"),
        ("工程师", "我的工单 → 待验证工单 → 到期 PM → 点检 → 8D 进行中 → 相似故障推荐"),
        ("QA", "审核中文档 → 复审到期 → 附加修正审批 → 8D 报告 → 安全检查 → 表单审核清单"),
        ("操作员", "设备状态总览 → 我的工单 → 今日点检 → 润滑到期 → 常用设备快速入口"),
    ]
    y = Emu(1100000)
    for role, cards in roles:
        _add_rect(slide, Emu(457200), y, Emu(11277000), Emu(914400), LIGHT_GRAY)
        _add_text(slide, Emu(640080), y + Emu(137160), Emu(1500000), Emu(640080),
                  role, 14, DARK_BLUE, bold=True)
        _add_text(slide, Emu(2200000), y + Emu(137160), Emu(9400000), Emu(640080),
                  cards, 12, DARK_TEXT)
        y += Emu(1000000)
    _add_text(slide, Emu(457200), Emu(5400000), Emu(11277000), Emu(457200),
              "支持：明色青绿 / 暗色霓虹 / 跟随系统 三套模式一键切换；暗色模式做了亮度压暗专项补丁，夜班大屏久看不累。",
              13, TEAL, bold=True)


def build_device_mgmt(prs, screenshot=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "① 设备管理 & DOWN→工单自动派发", "Equipment Management")
    if screenshot:
        slide.shapes.add_picture(io.BytesIO(screenshot),
            Emu(365760), Emu(1280000), Emu(5600000), Emu(4800000))
        panel_left = Emu(6200000)
        panel_w = Emu(5600000)
    else:
        panel_left = Emu(457200)
        panel_w = Emu(11277000)
    _add_card(slide, panel_left, Emu(1280000), panel_w, Emu(1400000),
              "设备台账 & 状态切换",
              ["厂区/区域分类、附件上传（说明书/图纸/SOP）",
               "任意角色切 DOWN → 自动创建 REPAIR 工单 → 自动派工程师",
               "状态：RUN / IDLE / DOWN / PM / ENGINEERING / PROCESS_VALIDATION / OFFLINE"])
    _add_card(slide, panel_left, Emu(2800000), panel_w, Emu(1400000),
              "设备全生命周期 T0-T3",
              ["T0 选型(URS/候选供应商) → T1 采购(PO/金额/交付)",
               "T2 安装调试(FAT/SAT) → T3 量产移交(验收结果)",
               "时间线视图一眼看全流程进度"])
    _add_card(slide, panel_left, Emu(4320000), panel_w, Emu(1400000),
              "润滑管理（五定）",
              ["定点/定人/定时/定质/定量",
               "自动推算下次润滑日，到期提醒高亮"])


def build_daily_ops(prs, screenshot=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "② 日常运维闭环 & 工单 SLA", "Daily Operations & Work Order SLA")
    if screenshot:
        slide.shapes.add_picture(io.BytesIO(screenshot),
            Emu(365760), Emu(1280000), Emu(5600000), Emu(4800000))
        panel_left = Emu(6200000)
        panel_w = Emu(5600000)
    else:
        panel_left = Emu(457200)
        panel_w = Emu(11277000)
    _add_card(slide, panel_left, Emu(1280000), panel_w, Emu(1300000),
              "工单管理",
              ["创建/跟踪、紧急度标签、持续时长实时显示、关键词全文检索",
               "详情：5Why 根因分析 + 备件领用 + 状态流转(OPEN→CLOSED)"])
    _add_card(slide, panel_left, Emu(2700000), panel_w, Emu(1300000),
              "SLA & 超期升级",
              ["SLA 目标响应/解决时长（按紧急度预设）",
               "实际时长自动计算 → SLA 达成率统计 → 超期自动升级指派"])
    _add_card(slide, panel_left, Emu(4120000), panel_w, Emu(1600000),
              "PM / 点检 / 备件",
              ["PM：周/双周/月/季 周期，批量生成到期 PM 工单",
               "点检：模板+检查项+历史记录",
               "备件：库存/出入库/低库存告警/设备易损件绑定"])


def build_process_control(prs, screenshot=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "③ 工艺文控闭环", "Process Document Control")
    if screenshot:
        slide.shapes.add_picture(io.BytesIO(screenshot),
            Emu(365760), Emu(1280000), Emu(5600000), Emu(4800000)
        )
        panel_left = Emu(6200000)
        panel_w = Emu(5600000)
    else:
        panel_left = Emu(457200)
        panel_w = Emu(11277000)
    _add_card(slide, panel_left, Emu(1280000), panel_w, Emu(1300000),
              "工艺文件 & 结构化电子表单",
              ["指导性文件：版本管理，生效版唯一，同组旧版自动作废",
               "管理员定义模板（9 种字段类型）→ 操作员动态填写 → JSON/CSV 导出"])
    _add_card(slide, panel_left, Emu(2700000), panel_w, Emu(1500000),
              "文控系统（符合体系标准）",
              ["三级电子签名审批链：编制→审核→批准（二次密码校验 + SHA256 指纹）",
               "状态机白名单：草稿↔审核中→生效→作废",
               "修订记录：字段级 before/after 对比"])
    _add_card(slide, panel_left, Emu(4320000), panel_w, Emu(1400000),
              "受控管理 & 复审告警",
              ["分发收回台账、PDF 自动加盖受控章水印",
               "表单审核锁定 + 附加修正流程",
               "复审周期告警：30 天内到期 / 已过期 Badge"])


def build_quality(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "④ 品管合规闭环", "Quality & Compliance")
    _add_card(slide, Emu(457200), Emu(1100000), Emu(5600000), Emu(2400000),
              "8D 报告",
              [("D1-D8 ", "8 步完整建模（团队→问题→遏制→根因→措施→验证→预防→表彰）"),
               ("关联工单", "（可选），状态：DRAFT / IN_PROGRESS / CLOSED"),
               ("⭐ 一键归档知识库", "：D0+D2→现象 / D4→根因 / D5→处置 / D7→预防")])
    _add_card(slide, Emu(6200000), Emu(1100000), Emu(5600000), Emu(2400000),
              "FMEA & 环境核查",
              [("FMEA", "：失效模式 / 影响 / S×O×D → RPN，RPN>100 高亮"),
               ("环境核查", "：温度/湿度/洁净度/VOC/ESD 参数录入与趋势")])
    _add_card(slide, Emu(457200), Emu(3600000), Emu(11277000), Emu(2400000),
              "安全检查（四类）",
              ["safety_device 安全装置（防护罩/急停/联锁）/ 特种设备 / 环保 / 消防",
               "按频率自动推算下次检查日，30 天内黄底、已过期红字",
               "特种设备证书到期 Badge，检查发现 / 整改措施跟踪"])


def build_knowledge_base(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "⑤ 故障知识库（双路径归档）", "Knowledge Base (Dual-Path Archive)")
    # Table-like layout for dual paths
    y = Emu(1100000)
    paths = [
        ("工单", "知识库→「从工单归档」", "description→现象 / fault_category→分类 / root_cause→根因 / solution→处置 / prevention→预防"),
        ("8D 报告", "知识库→「从 8D 归档」或 8D 列表→「归档至知识库」", "D0+D2→现象 / D4→根因 / D5→处置 / D7→预防"),
        ("手动", "新建条目", "任意填写"),
    ]
    _add_rect(slide, Emu(457200), y, Emu(11277000), Emu(365760), CARD_BG)
    for label, w in [("来源", Emu(1500000)), ("入口", Emu(4500000)), ("自动映射", Emu(5277000))]:
        pass
    _add_text(slide, Emu(640080), y + Emu(45720), Emu(1400000), Emu(274320), "来源", 12, DARK_BLUE, bold=True)
    _add_text(slide, Emu(2100000), y + Emu(45720), Emu(4200000), Emu(274320), "入口", 12, DARK_BLUE, bold=True)
    _add_text(slide, Emu(6400000), y + Emu(45720), Emu(5200000), Emu(274320), "自动映射", 12, DARK_BLUE, bold=True)
    y += Emu(365760)
    for source, entry, mapping in paths:
        _add_rect(slide, Emu(457200), y, Emu(11277000), Emu(548640), LIGHT_GRAY)
        _add_text(slide, Emu(640080), y + Emu(91440), Emu(1400000), Emu(365760), source, 12, TEAL, bold=True)
        _add_text(slide, Emu(2100000), y + Emu(91440), Emu(4200000), Emu(365760), entry, 11, DARK_TEXT)
        _add_text(slide, Emu(6400000), y + Emu(91440), Emu(5200000), Emu(365760), mapping, 11, DARK_TEXT)
        y += Emu(548640)
    _add_card(slide, Emu(457200), Emu(3300000), Emu(5600000), Emu(2700000),
              "相似案例推荐",
              [("同设备 + 同故障分类 ", "优先 → 退化匹配"),
               ("按复发次数/浏览量倒序", ""),
               ("新故障一眼看到以前怎么修的", "")])
    _add_card(slide, Emu(6200000), Emu(3300000), Emu(5600000), Emu(2700000),
              "复发追踪",
              [("同一根因复发 → ", "点「标记复发」→ recurrence_count +1"),
               ("看板/列表一眼看出哪些根因反复出现", ""),
               ("需要重点改进的问题一目了然", "")])


def build_lcc_oee(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "⑤ 数据价值（续）：LCC & OEE", "LCC & OEE Analysis")
    _add_card(slide, Emu(457200), Emu(1100000), Emu(5600000), Emu(4800000),
              "设备成本 LCC（全生命周期成本）",
              [("6 种成本类型：", ""),
               ("  采购 / 维护 / 备件 / 能耗 / 折旧 / 报废", ""),
               ("单设备汇总：", "各类成本占比 + 总成本"),
               ("全设备汇总：", "按类型统计（饼图）+ Top10 高成本设备排名"),
               ("单设备年度趋势：", "按年汇总（折线图）")])
    _add_card(slide, Emu(6200000), Emu(1100000), Emu(5600000), Emu(4800000),
              "OEE 分析",
              [("设备综合效率 = ", "可用率 × 性能率 × 良品率"),
               ("按设备 / 时间维度聚合", ""),
               ("按日/周/月维度趋势图展示", ""),
               ("设备对比分析视图", ""),
               ("OEE 组成拆解，瓶颈定位", "")])


def build_security(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "安全 & 审计 & 灾备", "Security, Audit & Backup")
    _add_card(slide, Emu(457200), Emu(1100000), Emu(5600000), Emu(2400000),
              "安全加固（面向局域网部署）",
              [("JWT 双令牌", "：2h access + 7d refresh，5 分钟前自动续期"),
               ("bcrypt", "(SHA-256 预哈希，修复 72 字节截断)"),
               ("密码 3/4 类字符复杂度 / 失败锁定 5 次/15min", ""),
               ("首次登录强制改密 / 用户名枚举防护", ""),
               ("CSP / X-Frame / CORS 白名单", "")])
    _add_card(slide, Emu(6200000), Emu(1100000), Emu(5600000), Emu(2400000),
              "审计日志入库（审核留痕）",
              [("11 类动作全入库：", ""),
               ("  LOGIN_OK / FAIL / LOCKED / LOGOUT", ""),
               ("  PASSWORD_CHANGED / RESET", ""),
               ("  USER_CREATE / UPDATE / DELETE / UNLOCK", ""),
               ("  RESTORE_BACKUP", ""),
               ("每条含：", "action / actor / target / ip / UA / detail / time")])
    _add_card(slide, Emu(457200), Emu(3600000), Emu(11277000), Emu(2400000),
              "灾备（3-2-1 策略）",
              [("应用内定时备份 + SQLite 热快照 + ZIP + AES-256(Fernet/PBKDF2 20万轮) + NAS/SMB/U盘异地副本", ""),
               ("系统级旁路备份脚本（sh/bat，不依赖后端进程也能跑）", ""),
               ("备份后自动烟雾还原测试：解压→打开 sqlite→查表行数→校验 uploads，永远知道备份能不能真还原", "")])


def build_deployment(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "部署方案选型", "Deployment Options")
    y = Emu(1100000)
    headers = [("场景", 2800000), ("推荐部署", 3000000), ("自启", 1600000), ("崩溃重启", 1800000), ("健康检查", 2077000)]
    x = Emu(457200)
    _add_rect(slide, x, y, Emu(11277000), Emu(365760), CARD_BG)
    for h, w in headers:
        _add_text(slide, x + Emu(91440), y + Emu(45720), Emu(w), Emu(274320), h, 11, DARK_BLUE, bold=True)
        x += Emu(w)
    y += Emu(365760)
    rows = [
        ("工厂 Linux 有 root", "systemd（方案 A）", "✅ 开机自启", "✅ on-failure", "✅ cron 2min→重启"),
        ("普通用户 / 受限 Linux", "watchdog_user.sh（方案 B）", "✅ @reboot", "✅ 每分钟 tick", "✅ 同上"),
        ("Windows Server / 工控机", "NSSM + 任务计划（方案 C）", "✅ 延迟自启", "✅ NSSM + SC failure", "✅ 2min 计划任务"),
        ("容器化 / 快速体验", "Docker Compose（方案 D）", "✅ unless-stopped", "✅ compose restart", "✅ healthcheck"),
    ]
    for row in rows:
        x = Emu(457200)
        _add_rect(slide, x, y, Emu(11277000), Emu(548640), LIGHT_GRAY)
        for val, (_, w) in zip(row, headers):
            _add_text(slide, x + Emu(91440), y + Emu(91440), Emu(w), Emu(365760), val, 10, DARK_TEXT)
            x += Emu(w)
        y += Emu(548640)
    _add_text(slide, Emu(457200), y + Emu(91440), Emu(11277000), Emu(731520),
              "入口级加固：端口占用友好提示 / SQLite WAL checkpoint 钩子 / uvicorn 优雅退出 / 日志双写（journald + 按天滚动 14/30 天）",
              12, TEAL, bold=True)


def build_accounts(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "默认账号 & 演示数据", "Default Accounts & Demo Data")
    accounts = [
        ("admin", "管理员", "admin123", "全功能 / 用户与权限 / 备份与恢复"),
        ("engineer1", "张工 工程师", "eng123", "工单处置 / 8D / 知识库录入"),
        ("process1", "周工艺 工艺员", "proc123", "工艺文件 / 电子表单 / 提交文控审核"),
        ("qa1", "陈品管 QA", "qa123", "文控审核批准 / 8D / FMEA / 安全检查"),
        ("operator1", "王操作 操作员", "op123", "切换设备状态 / 点检 / 电子表单填写"),
        ("viewer1", "孙查看 查看者", "view123", "只读浏览（设备/工单/文档/知识）"),
    ]
    y = Emu(1100000)
    _add_rect(slide, Emu(457200), y, Emu(11277000), Emu(365760), CARD_BG)
    cols = [("账号", 2000000), ("角色", 2500000), ("默认密码", 2000000), ("典型工作内容", 4777000)]
    x = Emu(457200)
    for h, w in cols:
        _add_text(slide, x + Emu(91440), y + Emu(45720), Emu(w), Emu(274320), h, 11, DARK_BLUE, bold=True)
        x += Emu(w)
    y += Emu(365760)
    for acc, role, pwd, work in accounts:
        x = Emu(457200)
        _add_rect(slide, x, y, Emu(11277000), Emu(457200), LIGHT_GRAY)
        for val, (_, w) in zip([acc, role, pwd, work], cols):
            _add_text(slide, x + Emu(91440), y + Emu(91440), Emu(w), Emu(274320), val, 10, DARK_TEXT)
            x += Emu(w)
        y += Emu(457200)
    _add_text(slide, Emu(457200), y + Emu(182880), Emu(11277000), Emu(731520),
              "演示数据开箱即有：10 台设备 + 27 条状态日志 + 16 附件 + 14 备件 + 12 工单 + 4 份 8D 报告 + 14 工艺文件 + 12 安全检查 + 9 润滑点 + 9 条故障知识 + 31 成本记录 等 30+ 类记录。",
              12, TEAL, bold=True)


def build_tech_stack(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "技术栈 & 扩展路线", "Tech Stack & Roadmap")
    _add_card(slide, Emu(457200), Emu(1100000), Emu(5600000), Emu(3500000),
              "当前技术栈",
              [("后端：", "Python 3.10+ / FastAPI / SQLAlchemy 2.0 / SQLite"),
               ("前端：", "Vue 3 / Vite / Element Plus / ECharts / Pinia"),
               ("安全：", "JWT / bcrypt(+SHA256 prehash) / CSP / CORS 白名单"),
               ("备份：", "SQLite hot backup / ZIP / Fernet AES-256"),
               ("打包：", "Docker / systemd / NSSM / PyInstaller")])
    _add_card(slide, Emu(6200000), Emu(1100000), Emu(5600000), Emu(1600000),
              "短期扩展路线（Q4 规划）",
              ["① 审计日志前端界面化（查询/导出 CSV）",
               "② 知识条目审核流（草稿→QA 审核→生效）",
               "③ 8D 报告通知（钉钉/飞书 Webhook）",
               "④ 备品低库存自动建采购申请单"])
    _add_card(slide, Emu(6200000), Emu(2900000), Emu(5600000), Emu(1700000),
              "中期扩展",
              ["对接 MES：设备状态 / 工艺参数自动上传",
               "对接 SCADA：关键机台参数趋势、阈值报警",
               "移动端（微信小程序/钉钉微应用）：点检 / 工单"])


def build_contact(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide, "文档索引 & 更新记录", "Documentation & Changelog")
    _add_card(slide, Emu(457200), Emu(1100000), Emu(5600000), Emu(2500000),
              "文档索引",
              [("README.md", " — 部署 & 架构详情 & 安全说明 & 灾备"),
               ("用户使用教程.md", " — 日常使用操作手册（按角色分步骤）"),
               ("SEMS_功能介绍_最新版.md", " — PPT 文字源（本文件）"),
               ("SEMS_UI设计评审与方案对比.pptx", " — UI 设计方案评审对比")])
    _add_card(slide, Emu(6200000), Emu(1100000), Emu(5600000), Emu(2500000),
              "更新记录",
              [("V2026.08：", "8D→知识库双归档、8 组分级菜单、角色看板、审计日志入库、bcrypt 截断修复"),
               ("V2026.06：", "文控系统(三级签名)、安全检查、润滑管理、工单 SLA、设备生命周期、知识库、LCC"),
               ("V2026.04：", "B+C 双皮肤主题、表单模板结构化电子表单、DOWN→工单自动派发")])
    _add_text(slide, Emu(457200), Emu(4400000), Emu(11277000), Emu(731520),
              "默认账号：admin / admin123（首次登录强制改密）",
              16, DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, Emu(457200), Emu(5200000), Emu(11277000), Emu(731520),
              "感谢观看！",
              28, TEAL, bold=True, align=PP_ALIGN.CENTER)


# ── Main ──

def main():
    orig_path = '/workspace/SEMS_功能介绍.pptx'
    out_path = '/workspace/SEMS_功能介绍.pptx'

    # Extract screenshot from original PPT
    screenshot = _extract_screenshot(orig_path)

    # Create new presentation
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Build all 16 slides
    build_cover(prs)                          # P1
    build_why_sems(prs)                       # P2
    build_five_loops(prs)                     # P3
    build_v2026_updates(prs)                  # P4
    build_dashboard_roles(prs)                # P5
    build_device_mgmt(prs, screenshot)        # P6
    build_daily_ops(prs, screenshot)          # P7
    build_process_control(prs, screenshot)    # P8
    build_quality(prs)                        # P9
    build_knowledge_base(prs)                 # P10
    build_lcc_oee(prs)                        # P11
    build_security(prs)                       # P12
    build_deployment(prs)                     # P13
    build_accounts(prs)                       # P14
    build_tech_stack(prs)                     # P15
    build_contact(prs)                        # P16

    prs.save(out_path)
    print(f"Generated {out_path} with {len(prs.slides)} slides")


if __name__ == '__main__':
    main()
