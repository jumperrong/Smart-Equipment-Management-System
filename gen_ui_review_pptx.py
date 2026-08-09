"""生成 SEMS UI 设计评审 + 3 套方案对比 PPT。
使用 python-pptx，无需嵌入字体——只要目标机器安装了中文 CJK 字体即可正常显示。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

OUT = "/workspace/SEMS_UI设计评审与方案对比.pptx"
IMGDIR = "/workspace/screenshots"

# ---- 色板 ----
BLUE       = RGBColor(0x25, 0x63, 0xEB)
INDIGO     = RGBColor(0x4F, 0x46, 0xE5)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
ORANGE     = RGBColor(0xF9, 0x73, 0x16)
RED        = RGBColor(0xEF, 0x44, 0x44)
CYAN       = RGBColor(0x22, 0xD3, 0xEE)
PURPLE     = RGBColor(0x63, 0x66, 0xF1)
SLATE      = RGBColor(0x0F, 0x17, 0x2A)
MUTED      = RGBColor(0x64, 0x74, 0x8B)
TEXT       = RGBColor(0x1F, 0x29, 0x37)
BORDER     = RGBColor(0xE2, 0xE8, 0xF0)
LIGHT_BG   = RGBColor(0xF8, 0xFA, 0xFC)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)
ROSE       = RGBColor(0xF4, 0x3F, 0x5E)
DARK_BASE  = RGBColor(0x0B, 0x12, 0x20)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None, radius=None):
    if radius is None:
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    else:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shp.adjustments[0] = radius
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w: shp.line.width = line_w
    return shp

def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Microsoft YaHei",
             spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing: p.space_after = spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        # 让中文 fallback 到中文字体
        rPr = r._r.get_or_add_rPr()
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', font)
    return tb

def add_page_banner(slide, title, subtitle="", page_no=None, total=None):
    # 顶部渐变横条
    add_rect(slide, 0, 0, SW, Inches(0.9), fill=SLATE)
    add_rect(slide, 0, Inches(0.82), SW, Inches(0.08), fill=BLUE)
    add_text(slide, Inches(0.5), Inches(0.2), Inches(10), Inches(0.5),
             title, size=24, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.5), Inches(10), Inches(0.3),
                 subtitle, size=11, color=RGBColor(0x94,0xA3,0xB8),
                 anchor=MSO_ANCHOR.MIDDLE)
    if page_no:
        pg = f"{page_no} / {total}" if total else str(page_no)
        add_text(slide, SW - Inches(1.8), Inches(0.3), Inches(1.3), Inches(0.4),
                 pg, size=11, color=RGBColor(0x94,0xA3,0xB8),
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    # 底部
    add_rect(slide, 0, SH - Inches(0.45), SW, Inches(0.45), fill=LIGHT_BG)
    add_rect(slide, 0, SH - Inches(0.45), SW, Inches(0.03), fill=BORDER)
    add_text(slide, Inches(0.5), SH - Inches(0.43), Inches(12), Inches(0.4),
             "SEMS 半导体设备管理系统  ·  UI 设计评审  ·  2026-08-09",
             size=9, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)

def add_bullets(slide, x, y, w, h, items, *, size=13, color=TEXT, bullet="•",
                line_space=6, bold_first=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(line_space)
        r = p.add_run()
        r.text = f"{bullet}  "
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(size)
        r.font.color.rgb = color
        rPr = r._r.get_or_add_rPr()
        ea = rPr.find(qn('a:ea'))
        if ea is None: ea = etree.SubElement(rPr, qn('a:ea'))
        ea.set('typeface', "Microsoft YaHei")
        # 内容，可 bold_first 模式（冒号前加粗）
        if bold_first and "：" in it:
            head, _, rest = it.partition("：")
            r1 = p.add_run(); r1.text = head + "："
            r1.font.name="Microsoft YaHei"; r1.font.size=Pt(size); r1.font.bold=True; r1.font.color.rgb=color
            rPr1 = r1._r.get_or_add_rPr()
            ea1 = rPr1.find(qn('a:ea'))
            if ea1 is None: ea1 = etree.SubElement(rPr1, qn('a:ea'))
            ea1.set('typeface', "Microsoft YaHei")
            r2 = p.add_run(); r2.text = rest
            r2.font.name="Microsoft YaHei"; r2.font.size=Pt(size); r2.font.color.rgb=color
            rPr2 = r2._r.get_or_add_rPr()
            ea2 = rPr2.find(qn('a:ea'))
            if ea2 is None: ea2 = etree.SubElement(rPr2, qn('a:ea'))
            ea2.set('typeface', "Microsoft YaHei")
        else:
            r2 = p.add_run(); r2.text = it
            r2.font.name="Microsoft YaHei"; r2.font.size=Pt(size); r2.font.color.rgb=color
            rPr2 = r2._r.get_or_add_rPr()
            ea2 = rPr2.find(qn('a:ea'))
            if ea2 is None: ea2 = etree.SubElement(rPr2, qn('a:ea'))
            ea2.set('typeface', "Microsoft YaHei")
    return tb

def add_swatch(slide, x, y, color, label, desc):
    add_rect(slide, x, y, Inches(1.0), Inches(0.9), fill=color, radius=0.15)
    # hex 文本（浅底用深色字）
    lightness = sum([color[0], color[1], color[2]]) / 3
    tx = RGBColor(0x0F,0x17,0x2A) if lightness > 140 else WHITE
    add_text(slide, x, y, Inches(1.0), Inches(0.9),
             "#{:02X}{:02X}{:02X}".format(*color),
             size=10, bold=True, color=tx,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x, y + Inches(0.95), Inches(1.0), Inches(0.3),
             label, size=10, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(slide, x - Inches(0.2), y + Inches(1.25), Inches(1.4), Inches(0.6),
             desc, size=9, color=MUTED, align=PP_ALIGN.CENTER)

def add_table(slide, x, y, w, h, data, *, header_fill=SLATE,
              header_color=WHITE, col_widths=None,
              alt_fill=LIGHT_BG, first_col_fill=None,
              first_col_bold=True, cell_size=11, header_size=12,
              cell_color=TEXT):
    rows, cols = len(data), len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    tbl = tbl_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths): tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Emu(60000); cell.margin_right = Emu(60000)
            cell.margin_top  = Emu(40000); cell.margin_bottom = Emu(40000)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(data[r][c])
            run.font.name = "Microsoft YaHei"
            rPr = run._r.get_or_add_rPr()
            ea = rPr.find(qn('a:ea'))
            if ea is None: ea = etree.SubElement(rPr, qn('a:ea'))
            ea.set('typeface', "Microsoft YaHei")
            is_header = (r == 0)
            is_first_col = (c == 0)
            if is_header:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
                run.font.bold = True
                run.font.size = Pt(header_size)
                run.font.color.rgb = header_color
                p.alignment = PP_ALIGN.CENTER
            else:
                if first_col_fill and is_first_col:
                    cell.fill.solid(); cell.fill.fore_color.rgb = first_col_fill
                elif r % 2 == 0 and alt_fill:
                    cell.fill.solid(); cell.fill.fore_color.rgb = alt_fill
                else:
                    cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
                run.font.size = Pt(cell_size)
                run.font.color.rgb = cell_color
                if first_col_bold and is_first_col:
                    run.font.bold = True
    return tbl_shape

def add_image(img_path, slide, x, y, w, h=None, caption=None, caption_color=MUTED):
    if not os.path.exists(img_path): return None
    shp = slide.shapes.add_picture(img_path, x, y, width=w, height=h)
    if caption:
        cw = shp.width
        add_text(slide, x, y + shp.height + Inches(0.05), cw, Inches(0.35),
                 caption, size=10, color=caption_color, align=PP_ALIGN.CENTER)
    return shp

# ========= 页 1 封面 =========
TOTAL = 22
s = prs.slides.add_slide(BLANK)
# 背景
add_rect(s, 0, 0, SW, SH, fill=SLATE)
# 渐变条装饰
add_rect(s, 0, Inches(3.3), SW, Inches(0.06), fill=BLUE)
add_rect(s, 0, Inches(3.38), Inches(4.5), Inches(0.06), fill=GREEN)
add_rect(s, 0, Inches(3.46), Inches(2.8), Inches(0.06), fill=CYAN)

add_text(s, Inches(0.8), Inches(1.4), Inches(12), Inches(0.8),
         "UI / UX 设计评审报告", size=18, bold=True,
         color=RGBColor(0x94,0xA3,0xB8))
add_text(s, Inches(0.8), Inches(2.0), Inches(12), Inches(1.3),
         "SEMS 半导体设备管理系统", size=44, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.7), Inches(12), Inches(0.6),
         "改版方案对比 · 现状评审 · 三套设计样图 · 落地路线",
         size=22, color=BLUE)

# 三张小卡片
cards = [
    (Inches(0.8),  BLUE,   "方案 A", "现代科技蓝", "渐进升级 · 低迁移成本"),
    (Inches(4.5),  GREEN,  "方案 B", "极简青绿洁净风", "Fab 洁净室心理暗示"),
    (Inches(8.2),  PURPLE, "方案 C", "暗色科技大屏风", "车间 55 寸监控屏"),
]
for x, color, tag, title, desc in cards:
    add_rect(s, x, Inches(4.7), Inches(3.4), Inches(1.9), fill=color, radius=0.1)
    add_text(s, x + Inches(0.25), Inches(4.85), Inches(3), Inches(0.4),
             tag, size=11, bold=True, color=WHITE)
    add_text(s, x + Inches(0.25), Inches(5.25), Inches(3), Inches(0.6),
             title, size=20, bold=True, color=WHITE)
    add_text(s, x + Inches(0.25), Inches(5.95), Inches(3), Inches(0.4),
             desc, size=11, color=RGBColor(0xE0,0xE7,0xFF))

add_text(s, Inches(0.8), SH - Inches(0.8), Inches(12), Inches(0.4),
         "评审版本 v1.0   ·   2026-08-09   ·   基于 SEMS v1.2.0 现场采集",
         size=11, color=RGBColor(0x64,0x74,0x8B))

# ========= 页 2 目录 =========
s = prs.slides.add_slide(BLANK)
add_page_banner(s, "报告目录", "Contents Overview", 2, TOTAL)
toc_data = [
    ["01", "一、现状证据采集", "登录页 / 看板 / 设备台账 —— 真实系统 Before 截图"],
    ["02", "二、问题清单与改造方向", "视觉层 / 信息层 / 组件层 —— 17 项具体问题"],
    ["03", "三、方案 A · 现代科技蓝", "色板 / 样图 / 适用场景"],
    ["04", "四、方案 B · 极简青绿洁净风", "色板 / 样图 / 适用场景"],
    ["05", "五、方案 C · 暗色科技大屏风", "色板 / 样图 / 适用场景"],
    ["06", "六、三套方案横向对比", "成本 / 工时 / 可读性 / 契合度 —— 10 项打分表"],
    ["07", "七、落地路线 & 决策建议", "5 阶段实施 + 方案选 A / B / C / B+C 的建议"],
]
for i, (no, title, desc) in enumerate(toc_data):
    y = Inches(1.25) + Inches(0.78) * i
    # 序号框
    add_rect(s, Inches(0.6), y, Inches(0.75), Inches(0.6), fill=BLUE, radius=0.2)
    add_text(s, Inches(0.6), y, Inches(0.75), Inches(0.6),
             no, size=16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.55), y, Inches(4.5), Inches(0.35),
             title, size=17, bold=True, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.55), y + Inches(0.35), Inches(11), Inches(0.35),
             desc, size=11, color=MUTED, anchor=MSO_ANCHOR.TOP)

# ========= 页 3 Before 总览 =========
s = prs.slides.add_slide(BLANK)
add_page_banner(s, "一、现状证据采集（Before）",
                "3 个核心页面真实系统截屏 · 作为改版基线", 3, TOTAL)

shots = [
    ("登录页  ·  #/login",         f"{IMGDIR}/before_login.png",     Inches(0.4), Inches(1.15)),
    ("看板总览  ·  #/dashboard",    f"{IMGDIR}/before_dashboard.png", Inches(0.4), Inches(3.95)),
]
for title, path, x, y in shots:
    add_text(s, x, y, Inches(5.5), Inches(0.35),
             title, size=11, bold=True, color=BLUE)
    add_image(path, s, x, y + Inches(0.35), Inches(5.8))

# 右侧：设备台账
add_text(s, Inches(6.5), Inches(1.15), Inches(6.5), Inches(0.35),
         "设备台账  ·  #/equipment", size=11, bold=True, color=BLUE)
add_image(f"{IMGDIR}/before_equipment.png", s,
          Inches(6.5), Inches(1.5), Inches(6.5))

# ========= 页 4 现状问题 1 视觉层 =========
s = prs.slides.add_slide(BLANK)
add_page_banner(s, "二、现状问题清单（1/3）",
                "视觉层 · 5 个维度逐一分析", 4, TOTAL)
vis_rows = [
    ["维度", "现状", "问题", "改造方向"],
    ["品牌主色", "单一 #409EFF 默认蓝", "未形成 6 色 Design Token 板\n缺少成功/警告/强调色语言", "建立 6 色令牌系统\n+ 状态色语义化"],
    ["侧栏", "深蓝 #001529 无渐变", "与 Element Plus 默认高度雷同\n缺乏品牌辨识度", "3 方案做差异化：\n渐变/纯白/霓虹"],
    ["卡片/阴影", "shadow=never 无阴影", "筛选/表头/数据区层级几乎相同\n信息结构扁平", "引入 2 级阴影体系\n卡片 12px 圆角"],
    ["留白节奏", "统一 18px padding", "KPI 卡与表格同间距\n密度不匹配内容", "分区 24px/卡片 16px/\n表格 12px 三级节奏"],
    ["圆角体系", "Element 默认 4px", "现代产品主流 6-12px 圆角\n亲和力更强", "卡片 12/按钮 8/\n表格 8 三档"],
]
add_table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.9), vis_rows,
          header_fill=BLUE,
          first_col_fill=RGBColor(0xDB,0xEA,0xFE),
          col_widths=[Inches(1.8), Inches(2.6), Inches(4.0), Inches(3.9)],
          cell_size=11)

# ========= 页 5 现状问题 2 信息层 =========
s = prs.slides.add_slide(BLANK)
add_page_banner(s, "二、现状问题清单（2/3）",
                "信息/交互层 · 6 项具体问题", 5, TOTAL)
add_rect(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.9),
         fill=LIGHT_BG, radius=0.05)
issues = [
    ("🔴", "DOWN / PM 顶部方块：视觉重量过大，数字与标签排版松散，缺少图标锚点。"),
    ("🟠", "状态 Tag 仅靠颜色区分（5 色 dark effect），色弱用户辨识度不足，缺形状通道。"),
    ("🟡", "斑马纹 + 粗边框表格：像「后台工具」，与半导体「精密生产系统」气质不一致。"),
    ("🟢", "筛选区 4 下拉 + 1 关键词 + 2 按钮挤在卡片头，换行后对齐非常混乱。"),
    ("🔵", "操作列 4 个文本链接并排：点击区域小、误触率高；未按主次拆分按钮。"),
    ("🟣", "面包屑 / 标题字号同级（16px），页面焦点弱，用户找不到「当前页锚点」。"),
]
for i, (emoji, it) in enumerate(issues):
    y = Inches(1.45) + Inches(0.92) * i
    add_rect(s, Inches(0.9), y, Inches(0.58), Inches(0.58),
             fill=WHITE, radius=0.3)
    add_text(s, Inches(0.9), y, Inches(0.58), Inches(0.58),
             emoji, size=18, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.7), y + Inches(0.06), Inches(10.8), Inches(0.5),
             it, size=14, color=TEXT, anchor=MSO_ANCHOR.TOP)

# ========= 页 6 现状问题 3 组件层 =========
s = prs.slides.add_slide(BLANK)
add_page_banner(s, "二、现状问题清单（3/3）",
                "组件层 · 统一升级方案（三方案均包含）", 6, TOTAL)
comp_rows = [
    ["组件", "现状", "统一改造方向（无论 A/B/C 哪套都做）"],
    ["按钮 Button", "Element 默认 8px 圆角", "主按钮加渐变 + 投影；主次按钮视觉差异加大"],
    ["输入框 Input", "2px 蓝色聚焦边框", "聚焦色跟随方案主色；禁用态灰阶统一"],
    ["状态标签 Tag", "统一 dark effect，仅靠颜色", "RUN=●绿点 / DOWN=▲红感叹 / PM=⏱橙时钟 / ENG=◆蓝菱形\n形状+颜色双通道，色弱人群可读"],
    ["空态 / 加载 Empty", "Element 默认文字", "补充品牌化插画空态；关键列表加骨架屏 Skeleton"],
    ["弹窗 Dialog", "居中白底 + 底部按钮居中", "16px 大圆角 + 底部按钮右对齐 + 危险操作红色突出确认"],
]
add_table(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.9), comp_rows,
          header_fill=SLATE,
          first_col_fill=RGBColor(0x1E,0x29,0x3B),
          first_col_bold=True,
          cell_color=TEXT,
          # 第一列（组件名）要白字
          col_widths=[Inches(2.0), Inches(3.6), Inches(6.7)])
# 手动给第一列（除表头外）字改白色
tbl_shape = s.shapes[-1]
tbl = tbl_shape.table
for r in range(1, len(comp_rows)):
    cell = tbl.cell(r, 0)
    for p in cell.text_frame.paragraphs:
        for run in p.runs:
            run.font.color.rgb = WHITE

# ========= 页 7 方案 A 介绍 + 色板 =========
s = prs.slides.add_slide(BLANK)
add_page_banner(s, "三、方案 A · 现代科技蓝（Modern Blue）",
                "推荐 · 普适 · 常规成本", 7, TOTAL)
# 标签
for i, (txt, clr, x) in enumerate([("★ 推荐", BLUE, Inches(0.6)),
                                    ("普适性", SLATE, Inches(1.85)),
                                    ("低成本", MUTED, Inches(2.9))]):
    add_rect(s, x, Inches(1.15), Inches(1.1), Inches(0.42), fill=clr, radius=0.35)
    add_text(s, x, Inches(1.15), Inches(1.1), Inches(0.42),
             txt, size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.6), Inches(1.7), Inches(6.3), Inches(1.2),
         "风格定位：延续现有蓝色基调，小步优化升级\n"
         "侧栏玻璃态渐变 / KPI 卡加入图标锚点 / 卡片 12px 圆角 + 柔和投影\n"
         "Element CSS 变量改 35% 即可达到 80% 效果 → 成本最低",
         size=13, color=TEXT)

add_text(s, Inches(0.6), Inches(3.15), Inches(6.3), Inches(0.4),
         "▎色板 Design Tokens", size=14, bold=True, color=SLATE)
# 7 个色板
swatchesA = [
    (BLUE,   "Primary",  "主按钮/激活菜单/链接"),
    (GREEN,  "Success",  "RUN 状态/完成"),
    (ORANGE, "Warning",  "PM 状态/预警卡"),
    (RED,    "Danger",   "DOWN 状态/删除确认"),
    (PURPLE, "Info",     "ENG 状态/关键指标"),
    (MUTED,  "Muted",    "次要文字/禁用态"),
    (LIGHT_BG,"Surface", "页面背景 #F1F5F9"),
]
for i, (clr, label, desc) in enumerate(swatchesA):
    col, row = i % 4, i // 4
    x = Inches(0.6) + Inches(1.55) * col
    y = Inches(3.65) + Inches(1.95) * row
    add_swatch(s, x, y, clr, label, desc)

# 右侧样图
add_text(s, Inches(7.15), Inches(1.2), Inches(5.8), Inches(0.35),
         "方案 A 样图 · 看板总览", size=11, bold=True, color=BLUE)
add_image(f"{IMGDIR}/design_scheme_A_modern_blue.jpg", s,
          Inches(7.15), Inches(1.55), Inches(5.85),
          caption="渐变靛青侧栏 + 柔和卡片阴影 + 圆角状态标签 + 图标化 KPI",
          caption_color=MUTED)

# ========= 页 8 方案 A 适配场景 =========
s = prs.slides.add_slide(BLANK)
add_page_banner(s, "三、方案 A · 适配场景与改造范围",
                "适合希望「保留操作习惯、马上能用」的场景", 8, TOTAL)

add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
         "✓ 适合场景", size=18, bold=True, color=BLUE)
scenesA = [
    "希望「保留现有操作习惯、渐进升级」：用户无需重新学习，颜色语言一致但更精致。",
    "主要使用场景为办公室 PC（22-27 寸显示器），不做车间 55 寸大屏投放。",
    "改造预算有限。Element Plus CSS 变量覆盖即可实现 70% 视觉效果，组件级改动量少。",
    "企业有现成 Element Plus 经验 / 历史项目，团队熟悉，落地风险最低。",
]
add_bullets(s, Inches(0.6), Inches(1.8), Inches(12), Inches(2.4), scenesA,
            size=14, bullet="✓", color=TEXT, line_space=8)

add_text(s, Inches(0.6), Inches(4.3), Inches(12), Inches(0.5),
         "▎改造范围（覆盖页面）", size=15, bold=True, color=SLATE)
scope = [
    ["类别", "模块（按优先级降序）"],
    ["P0 · 必须覆盖", "登录页 / 看板总览 / 设备台账 / 工单管理 / 侧栏+顶栏（全局）"],
    ["P1 · 建议覆盖", "PM维护计划 / 工艺文件 / 表单模板管理 / OEE分析 / 系统配置"],
    ["P2 · 可以延后", "点检巡检 / 备件管理 / 品管工具 / 环境核查 / 人员管理 / 资产管理"],
]
add_table(s, Inches(0.6), Inches(4.85), Inches(12.1), Inches(2.1), scope,
          header_fill=BLUE, first_col_fill=RGBColor(0xDB,0xEA,0xFE),
          col_widths=[Inches(2.4), Inches(9.7)], cell_size=12)

# ========= 页 9 方案 B =========
s = prs.slides.add_slide(BLANK)
add_page_banner(s, "四、方案 B · 极简青绿洁净风（Cleanroom Minimal）",
                "洁净室风格 · 低视觉噪音 · 阅读舒适", 9, TOTAL)
for i, (txt, clr, x) in enumerate([("★ Fab 洁净室", GREEN, Inches(0.6)),
                                    ("低视觉噪音", CYAN, Inches(2.0)),
                                    ("舒适阅读", PURPLE, Inches(3.45))]):
    add_rect(s, x, Inches(1.15), Inches(1.3), Inches(0.42), fill=clr, radius=0.35)
    add_text(s, x, Inches(1.15), Inches(1.3), Inches(0.42),
             txt, size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.6), Inches(1.7), Inches(6.3), Inches(1.2),
         "风格定位：对标半导体洁净车间心理暗示\n"
         "纯白基调 + 薄荷绿主色 + 大量留白 + 无斑马纹无粗边框\n"
         "绿色在 Fab 语义 =「RUN / 正常 / 绿灯」，契合行业习惯",
         size=13, color=TEXT)

add_text(s, Inches(0.6), Inches(3.15), Inches(6.3), Inches(0.4),
         "▎色板 Design Tokens", size=14, bold=True, color=SLATE)
swatchesB = [
    (GREEN,  "Primary",  "主按钮 / RUN状态 / 激活菜单"),
    (CYAN,   "Info",     "ENG状态 / 次级链接"),
    (AMBER,  "Warning",  "PM 状态（琥珀橙洁净面板）"),
    (ROSE,   "Danger",   "DOWN 状态 / 危险操作（玫瑰红更柔和）"),
    (TEXT,   "Text",     "正文深灰（非纯黑，不刺眼）"),
    (LIGHT_BG, "Surface","背景近白 / 卡片纯白"),
    (RGBColor(0xF0,0xFD,0xF4), "Tint", "选中行 / 成功提示底"),
]
for i, (clr, label, desc) in enumerate(swatchesB):
    col, row = i % 4, i // 4
    x = Inches(0.6) + Inches(1.55) * col
    y = Inches(3.65) + Inches(1.95) * row
    add_swatch(s, x, y, clr, label, desc)

add_text(s, Inches(7.15), Inches(1.2), Inches(5.8), Inches(0.35),
         "方案 B 样图 · 洁净室美学", size=11, bold=True, color=GREEN)
add_image(f"{IMGDIR}/design_scheme_B_minimal_green.jpg", s,
          Inches(7.15), Inches(1.55), Inches(5.85),
          caption="纯白侧栏 + 3px 绿激活条 + 数字点式 KPI + 超轻量级表格线",
          caption_color=MUTED)

# ========= 页 10 方案 B 场景 =========
s = prs.slides.add_slide(BLANK)
add_page_banner(s, "四、方案 B · 适配场景与元素改造亮点",
                "适合 Fab 工艺员/工程师长期使用场景", 10, TOTAL)

add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
         "✓ 适合场景", size=18, bold=True, color=GREEN)
scenesB = [
    "半导体 Fab 现场工程师 / 工艺员长期使用，长时间阅读疲劳最小。",
    "希望 UI 传递「洁净 / 精密 / 可信赖」品牌感知，契合半导体行业心理暗示。",
    "公司设计语言已在走极简风（飞书 / Notion / Linear 爱好者），内部审美接受度高。",
    "可以接受「与现状有一定视觉差异」，愿意花 1-2 天对用户做培训沟通。",
]
add_bullets(s, Inches(0.6), Inches(1.8), Inches(12), Inches(2.4), scenesB,
            size=14, bullet="✓", color=TEXT, line_space=8)

add_text(s, Inches(0.6), Inches(4.3), Inches(12), Inches(0.5),
         "▎改造亮点（与方案 A 的差异）", size=15, bold=True, color=SLATE)
diff = [
    ["模块", "方案 A（现代蓝）", "方案 B（极简青绿）"],
    ["侧栏", "深蓝 → 靛青 渐变玻璃态", "纯白 + 仅 3px 激活色条"],
    ["背景", "#F1F5F9 冷灰", "#FAFAFA 近白 + 卡片纯白"],
    ["表格", "轻斑马纹 / 卡片包裹", "无斑马纹、仅 1px 分隔线、卡片无边框"],
    ["KPI 卡", "渐变底 + 图标 + 阴影", "仅色点 + 数字 + 标签（大量留白）"],
    ["按钮", "实色渐变 + 投影", "圆角 16px 胶囊状，主色薄荷绿填充"],
]
add_table(s, Inches(0.6), Inches(4.85), Inches(12.1), Inches(2.3), diff,
          header_fill=GREEN,
          first_col_fill=RGBColor(0xD1,0xFA,0xE5),
          col_widths=[Inches(1.6), Inches(5.2), Inches(5.3)],
          cell_size=11.5)

# ========= 页 11 方案 C =========
s = prs.slides.add_slide(BLANK)
add_page_banner(s, "五、方案 C · 暗色科技大屏风（Dark Control Room）",
                "55寸大屏适配 · 霓虹发光 · 高对比度", 11, TOTAL)
for i, (txt, clr, x) in enumerate([("55寸大屏", PURPLE, Inches(0.6)),
                                    ("霓虹发光", CYAN, Inches(2.1)),
                                    ("高对比度", RED, Inches(3.45))]):
    add_rect(s, x, Inches(1.15), Inches(1.3), Inches(0.42), fill=clr, radius=0.35)
    add_text(s, x, Inches(1.15), Inches(1.3), Inches(0.42),
             txt, size=11, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.6), Inches(1.7), Inches(6.3), Inches(1.2),
         "风格定位：生产车间 55 寸监控大屏 / 指挥中心\n"
         "深海军蓝黑基底 + 霓虹青(#22D3EE) + 霓虹紫(#6366F1) 双强调色\n"
         "工业 SCADA 控制台质感；DOWN 呼吸脉冲红色边框；图表霓虹渐变线",
         size=12.5, color=TEXT)

add_text(s, Inches(0.6), Inches(3.15), Inches(6.3), Inches(0.4),
         "▎色板 Design Tokens", size=14, bold=True, color=SLATE)
swatchesC = [
    (DARK_BASE, "Base",     "全局背景（替代白）"),
    (SLATE,     "Surface",  "卡片/侧栏玻璃态叠加"),
    (CYAN,      "Neon Cyan","激活菜单/图表线/ENG状态"),
    (PURPLE,    "Neon Pur.","LOGO/主按钮/强强调"),
    (GREEN,     "Neon Green","RUN状态/正常KPI"),
    (RED,       "Neon Red", "DOWN警报/脉冲闪烁"),
    (AMBER,     "Neon Amber","PM/预警KPI"),
]
for i, (clr, label, desc) in enumerate(swatchesC):
    col, row = i % 4, i // 4
    x = Inches(0.6) + Inches(1.55) * col
    y = Inches(3.65) + Inches(1.95) * row
    add_swatch(s, x, y, clr, label, desc)

add_text(s, Inches(7.15), Inches(1.2), Inches(5.8), Inches(0.35),
         "方案 C 样图 · 控制大屏风格", size=11, bold=True, color=PURPLE)
add_image(f"{IMGDIR}/design_scheme_C_dark_neon.jpg", s,
          Inches(7.15), Inches(1.55), Inches(5.85),
          caption="霓虹发光状态卡 + 暗色玻璃态表格 + 55寸大屏远距离可读",
          caption_color=MUTED)

# ========= 页 12 方案 C 场景 =========
s = prs.slides.add_slide(BLANK)
add_page_banner(s, "五、方案 C · 适配场景与双皮肤架构建议",
                "适合车间大屏 + 值班岗双皮肤架构", 12, TOTAL)

add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
         "✓ 适合场景", size=18, bold=True, color=PURPLE)
scenesC = [
    "厂区指挥中心 / 生产线 55 寸监控大屏展示，远距离 3-5 米外需要高对比度可读。",
    "需要 24 小时亮屏：深色背景 OLED 屏幕更省电、像素灼烧风险显著降低。",
    "企业品牌科技感强，希望管理系统视觉像「工业 4.0 控制台」而不是「OA 后台」。",
    "现场值班岗 + 办公区工程师并行使用的场景（建议 B+C 双皮肤架构）。",
]
add_bullets(s, Inches(0.6), Inches(1.8), Inches(12), Inches(2.4), scenesC,
            size=14, bullet="✓", color=TEXT, line_space=8)

add_text(s, Inches(0.6), Inches(4.3), Inches(12), Inches(0.5),
         "▎推荐：B + C 双皮肤并行架构（仅额外 1-2 人日）", size=15, bold=True, color=SLATE)
# 架构图
add_rect(s, Inches(0.6), Inches(4.9), Inches(12.1), Inches(2.2),
         fill=SLATE, radius=0.05)
add_text(s, Inches(0.9), Inches(5.0), Inches(4), Inches(0.4),
         "顶层：一个 localStorage 开关按钮  ☾",
         size=13, bold=True, color=WHITE)
# 三个块
blocks = [
    (Inches(0.9),  GREEN,  "默认皮肤 B\n极简青绿\n\n办公 PC / 日常操作",   RGBColor(0x11,0x18,0x27)),
    (Inches(4.7),  PURPLE, "CSS 变量切换层\nElement Plus tokens 一键切换",   RGBColor(0x11,0x18,0x27)),
    (Inches(8.5),  CYAN,   "大屏模式皮肤 C\n暗色霓虹\n\n车间 55 寸 / 值班岗", DARK_BASE),
]
for bx, clr, t, cbg in blocks:
    add_rect(s, bx, Inches(5.45), Inches(3.5), Inches(1.45), fill=cbg, line=clr, line_w=Pt(1.5), radius=0.06)
    add_rect(s, bx, Inches(5.45), Inches(3.5), Inches(0.08), fill=clr)
    add_text(s, bx + Inches(0.2), Inches(5.55), Inches(3.1), Inches(1.3),
             t, size=11.5, color=WHITE, anchor=MSO_ANCHOR.TOP)

# ========= 页 13 对比表（上半） =========
s = prs.slides.add_slide(BLANK)
add_page_banner(s, "六、三套方案横向对比（1/2）",
                "10 项维度打分 · 表格化一目了然", 13, TOTAL)
cmp_rows1 = [
    ["维度", "方案 A · 现代蓝", "方案 B · 极简青绿", "方案 C · 暗色霓虹"],
    ["迁移成本",     "★☆☆ 低",      "★★☆ 中",       "★★★ 高"],
    ["开发工时(估)", "6 - 8 人日",   "10 - 14 人日",  "15 - 20 人日"],
    ["PC办公可读性", "★★★★★",       "★★★★★",        "★★★☆☆"],
    ["车间大屏可读", "★★★☆☆",       "★★★★☆",        "★★★★★"],
    ["长期用疲劳度", "较低",         "最低",          "中低（OLED屏幕）"],
    ["与现状差异度", "小（渐进优化）", "中（换色板+留白）", "大（彻底换视觉语言）"],
]
add_table(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(5.9), cmp_rows1,
          header_fill=SLATE,
          first_col_fill=LIGHT_BG,
          col_widths=[Inches(2.2), Inches(3.4), Inches(3.4), Inches(3.3)],
          cell_size=12.5)
# 分块填色
tbl = s.shapes[-1].table
for r in range(1, len(cmp_rows1)):
    tbl.cell(r, 1).fill.solid(); tbl.cell(r, 1).fill.fore_color.rgb = RGBColor(0xDB,0xEA,0xFE)
    tbl.cell(r, 2).fill.solid(); tbl.cell(r, 2).fill.fore_color.rgb = RGBColor(0xD1,0xFA,0xE5)
    tbl.cell(r, 3).fill.solid(); tbl.cell(r, 3).fill.fore_color.rgb = RGBColor(0x1E,0x29,0x3B)
    for p in tbl.cell(r, 3).text_frame.paragraphs:
        for run in p.runs: run.font.color.rgb = WHITE

# ========= 页 14 对比表（下半） =========
s = prs.slides.add_slide(BLANK)
add_page_banner(s, "六、三套方案横向对比（2/2）",
                "色弱友好度 / 行业契合 / 决策建议", 14, TOTAL)
cmp_rows2 = [
    ["维度", "方案 A · 现代蓝", "方案 B · 极简青绿", "方案 C · 暗色霓虹"],
    ["色弱人群辨识度", "良好（颜色+形状双通道）", "最佳（对比度 + 通道）", "良好（霓虹+形状）"],
    ["Fab 行业契合度", "通用企业后台风", "★★★★★ 高度契合（洁净室）", "★★★★★ 高度契合（监控大屏）"],
    ["Element 变量改动率", "~35%", "~55%", "~75%"],
    ["单皮肤推荐指数", "★★★★☆", "★★★★★", "★★★☆☆"],
    ["双皮肤搭档方案", "不建议搭配", "作为默认办公皮肤", "作为大屏模式皮肤"],
    ["搭配决策建议", "小步快跑 首选\n改得少 马上用", "洁净类 Fab 首选\n办公 PC 日常使用", "有大屏监控需求 /\n与 B 组合为双皮肤"],
]
add_table(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(5.9), cmp_rows2,
          header_fill=SLATE,
          first_col_fill=LIGHT_BG,
          col_widths=[Inches(2.4), Inches(3.3), Inches(3.3), Inches(3.3)],
          cell_size=12)
tbl = s.shapes[-1].table
for r in range(1, len(cmp_rows2)):
    tbl.cell(r, 1).fill.solid(); tbl.cell(r, 1).fill.fore_color.rgb = RGBColor(0xDB,0xEA,0xFE)
    tbl.cell(r, 2).fill.solid(); tbl.cell(r, 2).fill.fore_color.rgb = RGBColor(0xD1,0xFA,0xE5)
    tbl.cell(r, 3).fill.solid(); tbl.cell(r, 3).fill.fore_color.rgb = RGBColor(0x1E,0x29,0x3B)
    for p in tbl.cell(r, 3).text_frame.paragraphs:
        for run in p.runs: run.font.color.rgb = WHITE

# ========= 页 15 三方案样图并排 =========
s = prs.slides.add_slide(BLANK)
add_page_banner(s, "六、三套样图并排对比页",
                "A（左）/ B（中）/ C（右）看板样图视觉对比", 15, TOTAL)

triples = [
    (Inches(0.3),  f"{IMGDIR}/design_scheme_A_modern_blue.jpg",
     BLUE,   "方案 A · 现代蓝", "渐变靛青 + 柔和阴影"),
    (Inches(4.55), f"{IMGDIR}/design_scheme_B_minimal_green.jpg",
     GREEN,  "方案 B · 极简青绿", "洁净白 + 薄荷绿点缀"),
    (Inches(8.8),  f"{IMGDIR}/design_scheme_C_dark_neon.jpg",
     PURPLE, "方案 C · 暗色霓虹", "深海军黑底 + 霓虹发光"),
]
for x, path, clr, title, sub in triples:
    add_rect(s, x, Inches(1.15), Inches(4.2), Inches(0.45), fill=clr, radius=0.03)
    add_text(s, x + Inches(0.15), Inches(1.15), Inches(3.9), Inches(0.45),
             title, size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_image(path, s, x, Inches(1.65), Inches(4.2))
    add_text(s, x, Inches(6.3), Inches(4.2), Inches(0.45),
             sub, size=11, color=MUTED, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

# ========= 页 16 实施阶段 1-3 =========
s = prs.slides.add_slide(BLANK)
add_page_banner(s, "七、落地实施路线（1/2）",
                "5 个阶段 · 线性推进 · 阶段 1 可做 Demo 走查", 16, TOTAL)
phA = [
    ("P1", "阶段 1\n(1-2 天)", BLUE,
     "设计令牌 + 变量覆盖\n\n"
     "① 建立 styles/design-tokens.css\n"
     "    (色板/圆角/阴影/间距令牌)\n"
     "② 全局覆盖 Element Plus CSS 变量\n"
     "③ 登录页 + 侧栏 + 顶栏先行\n"
     "    → 产出视觉走查 Demo"),
    ("P2", "阶段 2\n(2-4 天)", INDIGO,
     "组件级样式升级\n\n"
     "① 表格统一（行高/hover/斑马纹）\n"
     "② 状态标签/按钮/KPI 卡组件化\n"
     "③ 卡片阴影/圆角/内边距规范\n"
     "→ 形成 1 页「组件样例」作为开发标准"),
    ("P3", "阶段 3\n(1-8 天)", PURPLE,
     "页面逐个适配\n\n"
     "看板 → 台账 → 工单 → 工艺文件\n"
     "→ 表单模板 → OEE → 品管\n"
     "→ 环境核查 → 人员 → 资产 → 配置\n"
     "共 14 模块按优先级顺序覆盖"),
]
for i, (tag, title, clr, body) in enumerate(phA):
    x = Inches(0.5) + Inches(4.25) * i
    # 大框
    add_rect(s, x, Inches(1.15), Inches(4.0), Inches(6.0),
             fill=WHITE, line=BORDER, line_w=Pt(1), radius=0.04)
    # 顶部色带
    add_rect(s, x, Inches(1.15), Inches(4.0), Inches(0.9), fill=clr)
    # 小 P 标签
    add_rect(s, x + Inches(0.2), Inches(1.3), Inches(0.7), Inches(0.55),
             fill=WHITE, radius=0.25)
    add_text(s, x + Inches(0.2), Inches(1.3), Inches(0.7), Inches(0.55),
             tag, size=14, bold=True, color=clr,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(1.0), Inches(1.3), Inches(2.9), Inches(0.55),
             title.replace("\n", " "), size=14, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.25), Inches(2.2), Inches(3.5), Inches(4.8),
             body, size=12.5, color=TEXT, anchor=MSO_ANCHOR.TOP)

# 箭头连接
for i in range(2):
    ax = Inches(0.5) + Inches(4.25) * (i+1) - Inches(0.2)
    ay = Inches(4.2)
    add_text(s, ax, ay, Inches(0.4), Inches(0.5),
             "➜", size=18, bold=True, color=MUTED,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ========= 页 17 实施阶段 4-5 + 总览工时 =========
s = prs.slides.add_slide(BLANK)
add_page_banner(s, "七、落地实施路线（2/2）",
                "阶段 4-5 + 四档方案工时总览表", 17, TOTAL)

phB = [
    ("P4", "阶段 4\n(1-2 天)", GREEN,
     "空态 / 骨架 / 动画 + 回归测试\n\n"
     "① 加载骨架屏 + 404 插画 + 空态插画\n"
     "② 路由过渡动画 / 卡片微交互\n"
     "③ Chrome / Safari / Edge\n"
     "   + 大屏 (2K / 4K) 分辨率回归测试"),
    ("P5", "阶段 5\n(0.5-1 天)", AMBER,
     "大屏模式开关（选分支 2 时）\n\n"
     "① 顶栏 ☾ 切换按钮\n"
     "② 用 CSS 变量 + localStorage 记忆皮肤\n"
     "③ 与默认皮肤并行不冲突"),
]
for i, (tag, title, clr, body) in enumerate(phB):
    x = Inches(0.5) + Inches(6.4) * i
    add_rect(s, x, Inches(1.15), Inches(6.1), Inches(2.65),
             fill=WHITE, line=BORDER, line_w=Pt(1), radius=0.04)
    add_rect(s, x, Inches(1.15), Inches(6.1), Inches(0.75), fill=clr)
    add_rect(s, x + Inches(0.2), Inches(1.25), Inches(0.7), Inches(0.5),
             fill=WHITE, radius=0.25)
    add_text(s, x + Inches(0.2), Inches(1.25), Inches(0.7), Inches(0.5),
             tag, size=12, bold=True, color=clr,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(1.0), Inches(1.25), Inches(4.9), Inches(0.55),
             title.replace("\n", " "), size=13, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.25), Inches(2.0), Inches(5.6), Inches(1.7),
             body, size=12, color=TEXT, anchor=MSO_ANCHOR.TOP)

# 工时总览表
add_text(s, Inches(0.6), Inches(4.0), Inches(12), Inches(0.5),
         "▎四档方案工时总览（单位：人日）", size=16, bold=True, color=SLATE)
effort = [
    ["方案组合", "阶段1", "阶段2", "阶段3", "阶段4", "阶段5", "合计"],
    ["A · 现代蓝（单皮肤）", "1", "2", "2", "1", "—", "6"],
    ["B · 极简青绿（单皮肤）", "2", "3", "5", "2", "—", "12"],
    ["C · 暗色霓虹（单皮肤）", "2", "4", "8", "3", "1", "18"],
    ["B + C · 双皮肤并行 ★", "3", "5", "10", "3", "2", "23"],
]
eff_shape = add_table(s, Inches(0.6), Inches(4.55), Inches(12.1), Inches(2.7), effort,
          header_fill=BLUE,
          first_col_fill=RGBColor(0xDB,0xEA,0xFE),
          col_widths=[Inches(3.6), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.4), Inches(1.5)],
          cell_size=13)
tbl = eff_shape.table
# 合计列 高亮
for r in range(1, len(effort)):
    cell = tbl.cell(r, -1)
    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFE,0xF3,0xC7)
    for p in cell.text_frame.paragraphs:
        for run in p.runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x92,0x40,0x0E)

# ========= 页 18 推荐决策路径 =========
s = prs.slides.add_slide(BLANK)
add_page_banner(s, "七、推荐决策路径（Decision Tree）",
                "根据用户群体 + 部署场景，两条分支快速定位", 18, TOTAL)

# 分支 1
add_rect(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(2.8),
         fill=RGBColor(0xDB,0xEA,0xFE), line=BLUE, line_w=Pt(1.5), radius=0.06)
add_rect(s, Inches(0.6), Inches(1.2), Inches(0.18), Inches(2.8), fill=BLUE)
add_text(s, Inches(0.95), Inches(1.3), Inches(11), Inches(0.55),
         "分支 1 · 单皮肤普适（80% 客户选此）",
         size=18, bold=True, color=BLUE)
add_text(s, Inches(0.95), Inches(1.85), Inches(5.5), Inches(1.8),
         "🔹 更重视「改得少、马上用」\n"
         "    → 选择 【方案 A 现代蓝】\n"
         "    工时：6-8 人日\n"
         "    风险：最低，Element 变量覆盖即可",
         size=14, color=TEXT)
# 分隔线
add_rect(s, Inches(6.75), Inches(1.9), Inches(0.02), Inches(1.9), fill=BLUE)
add_text(s, Inches(6.95), Inches(1.85), Inches(5.5), Inches(1.8),
         "🔹 更重视「洁净感/品牌感/久读」\n"
         "    → 选择 【方案 B 极简青绿】\n"
         "    工时：10-14 人日\n"
         "    效果：Fab 洁净氛围行业最佳契合",
         size=14, color=TEXT)

# 分支 2
add_rect(s, Inches(0.6), Inches(4.2), Inches(12.1), Inches(3.0),
         fill=RGBColor(0xE0,0xE7,0xFF), line=PURPLE, line_w=Pt(1.5), radius=0.06)
add_rect(s, Inches(0.6), Inches(4.2), Inches(0.18), Inches(3.0), fill=PURPLE)
add_text(s, Inches(0.95), Inches(4.3), Inches(11), Inches(0.55),
         "分支 2 · 双皮肤并行（办公 + 车间大屏都有）",
         size=18, bold=True, color=PURPLE)
add_text(s, Inches(0.95), Inches(4.9), Inches(11.5), Inches(2.2),
         "✓ 推荐组合：【默认皮肤 B 极简青绿】 + 【大屏模式 C 暗色霓虹】\n"
         "    总工时：23 人日\n\n"
         "✓ 切换方式：顶栏「☾ 大屏模式」按钮一键切换，localStorage 记忆，无需重新登录\n\n"
         "✓ 覆盖人群：工艺员/工程师（B）+ 车间值班岗/指挥中心大屏（C） 两类都满意",
         size=14, color=TEXT)

# ========= 页 19 统一改造项 =========
s = prs.slides.add_slide(BLANK)
add_page_banner(s, "七、三方案均包含的「统一保留改造」",
                "无论选择哪套，都会先做的 4 件事", 19, TOTAL)
unified = [
    ("①", BLUE,   "状态标签升级：形状+颜色双通道",
     "RUN = ● 绿点 + 文字\n"
     "DOWN = ▲ 红色感叹三角形 + 脉冲\n"
     "PM   = ⏱ 橙色时钟 + 过程进度\n"
     "ENG  = ◆ 蓝色菱形\n"
     "→ 色弱用户也能 100% 区分 5 种状态"),
    ("②", INDIGO, "顶部 KPI 卡图标化 + 层次化",
     "DOWN 卡 / PM 卡：加入图标锚点、渐变背景、悬停微浮起\n"
     "数字从「大块」改为「数字+单位/图标双行」，节奏更高级\n"
     "DOWN 卡有红边警告闪烁（非频繁，避免打扰）"),
    ("③", GREEN,  "筛选区布局标准化",
     "改为「分组卡片」模式：\n"
     "筛选控件区居左对齐（自动换行整齐）\n"
     "操作按钮区居右对齐（查询/重置/新增操作组）\n"
     "14 模块筛选区统一此规范，不再杂乱"),
    ("④", ORANGE, "表格交互 + 操作列分级",
     "行高 +4px，新增整行 hover 高亮\n"
     "操作列按钮分级：档案=信息链接 / 编辑=主按钮 / 切换=次按钮 / 删除=danger 链接\n"
     "列表加入点击行可查看的 hover 提示，视觉点击范围更大"),
]
for i, (no, clr, title, body) in enumerate(unified):
    col, row = i % 2, i // 2
    x = Inches(0.5) + Inches(6.3) * col
    y = Inches(1.15) + Inches(3.05) * row
    add_rect(s, x, y, Inches(6.1), Inches(2.9),
             fill=WHITE, line=BORDER, line_w=Pt(1), radius=0.05)
    # 头色条
    add_rect(s, x, y, Inches(6.1), Inches(0.9), fill=clr)
    add_rect(s, x + Inches(0.2), y + Inches(0.15), Inches(0.6), Inches(0.6),
             fill=WHITE, radius=0.5)
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(0.6), Inches(0.6),
             no, size=18, bold=True, color=clr,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.95), y + Inches(0.15), Inches(5), Inches(0.6),
             title, size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.25), y + Inches(1.05), Inches(5.6), Inches(1.8),
             body, size=12.5, color=TEXT)

# ========= 页 20 下一步 =========
s = prs.slides.add_slide(BLANK)
add_page_banner(s, "七、下一步（等待您的决策）",
                "3 个决策入口 · 选定后立即进入阶段 1", 20, TOTAL)

steps = [
    ("STEP 1", Inches(0.6), BLUE,
     "选择方案",
     ["A. 现代蓝（6人日，小步快跑）",
      "B. 极简青绿（12人日，洁净首选）",
      "C. 暗色霓虹（18人日，大屏专用）",
      "B + C 双皮肤（23人日，全覆盖 ★推荐）"]),
    ("STEP 2", Inches(4.5), INDIGO,
     "可选项：局部微调",
     ["主色再改一下（例：B 的薄荷绿→蓝一点）",
      "侧栏样式就用 A 不变动",
      "KPI 卡风格照搬 B，其它用 A",
      "→ 任意混搭，告知我调整点即可"]),
    ("STEP 3", Inches(8.4), PURPLE,
     "进入实施：阶段 1 做视觉走查 Demo",
     ["阶段 1（1-2 天）交付「登录页+看板+侧栏」可交互 Demo",
      "Demo 走查通过 → 进入阶段 2-5",
      "Demo 需要改 → 当场调整参数后再推进",
      "全程小步快跑，避免一次性大翻车"]),
]
for step, x, clr, title, items in steps:
    add_rect(s, x, Inches(1.3), Inches(4.0), Inches(5.7),
             fill=WHITE, line=BORDER, line_w=Pt(1), radius=0.05)
    add_rect(s, x + Inches(0.25), Inches(1.5), Inches(3.5), Inches(1.0),
             fill=clr, radius=0.06)
    add_text(s, x + Inches(0.25), Inches(1.55), Inches(3.5), Inches(0.4),
             step, size=12, bold=True, color=RGBColor(0xE0,0xE7,0xFF))
    add_text(s, x + Inches(0.25), Inches(1.9), Inches(3.5), Inches(0.6),
             title, size=18, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, x + Inches(0.25), Inches(2.7), Inches(3.5), Inches(4.0), items,
                size=12.5, color=TEXT, bullet="›", line_space=8)

# ========= 页 21 前后对比（Before vs After A） =========
s = prs.slides.add_slide(BLANK)
add_page_banner(s, "附 · 前后对比：Before vs After（方案 A 示意）",
                "登录页 + 看板 + 设备台账视觉升级样例", 21, TOTAL)

# Before 列
add_rect(s, Inches(0.3), Inches(1.1), Inches(6.3), Inches(0.4), fill=MUTED, radius=0.03)
add_text(s, Inches(0.3), Inches(1.1), Inches(6.3), Inches(0.4),
         "● Before（现状）", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# After 列
add_rect(s, Inches(6.75), Inches(1.1), Inches(6.3), Inches(0.4), fill=BLUE, radius=0.03)
add_text(s, Inches(6.75), Inches(1.1), Inches(6.3), Inches(0.4),
         "● After（方案 A 现代蓝 · AI 示意）", size=14, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Before 图
add_text(s, Inches(0.35), Inches(1.6), Inches(6.2), Inches(0.3),
         "登录页 / 看板", size=10, bold=True, color=MUTED)
add_image(f"{IMGDIR}/before_dashboard.png", s,
          Inches(0.35), Inches(1.9), Inches(6.2),
          caption="（现有系统真实截图）",
          caption_color=MUTED)
# After 图
add_text(s, Inches(6.8), Inches(1.6), Inches(6.2), Inches(0.3),
         "方案 A 看板（AI 设计稿）", size=10, bold=True, color=BLUE)
add_image(f"{IMGDIR}/design_scheme_A_modern_blue.jpg", s,
          Inches(6.8), Inches(1.9), Inches(6.2),
          caption="（实际落地后细节会以真实组件渲染为准）",
          caption_color=MUTED)

# ========= 页 22 结束页 =========
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, fill=SLATE)
add_rect(s, 0, Inches(3.4), SW, Inches(0.05), fill=BLUE)
add_rect(s, 0, Inches(3.48), Inches(5), Inches(0.05), fill=GREEN)
add_rect(s, 0, Inches(3.56), Inches(3), Inches(0.05), fill=CYAN)
add_text(s, 0, Inches(2.1), SW, Inches(1.0),
         "请选择您倾向的方案", size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(3.8), SW, Inches(0.7),
         "A · 现代蓝   /   B · 极简青绿   /   C · 暗色霓虹   /   B + C 双皮肤",
         size=20, color=BLUE, align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(5.0), SW, Inches(0.4),
         "或者对任一方案提出具体修改点（色、布局、侧栏、卡片…）",
         size=14, color=RGBColor(0x94,0xA3,0xB8), align=PP_ALIGN.CENTER)
add_text(s, 0, SH - Inches(1.1), SW, Inches(0.4),
         "选定后立即进入阶段 1 · 产出可交互 Demo 视觉走查",
         size=14, color=GREEN, align=PP_ALIGN.CENTER)


# ========= 保存 =========
prs.save(OUT)
print(f"UI 评审 PPT 已生成: {OUT}")
print(f"总页数: {len(prs.slides)}")
