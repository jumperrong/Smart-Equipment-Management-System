#!/usr/bin/env python3
"""生成 SEMS 功能介绍 PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import os

SCREENSHOTS_DIR = "/workspace/screenshots"
OUTPUT_PATH = "/workspace/SEMS_功能介绍.pptx"

# 颜色定义
PRIMARY = RGBColor(0x0B, 0x3D, 0x91)       # 深蓝
ACCENT = RGBColor(0x00, 0x96, 0x88)         # 青绿
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x60, 0x60, 0x60)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)

# 页面定义
SLIDES = [
    {"type": "cover"},
    {"type": "overview"},
    {"type": "feature", "title": "看板总览", "subtitle": "Dashboard Overview",
     "image": "01_dashboard.png",
     "points": ["设备实时状态一览：RUN / IDLE / DOWN / PM / ENGINEERING",
                "状态变更轨迹可视化（from → to）",
                "超时预警提醒，PM 进行中实时跟踪",
                "按厂区、区域、状态多维度筛选"]},
    {"type": "feature", "title": "设备台账", "subtitle": "Equipment Registry",
     "image": "02_equipment.png",
     "points": ["完整的设备档案管理（名称、编号、厂区、区域、供应商）",
                "支持按关键字、厂区、区域、状态查询",
                "右上角一键切换设备状态（DOWN 自动派工单）",
                "设备详情、编辑、删除等操作便捷"]},
    {"type": "feature", "title": "点检巡检", "subtitle": "Inspection & Patrol",
     "image": "03_inspection.png",
     "points": ["点检模板管理：日检 / 周检 / 月检",
                "每台设备关联独立检查项（4-8 项）",
                "点检记录历史可追溯",
                "一键发起点检操作"]},
    {"type": "feature", "title": "工单管理", "subtitle": "Work Order Management",
     "image": "04_workorders.png",
     "points": ["REPAIR 故障维修 + PM 预防性维护，两大工单类型",
                "工单持续时长：列表 & 详情可见（创建到关闭）",
                "关键词检索：支持标题/描述/现象模糊搜索",
                "设备 DOWN 状态切换自动创建 REPAIR 工单"]},
    {"type": "feature", "title": "PM 维护计划", "subtitle": "PM Maintenance Plan",
     "image": "06_pm_plans.png",
     "points": ["周期性 PM 计划：周 / 双周 / 月 / 季度",
                "计划到期自动提醒（已到期 / 即将到期）",
                "维护项目清单明确（如腔体清洁、校准、更换配件）",
                "一键生成到期 PM 工单"]},
    {"type": "feature", "title": "备件管理", "subtitle": "Spare Parts Management",
     "image": "07_spare_parts.png",
     "points": ["备件库存管理：编号、名称、规格、库存数量",
                "出入库流水记录，追溯每一笔变动",
                "设备易损件清单关联",
                "低库存预警提醒"]},
    {"type": "feature", "title": "工艺文件", "subtitle": "Process Documents",
     "image": "08_process_docs.png",
     "points": ["两大分类：指导性文件 + 作业记录文件",
                "指导性文件版本管理：草稿→生效→作废",
                "作业记录支持 PDF 附件 + 结构化电子表单双模式",
                "电子表单一键查看填写详情 + 导出 JSON/CSV"]},
    {"type": "feature", "title": "表单模板管理", "subtitle": "Form Template Manager",
     "image": "02_equipment.png",
     "points": ["管理员定义模板：字段(9种类型)/选项/单位/范围",
                "支持上传 PDF/Excel 参考模板文件供填写时对照",
                "模板分类：作业记录类 / 通用表单类",
                "启用/停用控制，仅启用模板可生成记录"]},
    {"type": "feature", "title": "电子表单填写", "subtitle": "Structured Form Filling",
     "image": "02_equipment.png",
     "points": ["作业记录页「新建电子表单」：选模板→动态渲染→填写",
                "已有结构化记录：点击「填写」按钮继续编辑",
                "草稿 / 已提交 / 已作废 三态流转",
                "元数据：关联机台、批次、班次、生产日期自动归档"]},
    {"type": "feature", "title": "文控系统 - 审批链", "subtitle": "Document Control - Approval Chain",
     "image": "08_process_docs.png",
     "points": ["三级电子签名：编制提交→审核→批准生效",
                "每次签署需二次密码校验 + SHA256 签名指纹",
                "状态机白名单：草稿→审核中→生效→作废",
                "驳回路径：退回草稿 + 强制填写原因"]},
    {"type": "feature", "title": "文控系统 - 受控管理", "subtitle": "Document Control - Distribution",
     "image": "08_process_docs.png",
     "points": ["修订记录：字段级 before/after 对比 + 变更原因",
                "分发收回台账：USER/DEPT 批量分发 + 收回记录",
                "PDF 下载自动加盖受控章（编号+状态+用户+日期）",
                "复审周期告警：30 天内到期 / 已过期 Badge"]},
    {"type": "feature", "title": "文控系统 - 表单审核", "subtitle": "Document Control - Form Audit",
     "image": "08_process_docs.png",
     "points": ["表单审核锁定：已审核记录禁止原地修改",
                "附加修正流程：PENDING→APPROVED + 密码校验",
                "记录原值/修正值/修正原因全链路追溯",
                "QA 角色权限：独立于工艺员，审核批准分离"]},
    {"type": "feature", "title": "OEE 分析", "subtitle": "OEE Analysis",
     "image": "09_oee.png",
     "points": ["设备综合效率统计（可用性 × 性能 × 质量）",
                "按日/周/月维度趋势图展示",
                "设备对比分析视图",
                "OEE 组成拆解，瓶颈定位"]},
    {"type": "feature", "title": "品管工具 (8D/FMEA)", "subtitle": "Quality Tools",
     "image": "10_quality.png",
     "points": ["8D 报告：结构化问题分析与解决",
                "FMEA 分析：失效模式与影响分析",
                "可靠性指标：MTBF / MTTR 统计",
                "按设备关联品质事件"]},
    {"type": "feature", "title": "环境核查", "subtitle": "Environment Monitoring",
     "image": "11_environment.png",
     "points": ["温湿度、洁净度等环境参数记录",
                "按时间段查询历史数据",
                "异常数据自动标记",
                "支持按区域筛选"]},
    {"type": "feature", "title": "人员管理", "subtitle": "Personnel Management",
     "image": "12_personnel.png",
     "points": ["资质考核：设备操作等级认证（主操作/副操作/培训中）",
                "技能矩阵：人员与设备技能匹配一览",
                "培训计划：计划制定与执行跟踪",
                "资质到期自动提醒"]},
    {"type": "feature", "title": "资产管理", "subtitle": "Asset Management",
     "image": "13_asset.png",
     "points": ["资产盘点：创建盘点计划，逐项核查",
                "调拨报废申请：资产变动流程管理",
                "盘点进度可视化（如 10/10 完成）",
                "盘点历史可追溯"]},
    {"type": "feature", "title": "安全加固 & 灾备方案", "subtitle": "Security & Disaster Recovery",
     "image": "14_system_config.png",
     "points": ["局域网安全：密码策略/账户锁定/JWT双令牌/IP白名单",
                "3-2-1 灾备：AES加密 + NAS异地副本 + U盘冷备",
                "备份还原烟雾测试（每次备份都知道能不能还原）",
                "敏感操作审计日志（登录/改密/恢复等）"]},
    {"type": "feature", "title": "服务守护 & 自启动", "subtitle": "Service Robustness",
     "image": "14_system_config.png",
     "points": ["4 套守护方案：systemd / NSSM / Docker / 看门狗",
                "开机自启：systemd WantedBy / NSSM 延迟启动 / cron @reboot",
                "崩溃自重启：on-failure / NSSM restart / watchdog tick",
                "健康检查兜底：每 2 分钟检测，连续失败 3 次自动重启"]},
    {"type": "feature", "title": "系统配置", "subtitle": "System Configuration",
     "image": "14_system_config.png",
     "points": ["字典管理：自定义厂区、区域、状态等选项",
                "角色权限：精细化角色权限矩阵",
                "用户管理：账号创建、强制改密、账户解锁",
                "定时备份计划、一键备份/加密/恢复"]},
    {"type": "end"},
]


def get_image_size(img_path):
    """获取图片尺寸"""
    with Image.open(img_path) as img:
        return img.size  # (width, height)


def add_bg_shape(slide, color, left, top, width, height):
    """添加背景色块"""
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def create_cover(prs, slide_layout):
    """创建封面"""
    slide = prs.slides.add_slide(slide_layout)
    # 深蓝背景
    bg = add_bg_shape(slide, PRIMARY, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    # 标题
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(11.333), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SEMS 半导体设备管理系统"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    # 副标题
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.333), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Semiconductor Equipment Management System"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(0x4F, 0xC3, 0xF7)
    p2.alignment = PP_ALIGN.CENTER
    # 版本信息
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(0.6))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "功能介绍  |  v1.2.0  |  2026"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(0x90, 0xCA, 0xF9)
    p3.alignment = PP_ALIGN.CENTER
    # 装饰线
    line = add_bg_shape(slide, ACCENT, Inches(5), Inches(3.5), Inches(3.333), Pt(3))


def create_overview(prs, slide_layout):
    """创建功能总览页"""
    slide = prs.slides.add_slide(slide_layout)
    # 标题
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "功能总览"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # 功能模块网格
    modules = [
        ("看板总览", "实时状态监控"), ("设备台账", "设备档案+状态切换"),
        ("点检巡检", "模板与记录"), ("工单管理", "持续时长+关键词"),
        ("PM维护计划", "周期性保养"), ("备件管理", "库存与流水"),
        ("工艺文件", "版本+电子表单"), ("表单模板", "结构化字段定义"),
        ("文控系统", "审批链+受控章"), ("OEE 分析", "效率统计"), ("品管工具", "8D / FMEA"),
        ("环境核查", "参数监控"), ("人员管理", "资质与培训"),
        ("资产管理", "盘点与调拨"), ("安全灾备", "加固+3-2-1备份"),
        ("服务守护", "自启+崩溃重启"), ("系统配置", "权限与设置"),
    ]
    cols = 4
    card_w = Inches(2.8)
    card_h = Inches(1.3)
    start_x = Inches(0.6)
    start_y = Inches(1.4)
    gap_x = Inches(0.15)
    gap_y = Inches(0.2)

    for i, (name, desc) in enumerate(modules):
        row = i // cols
        col = i % cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        # 卡片背景
        card = add_bg_shape(slide, LIGHT_BG, x, y, card_w, card_h)
        # 标题
        txBox = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), card_w - Inches(0.4), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        # 描述
        txBox2 = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.65), card_w - Inches(0.4), Inches(0.5))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY


def create_feature_slide(prs, slide_layout, data):
    """创建功能展示页"""
    slide = prs.slides.add_slide(slide_layout)
    # 顶部标题栏
    header = add_bg_shape(slide, PRIMARY, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(8), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = data["title"]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    # 副标题
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.55), Inches(8), Inches(0.4))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = data["subtitle"]
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(0x90, 0xCA, 0xF9)

    img_path = os.path.join(SCREENSHOTS_DIR, data["image"])
    if os.path.exists(img_path):
        img_w, img_h = get_image_size(img_path)
        # 截图区域：左侧大图
        max_img_w = Inches(8.5)
        max_img_h = Inches(5.5)
        ratio = min(max_img_w / img_w, max_img_h / img_h)
        scaled_w = int(img_w * ratio)
        scaled_h = int(img_h * ratio)
        img_left = Inches(0.4)
        img_top = Inches(1.3) + (Inches(5.8) - scaled_h) // 2
        slide.shapes.add_picture(img_path, img_left, img_top, scaled_w, scaled_h)
    else:
        # 没有截图时左侧放置占位卡+文字
        place_x = Inches(0.4)
        place_y = Inches(1.3)
        place_w = Inches(8.5)
        place_h = Inches(5.8)
        add_bg_shape(slide, LIGHT_BG, place_x, place_y, place_w, place_h)
        # 占位图标（用大号文字）
        phBox = slide.shapes.add_textbox(place_x, place_y + Inches(2.2), place_w, Inches(0.8))
        phTf = phBox.text_frame
        phP = phTf.paragraphs[0]
        phP.text = "📸"
        phP.font.size = Pt(48)
        phP.alignment = PP_ALIGN.CENTER
        phP.font.color.rgb = GRAY
        # 占位说明
        phBox2 = slide.shapes.add_textbox(place_x, place_y + Inches(3.1), place_w, Inches(0.6))
        phTf2 = phBox2.text_frame
        phP2 = phTf2.paragraphs[0]
        phP2.text = "(截图占位)"
        phP2.font.size = Pt(16)
        phP2.alignment = PP_ALIGN.CENTER
        phP2.font.color.rgb = GRAY

    # 右侧功能要点
    right_x = Inches(9.2)
    right_y = Inches(1.3)
    right_w = Inches(3.8)
    # 要点背景
    add_bg_shape(slide, LIGHT_BG, right_x - Inches(0.1), right_y - Inches(0.1), right_w + Inches(0.2), Inches(5.8))

    # 要点标题
    txBox3 = slide.shapes.add_textbox(right_x, right_y + Inches(0.1), right_w, Inches(0.5))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "功能亮点"
    p3.font.size = Pt(18)
    p3.font.bold = True
    p3.font.color.rgb = PRIMARY

    # 要点列表
    txBox4 = slide.shapes.add_textbox(right_x, right_y + Inches(0.7), right_w, Inches(5.0))
    tf4 = txBox4.text_frame
    tf4.word_wrap = True
    for i, point in enumerate(data["points"]):
        if i == 0:
            p4 = tf4.paragraphs[0]
        else:
            p4 = tf4.add_paragraph()
        p4.text = f"  {point}"
        p4.font.size = Pt(13)
        p4.font.color.rgb = DARK
        p4.space_after = Pt(12)
        p4.line_spacing = Pt(1.4)


def create_end_slide(prs, slide_layout):
    """创建结尾页"""
    slide = prs.slides.add_slide(slide_layout)
    # 深蓝背景
    add_bg_shape(slide, PRIMARY, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    # 装饰线
    add_bg_shape(slide, ACCENT, Inches(5), Inches(3.0), Inches(3.333), Pt(3))
    # 感谢
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3.3), Inches(11.333), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢观看"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    # 技术栈
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(0.6))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "FastAPI  +  Vue 3  +  Element Plus  +  SQLite"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0x4F, 0xC3, 0xF7)
    p2.alignment = PP_ALIGN.CENTER
    # 默认账号
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(0.6))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "默认账号：admin / admin123"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0x90, 0xCA, 0xF9)
    p3.alignment = PP_ALIGN.CENTER


def main():
    prs = Presentation()
    # 16:9 宽屏
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_layout = prs.slide_layouts[6]  # 空白布局

    for slide_data in SLIDES:
        stype = slide_data["type"]
        if stype == "cover":
            create_cover(prs, slide_layout)
        elif stype == "overview":
            create_overview(prs, slide_layout)
        elif stype == "feature":
            create_feature_slide(prs, slide_layout, slide_data)
        elif stype == "end":
            create_end_slide(prs, slide_layout)

    prs.save(OUTPUT_PATH)
    print(f"PPT 已生成: {OUTPUT_PATH}")
    print(f"总页数: {len(prs.slides)}")


if __name__ == "__main__":
    main()
